#!/usr/bin/env python3
"""Model-free grading renderer preflight for CI and local diagnostics."""

from __future__ import annotations

import base64
import contextlib
import hashlib
import io
import json
import os
import shutil
import subprocess
import sys
import tempfile
import types
from collections.abc import Mapping
from pathlib import Path
from typing import Any

BATCH_RUNNER_ROOT = Path(__file__).resolve().parent.parent
if str(BATCH_RUNNER_ROOT) not in sys.path:
    sys.path.insert(0, str(BATCH_RUNNER_ROOT))
if "core" not in sys.modules:
    core_package = types.ModuleType("core")
    core_package.__path__ = [str(BATCH_RUNNER_ROOT / "core")]
    core_package.__package__ = "core"
    sys.modules["core"] = core_package

from core.tools import get_renderer_fingerprint, read_deliverable
from core.tools.read_deliverable import MAX_IMAGE_BYTES

FONT_FAMILY = "Liberation Sans"
PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"
SUBPROCESS_TIMEOUT_SEC = 10


class RendererPreflightError(RuntimeError):
    """Raised when a renderer preflight invariant is not satisfied."""


def _minimal_subprocess_env(work_dir: Path) -> dict[str, str]:
    environment = {
        "PATH": os.environ.get("PATH", "/usr/bin:/bin"),
        "HOME": str(work_dir),
        "TMPDIR": str(work_dir),
    }
    for key, value in os.environ.items():
        if key == "LANG" or key.startswith("LC_"):
            environment[key] = value
    environment.setdefault("LC_ALL", "C.UTF-8")
    return environment


def _validate_liberation_sans(work_dir: Path) -> str:
    executable = shutil.which("fc-match")
    if executable is None:
        raise RendererPreflightError("fc-match executable not found")
    try:
        completed = subprocess.run(
            [executable, "--format", "%{family}\n", FONT_FAMILY],
            capture_output=True,
            text=True,
            timeout=SUBPROCESS_TIMEOUT_SEC,
            check=False,
            env=_minimal_subprocess_env(work_dir),
            shell=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise RendererPreflightError("fc-match timed out") from exc
    except OSError as exc:
        raise RendererPreflightError(f"fc-match failed: {exc}") from exc
    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout or "no output")[-300:]
        raise RendererPreflightError(
            f"fc-match failed (exit {completed.returncode}): {detail}"
        )
    families = [
        family.strip()
        for line in completed.stdout.splitlines()
        for family in line.split(",")
        if family.strip()
    ]
    if FONT_FAMILY not in families:
        raise RendererPreflightError(
            f"fc-match did not resolve {FONT_FAMILY}: {families!r}"
        )
    return FONT_FAMILY


def _write_xlsx_fixture(path: Path) -> None:
    from openpyxl import Workbook  # type: ignore
    from openpyxl.styles import Font  # type: ignore

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Renderer Preflight"
    worksheet["A1"] = "GDPVal grading renderer preflight"
    worksheet["A1"].font = Font(name=FONT_FAMILY, size=16, bold=True)
    worksheet.column_dimensions["A"].width = 40
    workbook.save(path)
    workbook.close()


def _write_pptx_fixture(path: Path) -> None:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    presentation = Presentation()
    for slide_number in (1, 2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1.5)
        )
        paragraph = text_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = f"GDPVal renderer preflight slide {slide_number}"
        run.font.name = FONT_FAMILY
        run.font.size = Pt(28)
        run.font.bold = True
    presentation.save(path)


def _write_docx_fixture(path: Path) -> None:
    from docx import Document  # type: ignore
    from docx.shared import Pt  # type: ignore

    document = Document()
    # A plain bold run rather than add_heading(): heading styles come from the
    # template and a missing one raises, which would report a python-docx
    # template gap as a LibreOffice failure.
    title = document.add_paragraph()
    title_run = title.add_run("GDPVal grading renderer preflight")
    title_run.font.name = FONT_FAMILY
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    body = document.add_paragraph()
    body_run = body.add_run(
        "Document body paragraph rendered through LibreOffice Writer."
    )
    body_run.font.name = FONT_FAMILY
    body_run.font.size = Pt(12)
    document.save(path)


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _validate_fingerprint(fingerprint: Any) -> dict[str, str]:
    if not isinstance(fingerprint, Mapping):
        raise RendererPreflightError("renderer fingerprint is not an object")
    required = (
        "libreoffice_binary",
        "libreoffice_version",
        "pymupdf_version",
    )
    normalized: dict[str, str] = {}
    for key in required:
        value = fingerprint.get(key)
        if not isinstance(value, str) or not value.strip():
            raise RendererPreflightError(
                f"renderer fingerprint missing non-empty {key}"
            )
        normalized[key] = value
    if Path(normalized["libreoffice_binary"]).name != normalized[
        "libreoffice_binary"
    ]:
        raise RendererPreflightError(
            "renderer fingerprint exposes an executable path"
        )
    return normalized


def _validate_render_result(
    result: Any,
    *,
    source_kind: str,
    scope: dict[str, int],
    fingerprint: dict[str, str],
) -> dict[str, Any]:
    if not isinstance(result, Mapping) or result.get("ok") is not True:
        raise RendererPreflightError(
            f"{source_kind} render failed: {result!r}"
        )
    data = result.get("data")
    if not isinstance(data, Mapping):
        raise RendererPreflightError(f"{source_kind} render data is not an object")
    if data.get("kind") != "image_png_base64":
        raise RendererPreflightError(
            f"{source_kind} render kind is not image_png_base64"
        )
    encoded = data.get("base64")
    if not isinstance(encoded, str) or not encoded:
        raise RendererPreflightError(f"{source_kind} render has no base64 PNG")
    try:
        decoded = base64.b64decode(encoded, validate=True)
    except (ValueError, TypeError) as exc:
        raise RendererPreflightError(
            f"{source_kind} render base64 is invalid"
        ) from exc
    if decoded[: len(PNG_SIGNATURE)] != PNG_SIGNATURE:
        raise RendererPreflightError(
            f"{source_kind} render does not have the full PNG signature"
        )
    byte_size = data.get("byte_size")
    if isinstance(byte_size, bool) or not isinstance(byte_size, int):
        raise RendererPreflightError(f"{source_kind} byte_size is not an integer")
    if byte_size != len(decoded):
        raise RendererPreflightError(
            f"{source_kind} byte_size mismatch: {byte_size} != {len(decoded)}"
        )
    if byte_size > MAX_IMAGE_BYTES:
        raise RendererPreflightError(
            f"{source_kind} render exceeds {MAX_IMAGE_BYTES} byte cap"
        )
    if data.get("source_kind") != source_kind:
        raise RendererPreflightError(
            f"{source_kind} render source_kind mismatch: {data.get('source_kind')!r}"
        )
    if data.get("scope") != scope:
        raise RendererPreflightError(
            f"{source_kind} render scope mismatch: {data.get('scope')!r}"
        )
    renderer = data.get("renderer")
    if not isinstance(renderer, Mapping):
        raise RendererPreflightError(f"{source_kind} renderer metadata missing")
    for key, value in fingerprint.items():
        if renderer.get(key) != value:
            raise RendererPreflightError(
                f"{source_kind} renderer fingerprint mismatch for {key}"
            )
    if renderer.get("converter") != "libreoffice":
        raise RendererPreflightError(
            f"{source_kind} converter is not libreoffice"
        )
    if renderer.get("rasterizer") != "pymupdf":
        raise RendererPreflightError(
            f"{source_kind} rasterizer is not pymupdf"
        )
    return {
        "source_kind": source_kind,
        "scope": scope,
        "byte_size": byte_size,
    }


def run_preflight() -> dict[str, Any]:
    with tempfile.TemporaryDirectory(
        prefix="gdpval-grading-renderer-preflight-"
    ) as temp:
        work_dir = Path(temp)
        font_match = _validate_liberation_sans(work_dir)
        fingerprint = _validate_fingerprint(get_renderer_fingerprint())

        xlsx_path = work_dir / "renderer_preflight.xlsx"
        pptx_path = work_dir / "renderer_preflight.pptx"
        docx_path = work_dir / "renderer_preflight.docx"
        _write_xlsx_fixture(xlsx_path)
        _write_pptx_fixture(pptx_path)
        _write_docx_fixture(docx_path)
        # docx is here because it is the format this preflight was blind to:
        # the grading runners installed libreoffice-core/-calc/-impress and no
        # -writer, so a Writer-shaped gap could not fail anything. xlsx and
        # pptx render through Calc and Impress and would keep passing.
        fixtures = (
            (xlsx_path, "xlsx", {"workbook_page": 1}),
            (pptx_path, "pptx", {"slide": 1}),
            (docx_path, "docx", {"page": 1}),
        )
        before_hashes = {path.name: _sha256_file(path) for path, _, _ in fixtures}
        renders: list[dict[str, Any]] = []
        for path, source_kind, scope in fixtures:
            result = read_deliverable(
                "render_to_image",
                path.name,
                base_dir=str(work_dir),
                scope=scope,
            )
            renders.append(
                _validate_render_result(
                    result,
                    source_kind=source_kind,
                    scope=scope,
                    fingerprint=fingerprint,
                )
            )
        after_hashes = {path.name: _sha256_file(path) for path, _, _ in fixtures}
        if before_hashes != after_hashes:
            raise RendererPreflightError("renderer mutated a source fixture")

    return {
        "ok": True,
        "font_family": font_match,
        "renderer_fingerprint": fingerprint,
        "renders": renders,
    }


def main() -> int:
    captured_stdout = io.StringIO()
    captured_stderr = io.StringIO()
    try:
        with contextlib.redirect_stdout(captured_stdout), contextlib.redirect_stderr(
            captured_stderr
        ):
            payload = run_preflight()
        exit_code = 0
    except Exception as exc:  # noqa: BLE001
        payload = {
            "ok": False,
            "error_type": type(exc).__name__,
            "error": str(exc)[:500],
        }
        exit_code = 1
    sys.stdout.write(
        json.dumps(payload, ensure_ascii=True, separators=(",", ":"), sort_keys=True)
        + "\n"
    )
    return exit_code


if __name__ == "__main__":
    sys.exit(main())