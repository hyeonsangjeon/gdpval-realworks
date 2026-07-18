"""Non-paid security tests for the agentic compute plane."""

from __future__ import annotations

import os
import json
import hashlib
import struct
import subprocess
import sys
import wave
from pathlib import Path
from types import SimpleNamespace

import pytest

from core.agentic_compute import AgenticDockerBackend, _container_path
from core import agentic_verifier
from sandbox.agentic_image_prepare import _is_forbidden_executable


SECCOMP = Path(__file__).resolve().parent.parent / "sandbox" / "agentic-seccomp.json"
LAUNCHER = Path(__file__).resolve().parent.parent / "core" / "agentic_python_launcher.py"


def _backend(tmp_path, reference_files=None):
    return AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=list(reference_files or []),
        occupation="Analyst",
        image="gdpval-agentic-sandbox:test",
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        require_approved_input_manifest=False,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
    )


def test_task_container_command_has_hardening_flags(tmp_path):
    backend = _backend(tmp_path)
    try:
        command = backend._task_container_command()
        joined = " ".join(command)

        for required in (
            "--network none", "--ipc none", "--read-only",
            "--cap-drop ALL", "--pids-limit 128", "--ulimit nofile=256:256",
            "--user 65532:65532", "--security-opt no-new-privileges",
            "type=volume", "dst=/work", "volume-nocopy",
            "readonly", "bind-propagation=rprivate",
        ):
            assert required in joined
        assert "--pid host" not in joined
    finally:
        backend.close()


def test_outer_seccomp_is_default_deny_allowlist():
    profile = json.loads(SECCOMP.read_text(encoding="utf-8"))
    allowed = {
        name
        for rule in profile["syscalls"]
        if rule["action"] == "SCMP_ACT_ALLOW"
        for name in rule["names"]
    }

    assert profile["defaultAction"] == "SCMP_ACT_ERRNO"
    assert "seccomp" in allowed
    for denied in (
        "bpf", "mount", "setns", "unshare", "ptrace", "keyctl",
        "io_uring_setup", "userfaultfd", "process_vm_readv",
    ):
        assert denied not in allowed


def test_inner_filter_denies_seccomp_and_queued_signals():
    from core.agentic_python_launcher import DENIED_SYSCALLS

    assert {"seccomp", "rt_sigqueueinfo", "rt_tgsigqueueinfo"} <= set(
        DENIED_SYSCALLS
    )


@pytest.mark.parametrize(
    "name",
    [
        "gcc", "gcc-12", "x86_64-linux-gnu-gcc", "x86_64-linux-gnu-gcc-12",
        "g++-12", "x86_64-linux-gnu-ld.bfd", "as", "objcopy", "pip3.11",
    ],
)
def test_image_preparation_detects_versioned_and_prefixed_build_tools(name):
    assert _is_forbidden_executable(name) is True


def test_production_backend_requires_pinned_sbom_identity():
    with pytest.raises(ValueError, match="SBOM hash is required"):
        AgenticDockerBackend(
            task_prompt="task",
            reference_files=[],
            occupation="Analyst",
            image="image@sha256:" + "a" * 64,
            approved_input_manifest={},
        )


def test_substrate_manifest_is_condition_neutral(tmp_path):
    first = _backend(tmp_path)
    second = _backend(tmp_path)
    try:
        first.image_id = second.image_id = "sha256:" + "a" * 64
        first.verifier_image_id = second.verifier_image_id = "sha256:" + "a" * 64
        assert first.substrate_manifest() == second.substrate_manifest()
        serialized = json.dumps(first.substrate_manifest(), sort_keys=True)
        assert "baseline" not in serialized
        assert "treatment" not in serialized
    finally:
        first.close()
        second.close()


def test_verifier_image_components_are_hashed_from_actual_image_id(tmp_path):
    backend = _backend(tmp_path)
    expected = backend._component_hashes()
    commands = []

    def run(command, **_):
        commands.append(command)
        if command[1:3] == ["image", "inspect"]:
            output = json.dumps([{"Id": "sha256:" + "f" * 64}]).encode()
        else:
            output = json.dumps({
                "verifier": expected["verifier"],
                "capabilities": expected["capabilities"],
                "core_tree": expected["core_tree"],
                "sbom": backend.sbom_sha256,
            }).encode()
        return SimpleNamespace(returncode=0, stdout=output, stderr=b"")

    backend._run = run
    try:
        backend._verify_verifier_image_components()
    finally:
        backend.close()

    verifier_command = commands[1]
    assert "sha256:" + "f" * 64 in verifier_command
    assert verifier_command[verifier_command.index("--network") + 1] == "none"
    assert "--read-only" in verifier_command


def test_verifier_component_probe_timeout_remains_registered_until_cleanup(
    tmp_path
):
    backend = _backend(tmp_path)
    commands = []

    def run(command, **kwargs):
        commands.append(command)
        if command[1:3] == ["image", "inspect"]:
            return SimpleNamespace(
                returncode=0,
                stdout=json.dumps([{"Id": "sha256:" + "f" * 64}]).encode(),
                stderr=b"",
            )
        if command[1] == "run":
            raise subprocess.TimeoutExpired(command, 1)
        return SimpleNamespace(returncode=0, stdout=b"", stderr=b"")

    backend._run = run
    try:
        with pytest.raises(subprocess.TimeoutExpired):
            backend._verify_verifier_image_components()
    finally:
        backend.close()

    run_command = next(command for command in commands if command[1] == "run")
    probe_name = run_command[run_command.index("--name") + 1]
    assert ["docker", "rm", "-f", probe_name] in commands
    assert probe_name not in backend._may_exist_containers


def test_startup_passes_one_deadline_to_every_stage(tmp_path, monkeypatch):
    backend = _backend(tmp_path)
    deadlines = []

    def stage(name):
        def call(deadline):
            deadlines.append((name, deadline))
            if name == "verifier":
                backend.verifier_image_id = "sha256:" + "f" * 64
                backend._verified_component_hashes = backend._component_hashes()
            if name == "volume":
                backend.work_volume_created = True
                backend._may_exist_volumes.add(backend.work_volume_name)
        return call

    monkeypatch.setattr(backend, "_verify_host_runtime", stage("host"))
    monkeypatch.setattr(
        backend, "_verify_verifier_image_components", stage("verifier")
    )
    monkeypatch.setattr(backend, "_stage_inputs", stage("inputs"))
    monkeypatch.setattr(backend, "_create_work_volume", stage("volume"))
    monkeypatch.setattr(backend, "_verify_runtime", stage("runtime"))
    monkeypatch.setattr(backend, "_assert_pid1_only", stage("pid"))
    monkeypatch.setattr(
        backend,
        "_run",
        lambda command, **kwargs: subprocess.CompletedProcess(
            command, 0, stdout=b"container-id\n", stderr=b""
        ),
    )
    try:
        result = backend.start(5)
        assert result["ok"] is True, result
    finally:
        backend._may_exist_containers.clear()
        backend._may_exist_volumes.clear()
        backend.container_started = False
        backend.work_volume_created = False
        backend.close()

    assert [name for name, _ in deadlines] == [
        "host", "verifier", "inputs", "volume", "runtime", "pid"
    ]
    first = deadlines[0][1]
    assert all(deadline == first for _, deadline in deadlines)


def test_startup_exhausted_deadline_stops_before_next_stage(
    tmp_path, monkeypatch
):
    backend = _backend(tmp_path)
    current = [100.0]
    calls = []
    monkeypatch.setattr(
        "core.agentic_compute.time.monotonic", lambda: current[0]
    )

    def verify_host(deadline):
        calls.append("host")
        current[0] = deadline + 0.01

    def verify_components(deadline):
        calls.append("verifier")

    monkeypatch.setattr(backend, "_verify_host_runtime", verify_host)
    monkeypatch.setattr(
        backend, "_verify_verifier_image_components", verify_components
    )

    try:
        result = backend.start(0.01)
    finally:
        backend.close()

    assert result["ok"] is False
    assert result["error_type"] == "container_preflight_failed"
    assert calls == ["host"]


def test_verifier_image_component_mismatch_fails_closed(tmp_path):
    backend = _backend(tmp_path)

    def run(command, **_):
        if command[1:3] == ["image", "inspect"]:
            output = json.dumps([{"Id": "sha256:" + "f" * 64}]).encode()
        else:
            output = json.dumps({
                "verifier": "0" * 64,
                "capabilities": "0" * 64,
                "core_tree": "0" * 64,
            }).encode()
        return SimpleNamespace(returncode=0, stdout=output, stderr=b"")

    backend._run = run
    try:
        with pytest.raises(RuntimeError, match="verifier image component"):
            backend._verify_verifier_image_components()
    finally:
        backend.close()


def test_production_verifier_image_sbom_mismatch_fails_closed(tmp_path):
    backend = _backend(tmp_path)
    backend.require_supply_chain_identity = True
    backend.sbom_sha256 = "a" * 64
    expected = backend._component_hashes()

    def run(command, **_):
        if command[1:3] == ["image", "inspect"]:
            output = json.dumps([{"Id": "sha256:" + "f" * 64}]).encode()
        else:
            output = json.dumps({
                "verifier": expected["verifier"],
                "capabilities": expected["capabilities"],
                "core_tree": expected["core_tree"],
                "sbom": "b" * 64,
            }).encode()
        return SimpleNamespace(returncode=0, stdout=output, stderr=b"")

    backend._run = run
    try:
        with pytest.raises(RuntimeError, match="SBOM identity mismatch"):
            backend._verify_verifier_image_components()
    finally:
        backend.close()


@pytest.mark.parametrize(
    "pid1_fds",
    [
        {"0": "/dev/null", "1": "pipe:[1]", "2": "pipe:[2]", "3": "/tmp/x"},
        {"0": "/dev/null", "1": "socket:[1]", "2": "pipe:[2]"},
    ],
)
def test_runtime_fd_probe_rejects_pid1_fd_or_socket(tmp_path, pid1_fds):
    backend = _backend(tmp_path)
    identity = {
        "uid": 65532,
        "gid": 65532,
        "groups": [],
        "probe_fds": {"0": "/dev/null", "1": "pipe:[1]", "2": "pipe:[2]"},
        "pid1_fds": pid1_fds,
    }
    backend._run = lambda *args, **kwargs: SimpleNamespace(
        returncode=0,
        stdout=json.dumps(identity).encode(),
        stderr=b"",
    )
    try:
        with pytest.raises(RuntimeError, match="pid1"):
            backend._verify_fds()
    finally:
        backend.close()


def test_input_staging_records_exact_hash_and_merkle_root(tmp_path):
    source = tmp_path / "source.txt"
    source.write_text("approved input", encoding="utf-8")
    backend = _backend(tmp_path, [source])
    try:
        backend._stage_inputs()

        assert len(backend.input_records) == 1
        assert backend.input_records[0]["model_path"] == "inputs/source.txt"
        assert len(backend.input_records[0]["sha256"]) == 64
        assert len(backend.input_merkle_root) == 64
        assert (backend.inputs_dir / "source.txt").read_bytes() == source.read_bytes()
        assert (backend.inputs_dir / "source.txt").stat().st_mode & 0o222 == 0
    finally:
        backend.close()


def test_input_staging_rejects_symlink_and_hardlink(tmp_path):
    source = tmp_path / "source.txt"
    source.write_text("input", encoding="utf-8")
    symlink = tmp_path / "link.txt"
    symlink.symlink_to(source)
    backend = _backend(tmp_path, [symlink])
    try:
        with pytest.raises(ValueError, match="input_type_or_link_violation"):
            backend._stage_inputs()
    finally:
        backend.close()


def test_input_staging_requires_exact_approved_manifest(tmp_path):
    source = tmp_path / "source.txt"
    source.write_text("approved input", encoding="utf-8")
    metadata = source.lstat()
    expected = {
        "documents/source.txt": {
            "path": "documents/source.txt",
            "type": "regular",
            "link_count": 1,
            "size_bytes": metadata.st_size,
            "source_allocated_bytes": metadata.st_blocks * 512,
            "sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
            "provider_classification": "approved_public_gdpval",
        }
    }
    backend = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[{
            "source_path": str(source),
            "relative_path": "documents/source.txt",
        }],
        occupation="Analyst",
        image="gdpval-agentic-sandbox:test",
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        approved_input_manifest=expected,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
    )
    try:
        backend._stage_inputs()
        assert backend.input_records[0]["sha256"] == expected[
            "documents/source.txt"
        ][
            "sha256"
        ]
        assert backend.input_records[0]["model_path"] == (
            "inputs/documents/source.txt"
        )
        assert backend.input_records[0]["staged_allocated_bytes"] > 0
    finally:
        backend.close()

    expected["documents/source.txt"]["sha256"] = "0" * 64
    backend = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[{
            "source_path": str(source),
            "relative_path": "documents/source.txt",
        }],
        occupation="Analyst",
        image="gdpval-agentic-sandbox:test",
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        approved_input_manifest=expected,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
    )
    try:
        with pytest.raises(ValueError, match="approved input identity mismatch"):
            backend._stage_inputs()
    finally:
        backend.close()

    hardlink = tmp_path / "hard.txt"
    os.link(source, hardlink)
    backend = _backend(tmp_path, [source])
    try:
        with pytest.raises(ValueError, match="input_type_or_link_violation"):
            backend._stage_inputs()
    finally:
        backend.close()


@pytest.mark.parametrize(
    ("path", "output"),
    [
        ("../escape", True),
        ("/etc/passwd", False),
        ("source.mov", False),
        ("inputs/source.mov", True),
        ("work/.hidden", True),
    ],
)
def test_container_path_rejects_escape_and_wrong_roots(path, output):
    with pytest.raises(ValueError):
        _container_path(path, output=output)


def test_ffmpeg_command_is_closed_and_shell_free(tmp_path):
    backend = _backend(tmp_path)
    try:
        command = backend._ffmpeg_command({
            "operation": "sample_frames",
            "input": "inputs/source.mp4",
            "output": "work/contact.png",
            "frame_count": 4,
            "width": 640,
            "start_seconds": 0,
            "duration_seconds": 10,
        })

        assert command[-1] == "/work/contact.png"
        assert "/usr/bin/ffmpeg" in command
        assert "-nostdin" in command
        assert "-n" in command
        assert command[command.index("-protocol_whitelist") + 1] == "file"
        assert not any(token in command for token in ("sh", "bash", "-c"))
        assert len([token for token in command if token == "/work/contact.png"]) == 1
        with pytest.raises(ValueError, match="unsupported ffmpeg input suffix"):
            backend._ffmpeg_command({
                "operation": "probe",
                "input": "inputs/playlist.m3u8",
            })
        with pytest.raises(ValueError, match="output extension"):
            backend._ffmpeg_command({
                "operation": "transcode_video",
                "input": "inputs/source.mov",
                "output": "work/report.webm",
                "container": "mp4",
                "video_codec": "h264",
                "audio_codec": "aac",
                "width": 640,
                "height": 360,
                "fps": 24,
                "start_seconds": 0,
                "duration_seconds": 1,
            })
        with pytest.raises(ValueError, match="output extension"):
            backend._ffmpeg_command({
                "operation": "extract_audio",
                "input": "inputs/source.mov",
                "output": "work/report.flac",
                "format": "wav",
                "sample_rate": 16000,
                "channels": 1,
                "start_seconds": 0,
                "duration_seconds": 1,
            })
    finally:
        backend.close()


def test_ffprobe_returns_only_allowlisted_bounded_metadata(tmp_path):
    backend = _backend(tmp_path)
    backend.container_started = True
    backend._assert_pid1_only = lambda deadline=None: None
    backend._run_tool = lambda *args, **kwargs: subprocess.CompletedProcess(
        args[0],
        0,
        stdout=json.dumps({
            "format": {
                "filename": "/inputs/private-name.mov",
                "format_name": "mov,mp4",
                "duration": "1.250000",
                "tags": {"comment": "secret"},
            },
            "streams": [{
                "index": 0,
                "codec_name": "h264",
                "codec_type": "video",
                "width": 640,
                "height": 360,
                "tags": {"handler_name": "private"},
            }],
        }).encode(),
        stderr=b"",
    )
    try:
        result = backend.run_ffmpeg({
            "operation": "probe",
            "input": "inputs/source.mov",
        })
    finally:
        backend.container_started = False
        backend.close()

    assert result["ok"] is True
    assert result["data"]["metadata"] == {
        "format": {"duration": "1.250000", "format_name": "mov,mp4"},
        "streams": [{
            "codec_name": "h264",
            "codec_type": "video",
            "height": 360,
            "index": 0,
            "width": 640,
        }],
    }
    assert "private" not in json.dumps(result)


def test_verifier_rechecks_contract_on_exact_selected_deliverables(
    tmp_path, monkeypatch
):
    (tmp_path / "report.pdf").write_bytes(b"full-workspace-primary")
    (tmp_path / "notes.txt").write_text("notes", encoding="utf-8")
    monkeypatch.setattr(agentic_verifier, "SNAPSHOT_ROOT", tmp_path)
    monkeypatch.setattr(
        agentic_verifier,
        "run_output_qa",
        lambda *args, **kwargs: SimpleNamespace(
            ok=True,
            render_reports=[],
            blocking_errors=[],
            warnings=[],
        ),
    )

    result = agentic_verifier.verify({
        "task_prompt": "Create a PDF report",
        "reference_hashes": [],
        "selected_deliverables": ["notes.txt"],
    })

    assert result["ok"] is False
    assert result["error_type"] == "artifact_verification_failed"
    assert result["data"]["artifact_count"] == 1
    assert result["data"]["artifacts"][0]["path"] == "notes.txt"


def test_tool_output_is_bounded_before_host_memory_capture():
    result = AgenticDockerBackend._run_tool(
        [sys.executable, "-c", "import sys;sys.stdout.write('x' * 1000000)"],
        timeout=10,
    )

    assert len(result.stdout) <= 32768
    assert result.returncode != 0


def test_inspection_preserves_oversized_snapshot_for_smaller_finalize_subset(
    tmp_path, monkeypatch
):
    first = tmp_path / "first.txt"
    second = tmp_path / "second.txt"
    first.write_bytes(b"123")
    second.write_bytes(b"456")
    artifacts = [
        {
            "path": path.name,
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }
        for path in (first, second)
    ]
    verification = {"ok": True, "data": {"artifacts": artifacts}}
    backend = _backend(tmp_path)
    backend.container_started = True
    monkeypatch.setattr("core.agentic_compute.MAX_TRANSFER_TOTAL", 4)
    backend._snapshot = lambda **kwargs: (tmp_path, verification)
    backend._strict_snapshot_files = (
        lambda snapshot, deadline=None: [first, second]
    )
    backend._verification_matches_snapshot = lambda *args, **kwargs: True
    try:
        inspected = backend.inspect_artifacts()
        assert inspected["ok"] is True
        assert backend.latest_snapshot == tmp_path
        assert backend.best_result() is None
        selected = {"ok": True, "data": {"artifacts": [artifacts[0]]}}
        result = backend._snapshot_result(
            tmp_path, selected, success=True, summary="selected"
        )
        assert result["files"] == [
            {"filename": "first.txt", "content": b"123"}
        ]
    finally:
        backend.container_started = False
        backend.close()


@pytest.mark.parametrize("failure", ["helper", "unpause"])
def test_snapshot_cleanup_failure_poisons_and_removes_task_container(
    tmp_path, failure
):
    backend = _backend(tmp_path)
    backend.container_started = True
    backend._may_exist_containers.add("snapshot-helper")
    commands = []

    def run(command, **kwargs):
        commands.append(command)
        if (
            failure == "helper"
            and command[1:3] == ["rm", "-f"]
            and command[-1] == "snapshot-helper"
        ):
            return subprocess.CompletedProcess(command, 1, b"", b"busy")
        if failure == "helper" and command[1] == "inspect":
            return subprocess.CompletedProcess(command, 0, b"{}", b"")
        if failure == "unpause" and command[1] == "unpause":
            return subprocess.CompletedProcess(command, 1, b"", b"failed")
        return subprocess.CompletedProcess(command, 0, b"", b"")

    backend._run = run
    with pytest.raises(RuntimeError, match="failed"):
        backend._cleanup_snapshot_resources("snapshot-helper")

    assert backend.poisoned is True
    assert backend.container_started is False
    assert ["docker", "rm", "-f", backend.container_name] in commands


def test_ambiguous_pause_timeout_unpauses_or_removes_task_container(tmp_path):
    backend = _backend(tmp_path)
    backend.container_started = True
    backend._may_exist_containers.add(backend.container_name)
    commands = []

    def run(command, **kwargs):
        commands.append(command)
        if command[1] == "pause":
            raise subprocess.TimeoutExpired(command, kwargs["timeout"])
        if command[1] == "exec":
            return subprocess.CompletedProcess(command, 0, b"ok\n", b"")
        return subprocess.CompletedProcess(command, 0, b"", b"")

    backend._run = run
    try:
        with pytest.raises(subprocess.TimeoutExpired):
            backend._snapshot()

        assert ["docker", "unpause", backend.container_name] in commands
        assert backend.container_started is True
    finally:
        backend.close()


def test_snapshot_rehash_stops_at_absolute_deadline(tmp_path, monkeypatch):
    artifact = tmp_path / "large.bin"
    artifact.write_bytes(b"x" * (2 * 1024 * 1024))
    verification = {
        "ok": True,
        "data": {
            "artifacts": [{
                "path": artifact.name,
                "sha256": hashlib.sha256(artifact.read_bytes()).hexdigest(),
            }]
        },
    }
    times = iter([100.0, 100.0, 102.0])
    monkeypatch.setattr(
        "core.agentic_compute.time.monotonic", lambda: next(times)
    )
    backend = _backend(tmp_path)
    try:
        with pytest.raises(TimeoutError, match="wall time exhausted"):
            backend._verification_matches_snapshot(
                verification, tmp_path, deadline=101.0
            )
    finally:
        backend.close()


@pytest.mark.parametrize("resource", ["container", "volume"])
def test_cleanup_verification_requires_explicit_docker_not_found(
    tmp_path, resource
):
    backend = _backend(tmp_path)
    if resource == "container":
        backend.container_started = True
    else:
        backend.work_volume_created = True

    def run(command, **kwargs):
        return subprocess.CompletedProcess(
            command, 1, stdout=b"", stderr=b"daemon unavailable"
        )

    backend._run = run
    method = (
        backend._remove_container
        if resource == "container"
        else backend._remove_work_volume
    )

    with pytest.raises(RuntimeError, match="could not be verified"):
        method()

    assert backend.poisoned is True
    if resource == "container":
        backend.container_started = False
    else:
        backend.work_volume_created = False
    backend.close()


def test_host_root_cleanup_failure_is_not_ignored(tmp_path, monkeypatch):
    backend = _backend(tmp_path)
    original = __import__("shutil").rmtree

    def fail_target(path, *args, **kwargs):
        if Path(path) == backend.root:
            raise PermissionError("read-only host root")
        return original(path, *args, **kwargs)

    monkeypatch.setattr("core.agentic_compute.shutil.rmtree", fail_target)

    with pytest.raises(PermissionError, match="read-only host root"):
        backend.close()

    original(backend.root)


def test_close_removes_may_exist_resources_without_creation_flags(tmp_path):
    backend = _backend(tmp_path)
    helper = f"{backend.container_name}-snapshot-lost"
    verifier = f"{backend.container_name}-verify-lost"
    backend._may_exist_containers.update({
        backend.container_name, helper, verifier,
    })
    backend._may_exist_volumes.add(backend.work_volume_name)
    backend.container_started = False
    backend.work_volume_created = False
    commands = []

    def run(command, **kwargs):
        commands.append(command)
        return subprocess.CompletedProcess(command, 0, b"", b"")

    backend._run = run
    backend.close()

    for name in (backend.container_name, helper, verifier):
        assert ["docker", "rm", "-f", name] in commands
    assert [
        "docker", "volume", "rm", "-f", backend.work_volume_name
    ] in commands
    assert backend._may_exist_containers == set()
    assert backend._may_exist_volumes == set()


@pytest.mark.skipif(sys.platform != "linux", reason="libseccomp launcher is Linux-only")
def test_generated_python_launcher_denies_exec_and_network():
    source = b"""
import os
import socket

denied = []
for action in (
    lambda: os.system('/bin/true'),
    lambda: socket.socket(socket.AF_INET, socket.SOCK_STREAM),
):
    try:
        action()
    except OSError as exc:
        denied.append(exc.errno)
assert len(denied) == 2, denied
print('blocked')
"""
    harness = (
        "import sys;"
        "from core.agentic_python_launcher import _install_filter;"
        "code=compile(sys.stdin.buffer.read(),'<security-test>','exec');"
        "_install_filter();exec(code,{'__name__':'__main__'})"
    )
    result = subprocess.run(
        [sys.executable, "-c", harness],
        input=source,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=20,
        check=False,
        cwd=Path(__file__).resolve().parent.parent,
    )
    if b"seccomp TSYNC unavailable" in result.stderr:
        pytest.skip("host runtime does not expose seccomp TSYNC; Docker gate remains required")
    assert result.returncode == 0, result.stderr.decode("utf-8", errors="replace")
    assert result.stdout.strip() == b"blocked"


@pytest.mark.integration
def test_agentic_docker_backend_end_to_end(tmp_path):
    image = os.environ.get(
        "AGENTIC_TEST_IMAGE", "gdpval-agentic-sandbox:local"
    )
    inspect = subprocess.run(
        ["docker", "image", "inspect", image],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    )
    if inspect.returncode != 0:
        pytest.skip(f"agentic test image is unavailable: {image}")

    source = tmp_path / "input.wav"
    with wave.open(str(source), "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(16000)
        audio.writeframes(struct.pack("<h", 0) * 16000)
    backend = AgenticDockerBackend(
        task_prompt="Create a professional WAV audio report named report.wav",
        reference_files=[str(source)],
        occupation="Analyst",
        image=image,
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        require_approved_input_manifest=False,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
        enforce_cpu_limit=False,
        enforce_pid_limit=False,
        enforce_outer_seccomp=False,
        enforce_procfs_policy=False,
        local_root_parent=os.environ.get("AGENTIC_TEST_LOCAL_ROOT_PARENT"),
        docker_root_parent=os.environ.get("AGENTIC_TEST_DOCKER_ROOT_PARENT"),
    )
    container_name = backend.container_name
    try:
        startup = backend.start()
        assert startup["ok"] is True, startup
        assert startup["data"]["input_count"] == 1
        assert len(startup["data"]["input_merkle_root"]) == 64
        assert backend.inspect_environment()["ok"] is True

        first = backend.run_ffmpeg({
            "operation": "extract_audio",
            "input": "inputs/input.wav",
            "output": "work/stage.flac",
            "format": "flac",
            "sample_rate": 16000,
            "channels": 1,
            "start_seconds": 0,
            "duration_seconds": 0.5,
        })
        assert first["ok"] is True, first
        workspace = backend.inspect_workspace()
        assert workspace["ok"] is True, workspace
        assert any(
            item["path"] == "work/stage.flac"
            for item in workspace["data"]["work"]
        )

        second = backend.run_ffmpeg({
            "operation": "transcode_audio",
            "input": "work/stage.flac",
            "output": "work/report.wav",
            "format": "wav",
            "sample_rate": 16000,
            "channels": 1,
            "start_seconds": 0,
            "duration_seconds": 0.4,
        })
        assert second["ok"] is True, second

        inspection = backend.inspect_artifacts()
        assert inspection["ok"] is True, inspection
        assert inspection["data"]["artifact_count"] >= 2

        finalized = backend.finalize(
            ["report.wav"], "Verified professional audio report"
        )
        assert finalized["ok"] is True, finalized
        result = backend.best_result()
        assert result is not None and result["success"] is True
        assert len(result["files"]) == 1
        assert result["files"][0]["filename"] == "report.wav"
        assert len(result["files"][0]["content"]) > 44
    finally:
        backend.close()

    remaining = subprocess.run(
        ["docker", "ps", "-a", "--filter", f"name=^{container_name}$", "-q"],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
        text=True,
    )
    assert remaining.stdout.strip() == ""


@pytest.mark.integration
@pytest.mark.parametrize("extension", ["xlsx", "docx", "pptx"])
def test_agentic_verifier_renders_primary_office_artifacts(
    tmp_path, extension
):
    image = os.environ.get(
        "AGENTIC_TEST_IMAGE", "gdpval-agentic-sandbox:local"
    )
    if subprocess.run(
        ["docker", "image", "inspect", image],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    ).returncode != 0:
        pytest.skip(f"agentic test image is unavailable: {image}")
    artifact = tmp_path / f"report.{extension}"
    if extension == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.active["A1"] = "Professional report"
        workbook.active["A2"] = 42
        workbook.save(artifact)
    elif extension == "docx":
        from docx import Document

        document = Document()
        document.add_heading("Professional report", 0)
        document.add_paragraph("Verified document content.")
        document.save(artifact)
    else:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Professional report"
        presentation.save(artifact)

    backend = AgenticDockerBackend(
        task_prompt=(
            f"Create a professional {extension.upper()} report named "
            f"report.{extension}"
        ),
        reference_files=[],
        occupation="Analyst",
        image=image,
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        require_approved_input_manifest=False,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
        enforce_cpu_limit=False,
        enforce_pid_limit=False,
        enforce_outer_seccomp=False,
        enforce_procfs_policy=False,
        local_root_parent=os.environ.get("AGENTIC_TEST_LOCAL_ROOT_PARENT"),
        docker_root_parent=os.environ.get("AGENTIC_TEST_DOCKER_ROOT_PARENT"),
    )
    try:
        startup = backend.start()
        assert startup["ok"] is True, startup
        copied = subprocess.run(
            [
                "docker",
                "cp",
                str(artifact),
                f"{backend.container_name}:/work/report.{extension}",
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=False,
        )
        assert copied.returncode == 0, copied.stderr.decode(
            "utf-8", errors="replace"
        )
        inspection = backend.inspect_artifacts(180)
        assert inspection["ok"] is True, inspection
        assert len(inspection["data"]["artifacts"]) == 1
        assert inspection["data"]["artifacts"][0]["path"] == (
            f"report.{extension}"
        )
        assert inspection["data"]["artifacts"][0]["openable"] is True
        finalized = backend.finalize(
            [f"report.{extension}"], "Rendered primary artifact", 180
        )
        assert finalized["ok"] is True, finalized
        assert backend.best_result()["success"] is True
    finally:
        backend.close()


@pytest.mark.integration
def test_agentic_docker_generated_python_denies_exec_and_network(tmp_path):
    image = os.environ.get(
        "AGENTIC_TEST_IMAGE", "gdpval-agentic-sandbox:local"
    )
    source = tmp_path / "input.txt"
    source.write_text("approved input", encoding="utf-8")
    backend = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[str(source)],
        occupation="Analyst",
        image=image,
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        require_approved_input_manifest=False,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
        enforce_cpu_limit=False,
        enforce_pid_limit=False,
        enforce_outer_seccomp=False,
        enforce_procfs_policy=False,
        local_root_parent=os.environ.get("AGENTIC_TEST_LOCAL_ROOT_PARENT"),
        docker_root_parent=os.environ.get("AGENTIC_TEST_DOCKER_ROOT_PARENT"),
    )
    try:
        startup = backend.start()
        assert startup["ok"] is True, startup
        result = backend.run_python(
            """
from pathlib import Path
import socket
import subprocess
blocked = 0
for action in (
    lambda: socket.socket(socket.AF_INET, socket.SOCK_STREAM),
    lambda: subprocess.run(['/usr/bin/ffprobe', '-version']),
):
    try:
        action()
    except OSError:
        blocked += 1
assert blocked == 2, blocked
Path('report.txt').write_text('blocked')
""",
            30,
        )
        if "seccomp TSYNC unavailable" in result.get("data", {}).get(
            "stderr_tail", ""
        ):
            pytest.skip(
                "local kernel lacks seccomp TSYNC; generated Python stayed fail-closed"
            )
        assert result["ok"] is True, result
    finally:
        backend.close()


@pytest.mark.integration
def test_outer_seccomp_allows_inner_filter_and_blocks_raw_signal_syscalls(
    tmp_path
):
    image = os.environ.get(
        "AGENTIC_TEST_IMAGE", "gdpval-agentic-sandbox:local"
    )
    profile_probe = subprocess.run(
        [
            "docker", "run", "--rm", "--network", "none", "--ipc", "none",
            "--read-only", "--cap-drop", "ALL", "--user", "65532:65532",
            "--security-opt", "no-new-privileges",
            "--security-opt", f"seccomp={SECCOMP}",
            "--entrypoint", "python", image, "-I", "-B", "-c",
            "print('outer-seccomp-ready')",
        ],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    if b"seccomp profiles are not supported" in profile_probe.stderr:
        pytest.skip("local Docker daemon does not support custom seccomp profiles")
    assert profile_probe.returncode == 0, profile_probe.stderr.decode(
        "utf-8", errors="replace"
    )
    backend = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[],
        occupation="Analyst",
        image=image,
        seccomp_profile=str(SECCOMP),
        allow_unpinned_image=True,
        require_rootless_or_userns=False,
        require_approved_input_manifest=False,
        require_supply_chain_identity=False,
        require_dedicated_host=False,
        enforce_cpu_limit=False,
        enforce_pid_limit=False,
        enforce_outer_seccomp=True,
        enforce_procfs_policy=False,
        local_root_parent=os.environ.get("AGENTIC_TEST_LOCAL_ROOT_PARENT"),
        docker_root_parent=os.environ.get("AGENTIC_TEST_DOCKER_ROOT_PARENT"),
    )
    try:
        startup = backend.start()
        assert startup["ok"] is True, startup
        result = backend.run_python(
            """
import ctypes
import errno
from pathlib import Path

libc = ctypes.CDLL(None, use_errno=True)
blocked = 0
for number, arguments in (
    (317, (0, 0, 0)),
    (129, (1, 0, 0)),
    (297, (1, 1, 0, 0)),
):
    ctypes.set_errno(0)
    returned = libc.syscall(number, *arguments)
    if returned == -1 and ctypes.get_errno() == errno.EPERM:
        blocked += 1
assert blocked == 3, blocked
Path('report.txt').write_text('blocked')
""",
            30,
        )
        if "seccomp TSYNC unavailable" in result.get("data", {}).get(
            "stderr_tail", ""
        ):
            pytest.skip("local nested kernel lacks seccomp TSYNC")
        assert result["ok"] is True, result
    finally:
        backend.close()


@pytest.mark.integration
def test_agentic_production_runner_preflight(tmp_path):
    if os.environ.get("AGENTIC_PRODUCTION_PREFLIGHT") != "1":
        pytest.skip("dedicated production runner preflight not requested")
    image = os.environ["AGENTIC_PRODUCTION_IMAGE"]
    sbom_sha256 = os.environ["AGENTIC_PRODUCTION_SBOM_SHA256"]
    apparmor_profile = os.environ["AGENTIC_PRODUCTION_APPARMOR_PROFILE"]
    if "@sha256:" not in image:
        pytest.fail("production image must be digest-pinned")
    source = tmp_path / "source.txt"
    source.write_text("approved public fixture", encoding="utf-8")

    probe = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[str(source)],
        occupation="Analyst",
        image=image,
        verifier_image=image,
        apparmor_profile=apparmor_profile,
        sbom_sha256=sbom_sha256,
        allow_unpinned_image=False,
        require_rootless_or_userns=True,
        require_approved_input_manifest=False,
        require_supply_chain_identity=True,
        require_dedicated_host=True,
    )
    try:
        probe._stage_inputs()
        record = probe.input_records[0]
        expected_root = probe.input_merkle_root
        approved = {
            record["path"]: {
                key: record[key]
                for key in (
                    "path", "type", "link_count", "size_bytes",
                    "source_allocated_bytes", "sha256",
                    "provider_classification",
                )
            }
        }
    finally:
        probe.close()

    backend = AgenticDockerBackend(
        task_prompt="Create report.txt",
        reference_files=[{
            "source_path": str(source),
            "relative_path": "source.txt",
        }],
        occupation="Analyst",
        image=image,
        verifier_image=image,
        apparmor_profile=apparmor_profile,
        sbom_sha256=sbom_sha256,
        approved_input_manifest=approved,
        expected_input_merkle_root=expected_root,
    )
    container_name = backend.container_name
    volume_name = backend.work_volume_name
    try:
        startup = backend.start()
        assert startup["ok"] is True, startup
        result = backend.run_python(
            """
from pathlib import Path
import socket
import subprocess
blocked = 0
for action in (
    lambda: socket.socket(socket.AF_INET, socket.SOCK_STREAM),
    lambda: subprocess.run(['/usr/bin/ffprobe', '-version']),
):
    try:
        action()
    except OSError:
        blocked += 1
assert blocked == 2, blocked
Path('report.txt').write_text('verified production preflight')
""",
            30,
        )
        assert result["ok"] is True, result
        inspection = backend.inspect_artifacts()
        assert inspection["ok"] is True, inspection
        finalized = backend.finalize(
            ["report.txt"], "production preflight"
        )
        assert finalized["ok"] is True, finalized
    finally:
        backend.close()

    assert subprocess.run(
        ["docker", "container", "inspect", container_name],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    ).returncode != 0
    assert subprocess.run(
        ["docker", "volume", "inspect", volume_name],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        check=False,
    ).returncode != 0