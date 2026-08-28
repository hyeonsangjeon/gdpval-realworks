"""Tests for ``core.tools.read_deliverable`` (PR2 task 201)."""

from __future__ import annotations

import ast
import base64
import importlib
import io
import os
from pathlib import Path

import pytest

from core.tools import (
    get_renderer_fingerprint,
    MODEL_READ_DELIVERABLE_OPS,
    MODEL_READ_DELIVERABLE_TOOL_SCHEMA,
    READ_DELIVERABLE_OPS,
    READ_DELIVERABLE_TOOL_SCHEMA,
    RendererDependencyError,
    read_deliverable,
)

read_deliverable_module = importlib.import_module("core.tools.read_deliverable")


# ── Fixtures ─────────────────────────────────────────────────────────


@pytest.fixture
def base_dir(tmp_path: Path) -> Path:
    return tmp_path


@pytest.fixture
def xlsx_file(base_dir: Path) -> Path:
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "Item"
    ws["B1"] = "Value"
    ws["A2"] = "alpha"
    ws["B2"] = 1
    ws["A3"] = "beta"
    ws["B3"] = 2
    # Add a bit of formatting so inspect_formatting has something to see.
    from openpyxl.styles import Font, PatternFill, Border, Side
    ws["A1"].font = Font(bold=True)
    ws["A1"].fill = PatternFill("solid", fgColor="FFFF00")
    ws["A1"].border = Border(bottom=Side(style="thin"))
    ws.merge_cells("A1:B1")
    p = base_dir / "report.xlsx"
    wb.save(p)
    return p


@pytest.fixture
def docx_file(base_dir: Path) -> Path:
    pytest.importorskip("docx")
    from docx import Document
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("First paragraph body.")
    doc.add_paragraph("Second paragraph body.")
    p = base_dir / "memo.docx"
    doc.save(p)
    return p


@pytest.fixture
def pdf_file(base_dir: Path) -> Path:
    """Generate a tiny 1-page PDF via reportlab (already in requirements)."""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    p = base_dir / "doc.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 750, "Hello PDF world.")
    c.showPage()
    c.save()
    return p


@pytest.fixture
def pptx_file(base_dir: Path) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation

    presentation = Presentation()
    for title in ("Overview", "Details"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
    p = base_dir / "slides.pptx"
    presentation.save(p)
    return p


@pytest.fixture
def png_file(base_dir: Path) -> Path:
    pytest.importorskip("PIL")
    from PIL import Image
    p = base_dir / "img.png"
    Image.new("RGB", (16, 16), color="red").save(p)
    return p


@pytest.fixture
def txt_file(base_dir: Path) -> Path:
    p = base_dir / "notes.txt"
    p.write_text("hello world\nline 2\n", encoding="utf-8")
    return p


def _fake_convert_to_pdf(source: Path, out_dir: Path, pages: int = 2) -> Path:
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    output = out_dir / f"{source.stem}.pdf"
    document = canvas.Canvas(str(output))
    for page in range(1, pages + 1):
        document.drawString(100, 750, f"Rendered page {page}")
        document.showPage()
    document.save()
    return output


# ── Schema / surface ─────────────────────────────────────────────────


def test_ops_constant_matches_schema_enum():
    enum = READ_DELIVERABLE_TOOL_SCHEMA["parameters"]["properties"]["op"]["enum"]
    assert tuple(enum) == READ_DELIVERABLE_OPS
    assert len(READ_DELIVERABLE_OPS) == 6


def test_model_schema_excludes_harness_only_render_op():
    enum = MODEL_READ_DELIVERABLE_TOOL_SCHEMA[
        "parameters"
    ]["properties"]["op"]["enum"]
    assert tuple(enum) == MODEL_READ_DELIVERABLE_OPS
    assert "render_to_image" in READ_DELIVERABLE_OPS
    assert "render_to_image" not in MODEL_READ_DELIVERABLE_OPS
    assert "base64" not in str(MODEL_READ_DELIVERABLE_TOOL_SCHEMA).lower()


# ── Path safety ──────────────────────────────────────────────────────


def test_rejects_traversal_outside_base(base_dir, txt_file):
    # Create a sibling outside the base
    outside = base_dir.parent / "secret.txt"
    outside.write_text("nope")
    result = read_deliverable(
        "read_content",
        f"../{outside.name}",
        base_dir=str(base_dir),
    )
    assert result["ok"] is False
    assert result["error_type"] == "bad_path"


def test_rejects_absolute_outside_base(base_dir):
    result = read_deliverable(
        "read_content",
        "/etc/passwd",
        base_dir=str(base_dir),
    )
    assert result["ok"] is False
    assert result["error_type"] == "bad_path"


def test_rejects_missing_file(base_dir):
    result = read_deliverable(
        "read_content",
        "no_such_file.txt",
        base_dir=str(base_dir),
    )
    assert result["ok"] is False
    assert result["error_type"] == "bad_path"


def test_rejects_unknown_op(base_dir, txt_file):
    result = read_deliverable(
        "rm_rf",
        txt_file.name,
        base_dir=str(base_dir),
    )
    assert result["ok"] is False
    assert result["error_type"] == "bad_op"


# ── inspect_structure ────────────────────────────────────────────────


def test_inspect_structure_txt(base_dir, txt_file):
    r = read_deliverable("inspect_structure", txt_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["kind"] == "txt"
    assert r["data"]["size_bytes"] > 0


def test_inspect_structure_xlsx(base_dir, xlsx_file):
    r = read_deliverable("inspect_structure", xlsx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    d = r["data"]
    assert d["kind"] == "xlsx"
    assert d["sheet_count"] == 1
    assert d["sheets"][0]["name"] == "Summary"


def test_inspect_structure_docx(base_dir, docx_file):
    r = read_deliverable("inspect_structure", docx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["kind"] == "docx"
    assert r["data"]["paragraph_count"] >= 2


def test_inspect_structure_pdf(base_dir, pdf_file):
    r = read_deliverable("inspect_structure", pdf_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["kind"] == "pdf"
    assert r["data"]["page_count"] == 1


# ── read_content ─────────────────────────────────────────────────────


def test_read_content_txt(base_dir, txt_file):
    r = read_deliverable("read_content", txt_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "hello world" in r["data"]["text"]


def test_read_content_xlsx_full(base_dir, xlsx_file):
    r = read_deliverable("read_content", xlsx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "alpha" in r["data"]["text"]
    assert "beta" in r["data"]["text"]


def test_read_content_xlsx_scope_sheet(base_dir, xlsx_file):
    r = read_deliverable(
        "read_content",
        xlsx_file.name,
        base_dir=str(base_dir),
        scope={"sheet": "Summary"},
    )
    assert r["ok"] is True
    assert "Summary" in r["data"]["text"]


def test_read_content_pdf(base_dir, pdf_file):
    r = read_deliverable("read_content", pdf_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "Hello PDF" in r["data"]["text"]


def test_read_content_docx(base_dir, docx_file):
    r = read_deliverable("read_content", docx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "First paragraph" in r["data"]["text"]


# ── read_content: a deck is more than its text frames ────────────────
#
# `_read_pptx_text` used to take `shape.has_text_frame` and nothing else, so a
# table was a graphic frame it walked straight past and a group was one shape
# whose children it never entered. Both are where real decks keep their
# content: a gold answer in the stage-1 corpus, `WorkStudy.pptx`, read as 186
# characters of slide titles while the fifteen-row table holding every activity
# category and percentage -- which a rubric item then asked about, and scored
# 0/1 -- was never shown to the judge at all.
#
# Each fixture below puts its content *only* in the shape under test, so a
# reader that regresses to text frames alone fails rather than passing on the
# title it can still see.


@pytest.fixture
def pptx_with_table(base_dir: Path) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Activity Analysis"
    table = slide.shapes.add_table(
        3, 2, Inches(1), Inches(2), Inches(6), Inches(2)
    ).table
    table.cell(0, 0).text = "Activity"
    table.cell(0, 1).text = "Share"
    table.cell(1, 0).text = "Material handling"
    table.cell(1, 1).text = "31.4%"
    table.cell(2, 0).text = "Machine setup"
    table.cell(2, 1).text = "12.7%"
    p = base_dir / "study.pptx"
    presentation.save(p)
    return p


@pytest.fixture
def pptx_with_group(base_dir: Path) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Findings"
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(
        Inches(1), Inches(2), Inches(4), Inches(1)
    )
    box.text_frame.text = "Grouped callout text"
    p = base_dir / "grouped.pptx"
    presentation.save(p)
    return p


@pytest.fixture
def pptx_with_chart(base_dir: Path) -> Path:
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Distribution"
    chart_data = CategoryChartData()
    chart_data.categories = ["Inspection", "Rework"]
    chart_data.add_series("Share", (0.62, 0.38))
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(2), Inches(6), Inches(4), chart_data
    )
    p = base_dir / "chart.pptx"
    presentation.save(p)
    return p


def test_read_content_pptx_reads_text_frames(base_dir, pptx_file):
    """The behaviour that already worked has to keep working."""
    r = read_deliverable("read_content", pptx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "Overview" in r["data"]["text"]
    assert "Details" in r["data"]["text"]


def test_read_content_pptx_reads_table_cells(base_dir, pptx_with_table):
    """The defect this fixes: cells live on a shape with no text frame."""
    r = read_deliverable(
        "read_content", pptx_with_table.name, base_dir=str(base_dir)
    )

    assert r["ok"] is True
    text = r["data"]["text"]
    assert "Material handling" in text
    assert "31.4%" in text
    assert "Machine setup" in text


def test_read_content_pptx_keeps_table_rows_together(base_dir, pptx_with_table):
    """A row read cell-by-cell loses which share belongs to which activity.

    Marked and joined the same way ``_read_docx_text`` already marks tables, so
    a judge meets one shape of table across both formats.
    """
    r = read_deliverable(
        "read_content", pptx_with_table.name, base_dir=str(base_dir)
    )

    text = r["data"]["text"]
    assert "[Table]" in text
    assert "Material handling | 31.4%" in text


def test_read_content_pptx_reads_inside_a_group(base_dir, pptx_with_group):
    """A group is one shape; its children were never visited."""
    r = read_deliverable(
        "read_content", pptx_with_group.name, base_dir=str(base_dir)
    )

    assert r["ok"] is True
    assert "Grouped callout text" in r["data"]["text"]


def test_read_content_pptx_reads_chart_categories(base_dir, pptx_with_chart):
    """Rubrics ask which categories a chart shows, which a picture cannot say."""
    r = read_deliverable(
        "read_content", pptx_with_chart.name, base_dir=str(base_dir)
    )

    assert r["ok"] is True
    text = r["data"]["text"]
    assert "Inspection" in text
    assert "Rework" in text


@pytest.fixture
def pptx_with_unmodelled_chart(base_dir: Path, pptx_with_chart: Path) -> Path:
    """A plot type python-pptx will not model, built from one it will.

    ``pie3DChart`` is not exotic -- four of the five charts in one stage-1 gold
    answer are that -- but python-pptx raises ``unsupported plot type`` on it,
    and there is no way to author one through the library. Retagging the plot
    element of a chart it *can* write produces the same file Excel would, and
    reproduces the failure exactly: chart_type, categories and series all raise.
    """
    pytest.importorskip("pptx")
    from pptx import Presentation

    namespace = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    presentation = Presentation(str(pptx_with_chart))
    shape = next(
        s
        for s in presentation.slides[0].shapes
        if getattr(s, "has_chart", False)
    )
    plot = shape.chart._chartSpace.find(".//c:plotArea/c:pieChart", namespace)
    plot.tag = "{%s}pie3DChart" % namespace["c"]
    p = base_dir / "pie3d.pptx"
    presentation.save(p)
    return p


def test_the_modelled_chart_accessors_really_do_refuse(pptx_with_unmodelled_chart):
    """The fixture is only worth anything if it reproduces the real failure.

    Without this, a python-pptx release that starts modelling 3-D pies would
    leave the fallback untested and the next test passing down the easy path.
    """
    from pptx import Presentation

    shape = next(
        s
        for s in Presentation(str(pptx_with_unmodelled_chart)).slides[0].shapes
        if getattr(s, "has_chart", False)
    )

    with pytest.raises(ValueError, match="unsupported plot type"):
        shape.chart.chart_type


def test_read_content_pptx_reads_a_chart_python_pptx_cannot_model(
    base_dir, pptx_with_unmodelled_chart
):
    """The cached points are in the XML whatever the plot type is."""
    r = read_deliverable(
        "read_content", pptx_with_unmodelled_chart.name, base_dir=str(base_dir)
    )

    assert r["ok"] is True
    text = r["data"]["text"]
    assert "pie3DChart" in text, "the chart type is itself a rubric fact"
    assert "Inspection" in text
    assert "Rework" in text
    assert "0.62" in text


def test_a_chart_that_cannot_describe_itself_does_not_break_the_read(
    base_dir, pptx_with_table, monkeypatch
):
    """One unreadable chart must cost its own text, not the whole deck's.

    python-pptx raises on chart types with no category axis, and this reader
    runs against arbitrary gold answers -- so the guard is the difference
    between losing a chart and losing every slide behind it.
    """
    class Exploding:
        has_text_frame = False
        has_table = False
        has_chart = True

        @property
        def chart(self):
            raise ValueError("no category axis on this chart type")

    assert read_deliverable_module._pptx_shape_text(Exploding()) == []

    # And the surrounding deck still reads.
    r = read_deliverable(
        "read_content", pptx_with_table.name, base_dir=str(base_dir)
    )
    assert r["ok"] is True
    assert "Material handling" in r["data"]["text"]


def test_read_content_truncation_flag(base_dir):
    big = base_dir / "big.txt"
    big.write_text("x" * (300_000))
    r = read_deliverable("read_content", big.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["truncated"] is True
    assert r["data"]["char_count"] <= 200_000


# ── inspect_formatting ───────────────────────────────────────────────


def test_inspect_formatting_xlsx(base_dir, xlsx_file):
    r = read_deliverable("inspect_formatting", xlsx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    d = r["data"]
    sheet = d["sheets"][0]
    assert sheet["merged_ranges"]
    assert sheet["styled_cells_count"] >= 1


def test_inspect_formatting_xlsx_preserves_color_types_without_descriptor_junk(
    base_dir,
):
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.styles import Color, Font, PatternFill

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    colors = {
        "A1": (Color(rgb="FF112233"), Color(rgb="FF445566")),
        "A2": (Color(theme=4, tint=0.4), Color(theme=5, tint=-0.25)),
        "A3": (Color(indexed=7), Color(indexed=8)),
        "A4": (Color(auto=True), Color(auto=True)),
    }
    for ref, (font_color, fill_color) in colors.items():
        sheet[ref] = ref
        sheet[ref].font = Font(color=font_color)
        sheet[ref].fill = PatternFill(fill_type="solid", fgColor=fill_color)
    sheet["A5"] = "plain"
    sheet["A6"] = 1
    sheet["A6"].number_format = "0.00"

    path = base_dir / "colors.xlsx"
    workbook.save(path)

    result = read_deliverable("inspect_formatting", path.name, base_dir=str(base_dir))

    assert result["ok"] is True
    styled = {
        cell["ref"]: cell
        for cell in result["data"]["sheets"][0]["styled_cells_sample"]
    }
    assert styled["A1"]["font_color"] == "FF112233"
    assert styled["A1"]["fill"] == "FF445566"
    assert styled["A2"]["font_color"] == "theme:4:tint:0.4"
    assert styled["A2"]["fill"] == "theme:5:tint:-0.25"
    assert styled["A3"]["font_color"] == "indexed:7"
    assert styled["A3"]["fill"] == "indexed:8"
    assert styled["A4"]["font_color"] == "auto"
    assert styled["A4"]["fill"] == "auto"
    assert set(styled) == set(colors)
    assert result["data"]["sheets"][0]["styled_cells_count"] == len(colors)
    assert "Values must be of type" not in str(result)


def _rgb_reads_in_core() -> list[tuple[str, str, int]]:
    """Every ``.rgb`` attribute read under ``core/``, with its enclosing function."""
    core_root = Path(read_deliverable_module.__file__).resolve().parents[1]
    reads: list[tuple[str, str, int]] = []
    for path in sorted(core_root.rglob("*.py")):
        tree = ast.parse(path.read_text(encoding="utf-8"))
        functions = [
            node for node in ast.walk(tree)
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef))
        ]
        for node in ast.walk(tree):
            if not (isinstance(node, ast.Attribute) and node.attr == "rgb"):
                continue
            enclosing = [
                fn for fn in functions
                if fn.lineno <= node.lineno <= (fn.end_lineno or fn.lineno)
            ]
            innermost = max(enclosing, key=lambda fn: fn.lineno).name if enclosing else "<module>"
            reads.append(
                (path.relative_to(core_root).as_posix(), innermost, node.lineno)
            )
    return reads


def test_openpyxl_color_rgb_is_only_read_behind_the_type_guard():
    """No module may read ``Color.rgb`` outside the type-dispatched helper.

    #156 fixed the read site that put openpyxl's validation message into judge
    evidence, but the fix is structural -- dispatch on ``color.type`` first --
    so it holds only as long as nothing else reads the descriptor blind. The
    behavioural test above cannot see a second offender in another module.
    """
    pytest.importorskip("openpyxl")
    from openpyxl.styles import Color

    # Why the rule is worth enforcing: on a non-RGB colour the descriptor
    # returns *itself* rather than raising, and its repr is the validation
    # message. It is not JSON-serialisable, so it only becomes visible once
    # something stringifies the evidence -- by which point it reads like an
    # observation. If a future openpyxl raises here instead, relax the rule.
    leaked = Color(theme=4).rgb
    assert not isinstance(leaked, str)
    assert "Values must be of type" in repr(leaked)

    assert [(module, function) for module, function, _ in _rgb_reads_in_core()] == [
        ("tools/read_deliverable.py", "_safe_cell_color")
    ], _rgb_reads_in_core()


def test_default_font_color_lookup_fails_soft_when_openpyxl_internals_change():
    class WorkbookWithoutStyleInternals:
        pass

    class CellWithoutStyleInternals:
        font = type("Font", (), {"color": "explicit"})()

    assert (
        read_deliverable_module._workbook_default_font_color(
            WorkbookWithoutStyleInternals()
        )
        is None
    )
    assert (
        read_deliverable_module._nondefault_font_color(
            CellWithoutStyleInternals(), None
        )
        is None
    )


def test_inspect_formatting_docx(base_dir, docx_file):
    r = read_deliverable("inspect_formatting", docx_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["kind"] == "docx"
    assert "style_histogram" in r["data"]


# ── render_to_image ──────────────────────────────────────────────────


def test_render_to_image_pdf(base_dir, pdf_file, monkeypatch):
    pytest.importorskip("fitz")
    monkeypatch.setattr(
        read_deliverable_module, "_pymupdf_runtime_version", lambda: "9.9.9"
    )
    r = read_deliverable(
        "render_to_image", pdf_file.name,
        base_dir=str(base_dir), scope={"page": 1},
    )
    assert r["ok"] is True
    payload = r["data"]
    assert payload["kind"] == "image_png_base64"
    raw = base64.b64decode(payload["base64"])
    assert raw.startswith(b"\x89PNG")
    assert payload["byte_size"] <= 5 * 1024 * 1024
    assert payload["scope"] == {"page": 1}
    assert payload["source_page_count"] == 1
    assert payload["renderer"]["pymupdf_version"] == "9.9.9"


def test_render_to_image_png(base_dir, png_file):
    r = read_deliverable(
        "render_to_image", png_file.name,
        base_dir=str(base_dir),
    )
    assert r["ok"] is True
    raw = base64.b64decode(r["data"]["base64"])
    assert raw.startswith(b"\x89PNG")


def test_render_to_image_xlsx_uses_original_path_and_is_read_only(
    base_dir, xlsx_file, monkeypatch
):
    openpyxl = pytest.importorskip("openpyxl")
    workbook = openpyxl.load_workbook(xlsx_file)
    workbook.create_sheet("Details")["A1"] = "second sheet"
    workbook.save(xlsx_file)
    workbook.close()
    source_before = xlsx_file.read_bytes()
    observed = {}

    def fake_convert(source: Path, out_dir: Path) -> Path:
        observed["source"] = source
        observed["bytes"] = source.read_bytes()
        return _fake_convert_to_pdf(source, out_dir, pages=2)

    monkeypatch.setattr(
        read_deliverable_module, "_convert_office_to_pdf", fake_convert
    )
    monkeypatch.setattr(
        read_deliverable_module,
        "get_renderer_fingerprint",
        lambda: {
            "libreoffice_binary": "soffice",
            "libreoffice_version": "LibreOffice 24.2.7.2",
            "pymupdf_version": "1.26.3",
        },
    )
    r = read_deliverable(
        "render_to_image", xlsx_file.name,
        base_dir=str(base_dir), scope={"workbook_page": 1},
    )
    assert r["ok"] is True
    payload = r["data"]
    assert payload["source_kind"] == "xlsx"
    assert payload["scope"] == {"workbook_page": 1}
    assert payload["source_sheet_count"] == 2
    assert payload["converted_page_count"] == 2
    assert payload["renderer"]["converter"] == "libreoffice"
    assert payload["renderer"]["libreoffice_binary"] == "soffice"
    assert payload["renderer"]["libreoffice_version"] == "LibreOffice 24.2.7.2"
    assert payload["renderer"]["pymupdf_version"] == "1.26.3"
    assert "/usr/" not in str(payload["renderer"])
    assert base64.b64decode(payload["base64"]).startswith(b"\x89PNG")
    assert observed["source"] == xlsx_file
    assert observed["bytes"] == source_before
    assert xlsx_file.read_bytes() == source_before


def test_render_to_image_pptx_selected_slide(
    base_dir, pptx_file, monkeypatch
):
    monkeypatch.setattr(
        read_deliverable_module, "_convert_office_to_pdf", _fake_convert_to_pdf
    )
    monkeypatch.setattr(
        read_deliverable_module,
        "get_renderer_fingerprint",
        lambda: {
            "libreoffice_binary": "libreoffice",
            "libreoffice_version": "LibreOffice 24.2.7.2",
            "pymupdf_version": "1.26.3",
        },
    )
    r = read_deliverable(
        "render_to_image", pptx_file.name,
        base_dir=str(base_dir), scope={"slide": 2},
    )
    assert r["ok"] is True
    payload = r["data"]
    assert payload["source_kind"] == "pptx"
    assert payload["scope"] == {"slide": 2}
    assert payload["source_slide_count"] == 2
    assert payload["converted_page_count"] == 2
    assert payload["renderer"]["libreoffice_binary"] == "libreoffice"
    assert payload["renderer"]["libreoffice_version"] == "LibreOffice 24.2.7.2"
    assert payload["renderer"]["pymupdf_version"] == "1.26.3"
    assert base64.b64decode(payload["base64"]).startswith(b"\x89PNG")


def test_render_xlsx_missing_libreoffice_is_actionable(
    base_dir, xlsx_file, monkeypatch
):
    monkeypatch.setattr(read_deliverable_module, "_find_soffice", lambda: None)
    r = read_deliverable(
        "render_to_image", xlsx_file.name,
        base_dir=str(base_dir), scope={"workbook_page": 1},
    )
    assert r["ok"] is False
    assert r["error_type"] == "dependency_missing"
    assert "LibreOffice" in r["error"]
    assert "install" in r["error"].lower()


def test_render_pdf_out_of_range(base_dir, pdf_file):
    pytest.importorskip("fitz")
    r = read_deliverable(
        "render_to_image", pdf_file.name,
        base_dir=str(base_dir), scope={"page": 99},
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "out of range" in r["error"]


@pytest.mark.parametrize("page", [0, -1, 1.5, True, "first"])
def test_render_pdf_rejects_invalid_page_scope(base_dir, pdf_file, page):
    r = read_deliverable(
        "render_to_image", pdf_file.name,
        base_dir=str(base_dir), scope={"page": page},
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "positive 1-based integer" in r["error"]


@pytest.mark.parametrize("legacy_scope", [{"sheet": "Summary"}, {"sheet_page": 1}])
def test_render_xlsx_rejects_named_sheet_scope(
    base_dir, xlsx_file, monkeypatch, legacy_scope
):
    monkeypatch.setattr(
        read_deliverable_module,
        "_convert_office_to_pdf",
        lambda *_: pytest.fail("conversion must not run for an invalid sheet"),
    )
    r = read_deliverable(
        "render_to_image", xlsx_file.name,
        base_dir=str(base_dir), scope=legacy_scope,
    )
    assert r["ok"] is False
    assert r["error_type"] == "unsupported_scope"
    assert "workbook_page" in r["error"]


def test_render_xlsx_unknown_key_beats_legacy_scope(
    base_dir, xlsx_file, monkeypatch
):
    monkeypatch.setattr(
        read_deliverable_module,
        "_convert_office_to_pdf",
        lambda *_: pytest.fail("conversion must not run for invalid scope"),
    )
    result = read_deliverable(
        "render_to_image", xlsx_file.name, base_dir=str(base_dir),
        scope={"sheet": "Summary", "bogus": 1},
    )
    assert result["ok"] is False
    assert result["error_type"] == "bad_scope"
    assert "bogus" in result["error"]


def test_render_xlsx_rejects_workbook_page_after_first(
    base_dir, xlsx_file, monkeypatch
):
    monkeypatch.setattr(
        read_deliverable_module,
        "_convert_office_to_pdf",
        lambda *_: pytest.fail("conversion must not run for unsupported page"),
    )
    r = read_deliverable(
        "render_to_image", xlsx_file.name,
        base_dir=str(base_dir), scope={"workbook_page": 2},
    )
    assert r["ok"] is False
    assert r["error_type"] == "unsupported_scope"
    assert "only workbook_page=1" in r["error"]


def test_render_pptx_rejects_out_of_range_slide(
    base_dir, pptx_file, monkeypatch
):
    monkeypatch.setattr(
        read_deliverable_module,
        "_convert_office_to_pdf",
        lambda *_: pytest.fail("conversion must not run for an invalid slide"),
    )
    r = read_deliverable(
        "render_to_image", pptx_file.name,
        base_dir=str(base_dir), scope={"slide": 3},
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "out of range" in r["error"]


@pytest.mark.parametrize(
    "fixture_name,scope",
    [
        ("pdf_file", {"slide": 1}),
        ("pptx_file", {"page": 1}),
        ("xlsx_file", {"page": 1}),
        ("png_file", {"page": 1}),
    ],
)
def test_render_rejects_unknown_scope_keys(
    request, base_dir, fixture_name, scope
):
    path = request.getfixturevalue(fixture_name)
    r = read_deliverable(
        "render_to_image", path.name, base_dir=str(base_dir), scope=scope,
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "unknown scope keys" in r["error"]


def test_render_rejects_non_object_scope(base_dir, pdf_file):
    r = read_deliverable(
        "render_to_image", pdf_file.name,
        base_dir=str(base_dir), scope=[1],  # type: ignore[arg-type]
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"


def test_render_to_image_still_rejects_unsupported_text(base_dir, txt_file):
    r = read_deliverable(
        "render_to_image", txt_file.name, base_dir=str(base_dir),
    )
    assert r["ok"] is False
    assert r["error_type"] == "unsupported_scope"
    assert "not supported for kind=txt" in r["error"]


def test_render_to_image_docx_page(base_dir, docx_file, monkeypatch):
    """docx renders through the same LibreOffice->PDF->PyMuPDF path.

    This used to assert the opposite. The judge's visual route could not
    render a document deliverable at all, so every visual rubric item on a
    docx-only task failed the run with
    `required_visual_render_target_unavailable` -- 73 items across 24 tasks of
    the sol-220 grading run, none of which ever reached a model.
    """
    source_before = docx_file.read_bytes()
    observed = {}

    def fake_convert(source: Path, out_dir: Path) -> Path:
        observed["source"] = source
        observed["bytes"] = source.read_bytes()
        return _fake_convert_to_pdf(source, out_dir, pages=2)

    monkeypatch.setattr(
        read_deliverable_module, "_convert_office_to_pdf", fake_convert
    )
    monkeypatch.setattr(
        read_deliverable_module,
        "get_renderer_fingerprint",
        lambda: {
            "libreoffice_binary": "soffice",
            "libreoffice_version": "LibreOffice 24.2.7.2",
            "pymupdf_version": "1.26.3",
        },
    )
    r = read_deliverable(
        "render_to_image", docx_file.name,
        base_dir=str(base_dir), scope={"page": 2},
    )
    assert r["ok"] is True
    payload = r["data"]
    assert payload["source_kind"] == "docx"
    assert payload["scope"] == {"page": 2}
    assert payload["converted_page_count"] == 2
    # A .docx carries no pagination of its own, so there is no source page
    # count to report -- only what the conversion produced.
    assert "source_page_count" not in payload
    assert payload["renderer"]["converter"] == "libreoffice"
    assert payload["renderer"]["rasterizer"] == "pymupdf"
    assert payload["renderer"]["libreoffice_binary"] == "soffice"
    assert payload["renderer"]["libreoffice_version"] == "LibreOffice 24.2.7.2"
    assert payload["renderer"]["pymupdf_version"] == "1.26.3"
    assert base64.b64decode(payload["base64"]).startswith(b"\x89PNG")
    # The original is handed to the converter untouched, as for xlsx.
    assert observed["source"] == docx_file
    assert observed["bytes"] == source_before
    assert docx_file.read_bytes() == source_before


def test_render_to_image_docx_defaults_to_first_page(
    base_dir, docx_file, monkeypatch
):
    monkeypatch.setattr(
        read_deliverable_module, "_convert_office_to_pdf", _fake_convert_to_pdf
    )
    monkeypatch.setattr(
        read_deliverable_module,
        "get_renderer_fingerprint",
        lambda: {
            "libreoffice_binary": "soffice",
            "libreoffice_version": "LibreOffice 24.2.7.2",
            "pymupdf_version": "1.26.3",
        },
    )
    r = read_deliverable(
        "render_to_image", docx_file.name, base_dir=str(base_dir),
    )
    assert r["ok"] is True
    assert r["data"]["scope"] == {"page": 1}


def test_render_docx_out_of_range_page_is_invalid_scope(
    base_dir, docx_file, monkeypatch
):
    """Range is checked against the conversion, the only page count there is."""
    monkeypatch.setattr(
        read_deliverable_module, "_convert_office_to_pdf", _fake_convert_to_pdf
    )
    monkeypatch.setattr(
        read_deliverable_module,
        "get_renderer_fingerprint",
        lambda: {
            "libreoffice_binary": "soffice",
            "libreoffice_version": "LibreOffice 24.2.7.2",
            "pymupdf_version": "1.26.3",
        },
    )
    r = read_deliverable(
        "render_to_image", docx_file.name,
        base_dir=str(base_dir), scope={"page": 9},
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "page 9 out of range 1..2" in r["error"]


def test_render_docx_rejects_unknown_scope_keys(base_dir, docx_file):
    r = read_deliverable(
        "render_to_image", docx_file.name,
        base_dir=str(base_dir), scope={"slide": 1},
    )
    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "unknown scope keys for docx" in r["error"]


def test_render_docx_missing_libreoffice_is_actionable(
    base_dir, docx_file, monkeypatch
):
    """The exact shape of the gap this fix closes: no Writer, clear error."""
    monkeypatch.setattr(read_deliverable_module, "_find_soffice", lambda: None)
    r = read_deliverable(
        "render_to_image", docx_file.name,
        base_dir=str(base_dir), scope={"page": 1},
    )
    assert r["ok"] is False
    assert r["error_type"] == "dependency_missing"
    assert "LibreOffice" in r["error"]


def test_libreoffice_conversion_uses_isolated_profile_and_allowlisted_env(
    base_dir, xlsx_file, monkeypatch
):
    observed = {}
    monkeypatch.setattr(
        read_deliverable_module, "_find_soffice", lambda: "/usr/bin/soffice"
    )
    monkeypatch.setenv("AZURE_OPENAI_API_KEY", "must-not-leak")
    monkeypatch.setenv("SECRET_TOKEN", "must-not-leak")

    def fake_run(command, **kwargs):
        observed["command"] = command
        observed["env"] = kwargs["env"]
        output = Path(command[command.index("--outdir") + 1]) / "report.pdf"
        output.write_bytes(b"%PDF-fake")
        return type("Completed", (), {"returncode": 0, "stdout": "", "stderr": ""})()

    monkeypatch.setattr(read_deliverable_module.subprocess, "run", fake_run)
    with read_deliverable_module.tempfile.TemporaryDirectory() as temp:
        read_deliverable_module._convert_office_to_pdf(xlsx_file, Path(temp))

    command = observed["command"]
    assert "--headless" in command
    assert "--safe-mode" in command
    assert "--norestore" in command
    assert any(arg.startswith("-env:UserInstallation=file:") for arg in command)
    assert all(
        key in {"PATH", "HOME", "TMPDIR", "LANG"} or key.startswith("LC_")
        for key in observed["env"]
    )
    assert "AZURE_OPENAI_API_KEY" not in observed["env"]
    assert "SECRET_TOKEN" not in observed["env"]


def test_renderer_fingerprint_is_cached_by_selected_executable_and_isolated(
    monkeypatch,
):
    observed = []
    read_deliverable_module._renderer_fingerprint_for_executable.cache_clear()
    monkeypatch.setattr(
        read_deliverable_module,
        "_find_soffice",
        lambda: "/opt/libreoffice/program/soffice",
    )
    monkeypatch.setattr(
        read_deliverable_module, "_pymupdf_runtime_version", lambda: "1.26.3"
    )
    monkeypatch.setenv("AZURE_OPENAI_API_KEY", "must-not-leak")
    monkeypatch.setenv("SECRET_TOKEN", "must-not-leak")

    def fake_run(command, **kwargs):
        observed.append((command, kwargs))
        return type(
            "Completed",
            (),
            {
                "returncode": 0,
                "stdout": "LibreOffice 24.2.7.2 build:abc\nextra\n",
                "stderr": "",
            },
        )()

    monkeypatch.setattr(read_deliverable_module.subprocess, "run", fake_run)

    first = get_renderer_fingerprint()
    second = get_renderer_fingerprint()

    assert first == second == {
        "libreoffice_binary": "soffice",
        "libreoffice_version": "LibreOffice 24.2.7.2 build:abc",
        "pymupdf_version": "1.26.3",
    }
    assert len(observed) == 1
    command, kwargs = observed[0]
    assert command == [
        "/opt/libreoffice/program/soffice", "--headless", "--version"
    ]
    assert kwargs["shell"] is False
    assert kwargs["check"] is False
    assert kwargs["timeout"] <= 10
    assert kwargs["env"]["HOME"] == kwargs["env"]["TMPDIR"]
    assert all(
        key in {"PATH", "HOME", "TMPDIR", "LANG"} or key.startswith("LC_")
        for key in kwargs["env"]
    )
    assert "AZURE_OPENAI_API_KEY" not in kwargs["env"]
    assert "SECRET_TOKEN" not in kwargs["env"]
    assert all("/opt/libreoffice" not in value for value in first.values())


def test_renderer_fingerprint_fails_closed_when_libreoffice_missing(monkeypatch):
    read_deliverable_module._renderer_fingerprint_for_executable.cache_clear()
    monkeypatch.setattr(read_deliverable_module, "_find_soffice", lambda: None)

    with pytest.raises(RendererDependencyError, match="LibreOffice executable"):
        get_renderer_fingerprint()


def test_renderer_fingerprint_fails_closed_on_bad_version_probe(monkeypatch):
    read_deliverable_module._renderer_fingerprint_for_executable.cache_clear()
    monkeypatch.setattr(
        read_deliverable_module, "_find_soffice", lambda: "/usr/bin/soffice"
    )
    monkeypatch.setattr(
        read_deliverable_module.subprocess,
        "run",
        lambda *args, **kwargs: type(
            "Completed",
            (),
            {"returncode": 1, "stdout": "", "stderr": "probe failed"},
        )(),
    )

    with pytest.raises(RendererDependencyError, match="exit 1"):
        get_renderer_fingerprint()


def test_renderer_fingerprint_bounds_version_to_first_line(monkeypatch):
    read_deliverable_module._renderer_fingerprint_for_executable.cache_clear()
    monkeypatch.setattr(
        read_deliverable_module, "_find_soffice", lambda: "/usr/bin/soffice"
    )
    monkeypatch.setattr(
        read_deliverable_module, "_pymupdf_runtime_version", lambda: "1.26.3"
    )
    monkeypatch.setattr(
        read_deliverable_module.subprocess,
        "run",
        lambda *args, **kwargs: type(
            "Completed",
            (),
            {
                "returncode": 0,
                "stdout": ("L" * 500) + "\nsecond line must not appear",
                "stderr": "",
            },
        )(),
    )

    fingerprint = get_renderer_fingerprint()

    assert fingerprint["libreoffice_version"] == "L" * 200
    assert "second line" not in fingerprint["libreoffice_version"]


# ── render_to_image size cap (downsample) ────────────────────────────


def test_render_to_image_size_cap(base_dir):
    PIL = pytest.importorskip("PIL")
    from PIL import Image
    # Generate a noisy PNG that exceeds the 5MB cap pre-downsample.
    big = base_dir / "big.png"
    import os as _os
    img = Image.effect_noise((3000, 3000), 64)
    img = img.convert("RGB")
    img.save(big, format="PNG")
    if _os.path.getsize(big) < 5 * 1024 * 1024:
        pytest.skip("noise PNG fortuitously small; cap path not exercised")
    r = read_deliverable("render_to_image", big.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert r["data"]["byte_size"] <= 5 * 1024 * 1024


# ── probe_audio / probe_video ────────────────────────────────────────


def test_probe_audio_wav(base_dir):
    pytest.importorskip("av")
    import wave
    import struct
    p = base_dir / "tone.wav"
    framerate = 8000
    duration_s = 1
    with wave.open(str(p), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(framerate)
        # ~440Hz square-ish wave so silence_ratio < 1
        samples = bytes()
        import math
        for i in range(framerate * duration_s):
            val = int(10000 * (1 if (i // 9) % 2 else -1))
            samples += struct.pack("<h", val)
        w.writeframes(samples)
    r = read_deliverable("probe_audio", p.name, base_dir=str(base_dir))
    assert r["ok"] is True
    data = r["data"]
    assert data["sample_rate"] == framerate
    assert data["channels"] == 1
    assert data["duration_s"] is not None
    assert 0.5 < data["duration_s"] < 1.5
    assert "peak_amplitude_normalized" in data


def test_probe_audio_rejects_non_audio(base_dir, txt_file):
    r = read_deliverable("probe_audio", txt_file.name, base_dir=str(base_dir))
    assert r["ok"] is True  # graceful, but note
    assert "note" in r["data"]


def test_probe_video_note_for_non_video(base_dir, txt_file):
    r = read_deliverable("probe_video", txt_file.name, base_dir=str(base_dir))
    assert r["ok"] is True
    assert "note" in r["data"]


# ── Archives: a container of readable files is not an unreadable file ─
#
# A `.zip` used to fall through to `kind: "unknown"`, and `read_content`
# answered every question about it with an empty string and the note "binary or
# unsupported for text read". One stage-1 gold answer is five WAV stems in an
# archive; it scored 2.00 of 62, and the single item it passed was "exactly one
# top-level ZIP archive is submitted" -- which passed *because* it is a zip.
# Everything else -- five "in WAV format", five "48,000 Hz exactly", five
# "24-bit PCM or IEEE float", the master's running time -- was asked about a
# file nothing had opened.
#
# The formats inside were never the problem. `probe_audio` has read WAV since
# PR2. What was missing was a way in.


@pytest.fixture
def wav_bytes() -> bytes:
    """A 24-bit 48 kHz stereo WAV, which is what the rubric asks about."""
    import io as _io
    import wave

    buffer = _io.BytesIO()
    with wave.open(buffer, "wb") as sound:
        sound.setnchannels(2)
        sound.setsampwidth(3)
        sound.setframerate(48_000)
        sound.writeframes(b"\x00" * 3 * 2 * 24_000)  # half a second
    return buffer.getvalue()


@pytest.fixture
def zip_file(base_dir: Path, wav_bytes: bytes) -> Path:
    """One real stem and the resource fork macOS files beside it."""
    import zipfile

    p = base_dir / "STEMS.zip"
    with zipfile.ZipFile(p, "w") as writing:
        writing.writestr("STEMS/MASTER.wav", wav_bytes)
        writing.writestr("STEMS/notes.txt", "mixed at 48k, 24-bit\n")
        writing.writestr("__MACOSX/STEMS/._MASTER.wav", b"\x00\x05\x16\x07")
        writing.writestr("STEMS/._notes.txt", b"\x00\x05\x16\x07")
    return p


def test_inspect_structure_lists_an_archive(base_dir, zip_file):
    r = read_deliverable("inspect_structure", zip_file.name, base_dir=str(base_dir))

    assert r["ok"] is True
    data = r["data"]
    assert data["kind"] == "zip"
    assert data["entry_count"] == 2
    assert [entry["name"] for entry in data["entries"]] == [
        "STEMS/MASTER.wav",
        "STEMS/notes.txt",
    ]
    # The kind of each member is what tells a judge which op to reach for.
    assert [entry["kind"] for entry in data["entries"]] == ["audio", "txt"]


def test_an_archive_reads_as_its_own_listing(base_dir, zip_file):
    """"Does it contain a Bass stem in WAV format" is answerable from names.

    A judge that never learns about the member scope still has to be able to
    answer that, so the listing is the text of an archive.
    """
    r = read_deliverable("read_content", zip_file.name, base_dir=str(base_dir))

    assert r["ok"] is True
    assert r["data"]["kind"] == "zip"
    assert "STEMS/MASTER.wav" in r["data"]["text"]
    assert "audio" in r["data"]["text"]


def test_resource_forks_are_hidden_but_counted(base_dir, zip_file):
    """Five real stems and five ``._`` twins reads as ten files of nothing.

    Hiding them is not the same as dropping them: the count says how many were
    withheld, so a listing that looks short can be checked rather than trusted.
    """
    listing = read_deliverable(
        "inspect_structure", zip_file.name, base_dir=str(base_dir)
    )["data"]

    assert "__MACOSX" not in str(listing["entries"])
    assert "._notes.txt" not in str(listing["entries"])
    assert listing["hidden_resource_fork_count"] == 2


def test_a_hidden_member_is_still_reachable_by_name(base_dir, zip_file):
    """Hidden from the listing, not removed from the archive.

    Something has to be able to look at a resource fork to establish that it is
    one, and a reader that can only report what it chose to show cannot be
    argued with.
    """
    r = read_deliverable(
        "read_content", zip_file.name, base_dir=str(base_dir),
        scope={"member": "__MACOSX/STEMS/._MASTER.wav"},
    )

    assert r["ok"] is True


def test_probe_audio_reaches_a_stem_inside_an_archive(base_dir, zip_file):
    """The three rubric items the archive gap cost, on one member.

    Format, sample rate and bit depth are each five items on the real answer --
    thirty of the sixty points it lost -- and all three come off a WAV header
    that was always readable, once something opens the container.
    """
    pytest.importorskip("av")

    r = read_deliverable(
        "probe_audio", zip_file.name, base_dir=str(base_dir),
        scope={"member": "STEMS/MASTER.wav"},
    )

    assert r["ok"] is True, r
    data = r["data"]
    assert data["sample_rate"] == 48_000
    assert data["channels"] == 2
    assert data["codec"] == "pcm_s24le"  # 24-bit PCM, little-endian
    assert 0.4 < data["duration_s"] < 0.6


def test_read_content_reaches_a_text_member(base_dir, zip_file):
    """Any op, not just the audio one: a member is addressed like a file."""
    r = read_deliverable(
        "read_content", zip_file.name, base_dir=str(base_dir),
        scope={"member": "STEMS/notes.txt"},
    )

    assert r["ok"] is True
    assert r["data"]["kind"] == "txt"
    assert "24-bit" in r["data"]["text"]


def test_a_member_keeps_its_extension_while_being_read(base_dir, zip_file):
    """Extraction under a random temp name would make every member unknown.

    Every op here dispatches on suffix, so a member written to disk as
    ``tmpXXXX`` is a member no op can identify -- the container would be open
    and its contents still unreadable.
    """
    with read_deliverable_module._extracted_zip_member(
        zip_file, "STEMS/MASTER.wav"
    ) as extracted:
        assert extracted.suffix == ".wav"
        assert read_deliverable_module._kind_of(extracted) == "audio"
        assert extracted.is_file()

    # And it does not outlive the op that asked for it.
    assert not extracted.exists()


def test_an_unknown_member_is_refused_with_the_names_that_exist(
    base_dir, zip_file
):
    r = read_deliverable(
        "probe_audio", zip_file.name, base_dir=str(base_dir),
        scope={"member": "STEMS/BASS.wav"},
    )

    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "STEMS/MASTER.wav" in r["error"]


def test_a_member_name_cannot_escape_the_archive(base_dir, zip_file):
    """Traversal is refused by lookup, not by sanitising.

    ``getinfo`` matches names exactly, so a member that is not in the archive
    cannot be named -- there is no rule to get wrong and no encoding to slip
    past. This checks the property rather than a blocklist.
    """
    for attempt in ("../../etc/passwd", "/etc/passwd", "STEMS/../../../secret"):
        r = read_deliverable(
            "read_content", zip_file.name, base_dir=str(base_dir),
            scope={"member": attempt},
        )
        assert r["ok"] is False, attempt
        assert r["error_type"] == "bad_scope", attempt


def test_the_member_scope_is_refused_on_a_file_that_is_not_an_archive(
    base_dir, txt_file
):
    """A scope key that silently does nothing is worse than one that errors."""
    r = read_deliverable(
        "read_content", txt_file.name, base_dir=str(base_dir),
        scope={"member": "anything"},
    )

    assert r["ok"] is False
    assert r["error_type"] == "bad_scope"
    assert "zip" in r["error"]


def test_an_oversized_member_is_refused_before_it_is_written(
    base_dir, zip_file, monkeypatch
):
    """The cap bounds temp disk, so it has to be checked before extracting.

    The real archive is 179.7 MB compressed; a hostile or merely broken one is
    the reason this is a number rather than a hope.
    """
    monkeypatch.setattr(read_deliverable_module, "MAX_ZIP_MEMBER_BYTES", 16)

    r = read_deliverable(
        "read_content", zip_file.name, base_dir=str(base_dir),
        scope={"member": "STEMS/notes.txt"},
    )

    assert r["ok"] is False
    assert r["error_type"] == "op_error"
    assert "cap" in r["error"]


def test_a_long_listing_is_cut_and_says_so(base_dir, monkeypatch):
    """A truncated listing that looks complete is a wrong answer, not a short one."""
    import zipfile

    monkeypatch.setattr(read_deliverable_module, "MAX_ZIP_ENTRIES", 3)
    p = base_dir / "many.zip"
    with zipfile.ZipFile(p, "w") as writing:
        for i in range(10):
            writing.writestr(f"file{i}.txt", "x")

    listing = read_deliverable("inspect_structure", p.name, base_dir=str(base_dir))
    text = read_deliverable("read_content", p.name, base_dir=str(base_dir))

    assert listing["data"]["entry_count"] == 3
    assert listing["data"]["truncated"] is True
    assert "truncated" in text["data"]["text"]


def test_a_corrupt_archive_reports_the_failure_instead_of_reading_empty(
    base_dir
):
    """The old behaviour was an empty read with a note, which reads as a weak
    answer rather than a broken file. An archive that will not open should look
    like what it is.
    """
    p = base_dir / "broken.zip"
    p.write_bytes(b"PK\x03\x04not really a zip")

    r = read_deliverable("read_content", p.name, base_dir=str(base_dir))
    structure = read_deliverable("inspect_structure", p.name, base_dir=str(base_dir))

    assert r["ok"] is False
    assert r["error_type"] == "exception"
    # inspect_structure keeps its own contract: it never raises, it annotates.
    assert structure["ok"] is True
    assert "inspection_error" in structure["data"]
