from __future__ import annotations

import base64
import json
import subprocess
import sys
from pathlib import Path

import pytest

from scripts import preflight_grading_renderer as preflight


FINGERPRINT = {
    "libreoffice_binary": "soffice",
    # The full line the probe actually returns, build suffix included, and the
    # value every published grade file carries. A fixture trimmed to the bare
    # version number would silently satisfy a prefix check and stop testing
    # the pin the moment the pin got stricter.
    "libreoffice_version": preflight.EXPECTED_LIBREOFFICE_VERSION,
    "pymupdf_version": "1.26.3",
}
PNG = preflight.PNG_SIGNATURE + b"mock-rendered-png"


def test_direct_script_entrypoint_resolves_core_import():
    batch_root = Path(__file__).resolve().parents[1]
    script = batch_root / "scripts" / "preflight_grading_renderer.py"

    completed = subprocess.run(
        [sys.executable, str(script)],
        cwd=batch_root,
        capture_output=True,
        text=True,
        timeout=30,
        check=False,
    )

    assert completed.returncode in {0, 1}
    assert completed.stderr == ""
    lines = completed.stdout.splitlines()
    assert len(lines) == 1
    payload = json.loads(lines[0])
    assert isinstance(payload.get("ok"), bool)
    assert "No module named 'core'" not in completed.stdout


def _successful_render(source_kind, scope):
    return {
        "ok": True,
        "data": {
            "kind": "image_png_base64",
            "source_kind": source_kind,
            "scope": scope,
            "renderer": {
                "converter": "libreoffice",
                "rasterizer": "pymupdf",
                "dpi": 150,
                **FINGERPRINT,
            },
            "base64": base64.b64encode(PNG).decode("ascii"),
            "byte_size": len(PNG),
        },
    }


def _mock_boundaries(monkeypatch, *, mutate_source=False):
    observed = {"fc_match": [], "renders": []}
    monkeypatch.setattr(
        preflight, "get_renderer_fingerprint", lambda: dict(FINGERPRINT)
    )
    monkeypatch.setattr(
        preflight.shutil,
        "which",
        lambda executable: "/usr/bin/fc-match" if executable == "fc-match" else None,
    )
    monkeypatch.setenv("AZURE_OPENAI_API_KEY", "must-not-leak")
    monkeypatch.setenv("SECRET_TOKEN", "must-not-leak")

    def fake_run(command, **kwargs):
        observed["fc_match"].append((command, kwargs))
        return type(
            "Completed",
            (),
            {"returncode": 0, "stdout": "Liberation Sans\n", "stderr": ""},
        )()

    def fake_read(op, path, *, base_dir, scope):
        observed["renders"].append((op, path, base_dir, scope))
        source_path = Path(base_dir) / path
        assert source_path.is_file()
        if path.endswith(".xlsx"):
            from openpyxl import load_workbook

            workbook = load_workbook(source_path, read_only=False)
            try:
                assert workbook.active["A1"].font.name == preflight.FONT_FAMILY
            finally:
                workbook.close()
            source_kind = "xlsx"
        elif path.endswith(".pptx"):
            from pptx import Presentation

            presentation = Presentation(str(source_path))
            assert len(presentation.slides) == 2
            for slide in presentation.slides:
                runs = [
                    run
                    for shape in slide.shapes
                    if shape.has_text_frame
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                ]
                assert runs and all(
                    run.font.name == preflight.FONT_FAMILY for run in runs
                )
            source_kind = "pptx"
        else:
            from docx import Document

            document = Document(str(source_path))
            runs = [
                run
                for paragraph in document.paragraphs
                for run in paragraph.runs
            ]
            assert runs and all(
                run.font.name == preflight.FONT_FAMILY for run in runs
            )
            source_kind = "docx"
        if mutate_source:
            source_path.write_bytes(source_path.read_bytes() + b"mutated")
        return _successful_render(source_kind, scope)

    monkeypatch.setattr(preflight.subprocess, "run", fake_run)
    monkeypatch.setattr(preflight, "read_deliverable", fake_read)
    return observed


def test_run_preflight_mocks_version_conversion_and_fc_match_boundaries(
    monkeypatch,
):
    observed = _mock_boundaries(monkeypatch)

    result = preflight.run_preflight()

    assert result["ok"] is True
    assert result["renderer_fingerprint"] == FINGERPRINT
    assert result["font_family"] == "Liberation Sans"
    assert [render["source_kind"] for render in result["renders"]] == [
        "xlsx", "pptx", "docx"
    ]
    assert [call[3] for call in observed["renders"]] == [
        {"workbook_page": 1}, {"slide": 1}, {"page": 1}
    ]
    assert len(observed["fc_match"]) == 1
    command, kwargs = observed["fc_match"][0]
    assert command == [
        "/usr/bin/fc-match", "--format", "%{family}\n", "Liberation Sans"
    ]
    assert kwargs["shell"] is False
    assert kwargs["timeout"] <= 10
    assert "AZURE_OPENAI_API_KEY" not in kwargs["env"]
    assert "SECRET_TOKEN" not in kwargs["env"]


def test_run_preflight_rejects_source_mutation(monkeypatch):
    _mock_boundaries(monkeypatch, mutate_source=True)

    with pytest.raises(preflight.RendererPreflightError, match="mutated"):
        preflight.run_preflight()


def test_pin_matches_every_published_grade_file():
    """The pin is only meaningful if it is the corpus's actual renderer.

    Read from the published artifacts rather than restated, because a pin
    typed by hand can drift from the thing it claims to pin and nothing
    would say so.
    """
    grades_dir = Path(__file__).resolve().parents[2] / "data" / "grades"
    observed = set()
    for path in sorted(grades_dir.glob("*.json")):
        payload = json.loads(path.read_text(encoding="utf-8"))
        fingerprint = payload.get("renderer_fingerprint")
        if isinstance(fingerprint, dict) and fingerprint.get(
            "libreoffice_version"
        ):
            observed.add(fingerprint["libreoffice_version"])

    assert observed, "no published grade file records a renderer fingerprint"
    assert observed == {preflight.EXPECTED_LIBREOFFICE_VERSION}


def test_version_drift_fails_closed_before_any_paid_work(monkeypatch):
    monkeypatch.delenv(preflight.ALLOW_VERSION_DRIFT_ENV, raising=False)
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)

    with pytest.raises(
        preflight.RendererPreflightError, match="not be comparable"
    ):
        preflight._validate_pinned_version(
            {**FINGERPRINT, "libreoffice_version": "LibreOffice 7.4.7.2 40"}
        )


def test_version_drift_is_not_satisfied_by_a_prefix(monkeypatch):
    """``24.2.7.2`` alone is a different build, not a looser spelling."""
    monkeypatch.delenv(preflight.ALLOW_VERSION_DRIFT_ENV, raising=False)
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)

    with pytest.raises(preflight.RendererPreflightError):
        preflight._validate_pinned_version(
            {**FINGERPRINT, "libreoffice_version": "LibreOffice 24.2.7.2"}
        )


def test_expected_version_env_repoints_the_pin(monkeypatch):
    monkeypatch.delenv(preflight.ALLOW_VERSION_DRIFT_ENV, raising=False)
    monkeypatch.setenv(preflight.EXPECTED_VERSION_ENV, "LibreOffice 25.8.1.2")

    assert (
        preflight._validate_pinned_version(
            {**FINGERPRINT, "libreoffice_version": "LibreOffice 25.8.1.2"}
        )
        is False
    )


def test_accepted_drift_warns_and_is_recorded_rather_than_hidden(
    monkeypatch, capsys
):
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)
    monkeypatch.setenv(preflight.ALLOW_VERSION_DRIFT_ENV, "1")

    accepted = preflight._validate_pinned_version(
        {**FINGERPRINT, "libreoffice_version": "LibreOffice 7.4.7.2 40"}
    )

    assert accepted is True
    assert "::warning::" in capsys.readouterr().err


@pytest.mark.parametrize("value", ["", "0", "true", "yes", " 1 x"])
def test_only_exactly_1_accepts_drift(monkeypatch, value):
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)
    monkeypatch.setenv(preflight.ALLOW_VERSION_DRIFT_ENV, value)

    with pytest.raises(preflight.RendererPreflightError):
        preflight._validate_pinned_version(
            {**FINGERPRINT, "libreoffice_version": "LibreOffice 7.4.7.2 40"}
        )


def test_matching_version_reports_no_drift(monkeypatch):
    monkeypatch.delenv(preflight.ALLOW_VERSION_DRIFT_ENV, raising=False)
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)

    assert preflight._validate_pinned_version(FINGERPRINT) is False


def test_successful_preflight_publishes_the_pin_it_enforced(monkeypatch):
    monkeypatch.delenv(preflight.ALLOW_VERSION_DRIFT_ENV, raising=False)
    monkeypatch.delenv(preflight.EXPECTED_VERSION_ENV, raising=False)
    _mock_boundaries(monkeypatch)

    result = preflight.run_preflight()

    assert (
        result["expected_libreoffice_version"]
        == preflight.EXPECTED_LIBREOFFICE_VERSION
    )
    assert result["renderer_version_drift_accepted"] is False


@pytest.mark.parametrize("family", ["DejaVu Sans", "Liberation Sans Narrow"])
def test_fc_match_must_return_exact_liberation_sans(monkeypatch, tmp_path, family):
    monkeypatch.setattr(preflight.shutil, "which", lambda _: "/usr/bin/fc-match")
    monkeypatch.setattr(
        preflight.subprocess,
        "run",
        lambda *args, **kwargs: type(
            "Completed",
            (),
            {"returncode": 0, "stdout": f"{family}\n", "stderr": ""},
        )(),
    )

    with pytest.raises(preflight.RendererPreflightError, match="did not resolve"):
        preflight._validate_liberation_sans(tmp_path)


def test_render_validation_rejects_wrong_kind():
    result = _successful_render("xlsx", {"workbook_page": 1})
    result["data"]["kind"] = "image_jpeg_base64"
    result["data"]["base64"] = "not-valid-base64"

    with pytest.raises(preflight.RendererPreflightError, match="render kind"):
        preflight._validate_render_result(
            result,
            source_kind="xlsx",
            scope={"workbook_page": 1},
            fingerprint=FINGERPRINT,
        )


@pytest.mark.parametrize(
    "encoded",
    ["%%%not-base64%%%", base64.b64encode(b"not-a-png").decode("ascii")],
)
def test_render_validation_rejects_invalid_png(encoded):
    result = _successful_render("xlsx", {"workbook_page": 1})
    result["data"]["base64"] = encoded
    result["data"]["byte_size"] = len(b"not-a-png")

    with pytest.raises(preflight.RendererPreflightError):
        preflight._validate_render_result(
            result,
            source_kind="xlsx",
            scope={"workbook_page": 1},
            fingerprint=FINGERPRINT,
        )


def test_main_prints_exactly_one_compact_json_on_success(monkeypatch, capsys):
    def noisy_success():
        print("captured library output")
        return {"ok": True, "renderer_fingerprint": FINGERPRINT}

    monkeypatch.setattr(preflight, "run_preflight", noisy_success)

    assert preflight.main() == 0
    captured = capsys.readouterr()
    assert captured.err == ""
    assert len(captured.out.splitlines()) == 1
    assert json.loads(captured.out)["ok"] is True
    assert '": ' not in captured.out
    assert '", ' not in captured.out


def test_main_prints_exactly_one_compact_json_on_failure(monkeypatch, capsys):
    def noisy_failure():
        print("captured library output")
        raise preflight.RendererPreflightError("renderer unavailable")

    monkeypatch.setattr(preflight, "run_preflight", noisy_failure)

    assert preflight.main() != 0
    captured = capsys.readouterr()
    assert captured.err == ""
    assert len(captured.out.splitlines()) == 1
    payload = json.loads(captured.out)
    assert payload == {
        "error": "renderer unavailable",
        "error_type": "RendererPreflightError",
        "ok": False,
    }
    assert '": ' not in captured.out
    assert '", ' not in captured.out