"""Tests for Reference File Reader

Usage:
    pytest tests/test_file_reader.py -v
"""

import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock

from core.file_reader import (
    read_reference_file,
    read_all_references,
    get_supported_extensions,
    _read_xlsx,
    _read_pdf,
    _read_docx,
    _read_pptx,
    _read_text,
    _describe_image,
    _describe_media,
)


class TestGetSupportedExtensions:
    """Test get_supported_extensions()"""

    def test_returns_all_extensions(self):
        exts = get_supported_extensions()
        assert ".xlsx" in exts
        assert ".pdf" in exts
        assert ".docx" in exts
        assert ".pptx" in exts
        assert ".txt" in exts
        assert ".png" in exts
        assert ".jpg" in exts
        assert ".wav" in exts
        assert ".mp3" in exts
        assert ".mp4" in exts
        assert ".zip" in exts
        assert ".step" in exts
        assert ".psd" in exts
        assert ".pages" in exts

    def test_text_extractable_marked_as_text(self):
        exts = get_supported_extensions()
        assert exts[".xlsx"] == "text"
        assert exts[".pdf"] == "text"
        assert exts[".docx"] == "text"
        assert exts[".pptx"] == "text"
        assert exts[".txt"] == "text"

    def test_binary_marked_as_description(self):
        exts = get_supported_extensions()
        assert exts[".png"] == "description"
        assert exts[".mp4"] == "description"
        assert exts[".zip"] == "description"


class TestReadText:
    """Test plain text reading"""

    def test_read_text_file(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("Hello World\nLine 2", encoding="utf-8")
        result = _read_text(str(f))
        assert "Hello World" in result
        assert "Line 2" in result


class TestReadXlsx:
    """Test Excel reading"""

    def test_read_xlsx_basic(self, tmp_path):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Value"])
        ws.append(["Alice", 100])
        ws.append(["Bob", 200])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))

        result = _read_xlsx(str(path))
        assert "[Sheet: Sheet1]" in result
        assert "Alice" in result
        assert "100" in result

    def test_read_xlsx_multiple_sheets(self, tmp_path):
        import openpyxl
        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Data"
        ws1.append(["col1", "col2"])
        ws2 = wb.create_sheet("Summary")
        ws2.append(["total", "300"])
        path = tmp_path / "multi.xlsx"
        wb.save(str(path))

        result = _read_xlsx(str(path))
        assert "[Sheet: Data]" in result
        assert "[Sheet: Summary]" in result


class TestReadDocx:
    """Test Word document reading"""

    def test_read_docx_basic(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")
        path = tmp_path / "test.docx"
        doc.save(str(path))

        result = _read_docx(str(path))
        assert "First paragraph" in result
        assert "Second paragraph" in result

    def test_read_docx_with_table(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Header")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        path = tmp_path / "table.docx"
        doc.save(str(path))

        result = _read_docx(str(path))
        assert "A1" in result
        assert "B2" in result


class TestReadPptx:
    """Test PowerPoint reading"""

    def test_read_pptx_basic(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide content"
        path = tmp_path / "test.pptx"
        prs.save(str(path))

        result = _read_pptx(str(path))
        assert "[Slide 1]" in result
        assert "Slide Title" in result
        assert "Slide content" in result


class TestDescribeImage:
    """Test image description"""

    def test_describe_png(self, tmp_path):
        from PIL import Image
        img = Image.new("RGB", (100, 200), color="red")
        path = tmp_path / "test.png"
        img.save(str(path))

        result = _describe_image(str(path))
        assert "test.png" in result
        assert "100x200" in result
        assert "RGB" in result

    def test_describe_image_missing_pillow(self, tmp_path):
        """Still returns description even if PIL fails"""
        path = tmp_path / "fake.png"
        path.write_bytes(b"not an image")

        result = _describe_image(str(path))
        assert "fake.png" in result


class TestDescribeMedia:
    """Test media file description"""

    def test_describe_media(self, tmp_path):
        path = tmp_path / "audio.wav"
        path.write_bytes(b"\x00" * 1024)

        result = _describe_media(str(path))
        assert "audio.wav" in result
        assert "KB" in result


class TestReadReferenceFile:
    """Test the unified read_reference_file()"""

    def test_file_not_found(self):
        result = read_reference_file("/nonexistent/file.xlsx")
        assert "File not found" in result

    def test_unsupported_extension(self, tmp_path):
        f = tmp_path / "data.xyz"
        f.write_bytes(b"data")
        result = read_reference_file(str(f))
        assert "Unsupported file" in result
        assert ".xyz" in result

    def test_text_file(self, tmp_path):
        f = tmp_path / "notes.txt"
        f.write_text("some notes")
        result = read_reference_file(str(f))
        assert result == "some notes"

    def test_error_handling(self, tmp_path):
        """Corrupted file returns error description instead of crashing"""
        f = tmp_path / "bad.xlsx"
        f.write_bytes(b"this is not excel")
        result = read_reference_file(str(f))
        assert "Error reading" in result


class TestReadAllReferences:
    """Test read_all_references()"""

    def test_no_reference_files(self):
        task = MagicMock()
        task.reference_files = []
        result = read_all_references(task)
        assert result == ""

    def test_with_text_file(self, tmp_path):
        # Create reference file
        ref_dir = tmp_path / "reference_files" / "abc123"
        ref_dir.mkdir(parents=True)
        (ref_dir / "notes.txt").write_text("Important note")

        task = MagicMock()
        task.reference_files = ["reference_files/abc123/notes.txt"]

        result = read_all_references(task, base_dir=str(tmp_path))
        assert "--- File: notes.txt ---" in result
        assert "Important note" in result

    def test_truncation(self, tmp_path):
        ref_dir = tmp_path / "reference_files" / "abc123"
        ref_dir.mkdir(parents=True)
        (ref_dir / "big.txt").write_text("x" * 100)

        task = MagicMock()
        task.reference_files = ["reference_files/abc123/big.txt"]

        result = read_all_references(task, base_dir=str(tmp_path), max_chars_per_file=50)
        assert "Truncated at 50 chars" in result

    def test_multiple_files(self, tmp_path):
        ref_dir = tmp_path / "reference_files" / "abc123"
        ref_dir.mkdir(parents=True)
        (ref_dir / "file1.txt").write_text("Content 1")
        (ref_dir / "file2.txt").write_text("Content 2")

        task = MagicMock()
        task.reference_files = [
            "reference_files/abc123/file1.txt",
            "reference_files/abc123/file2.txt",
        ]

        result = read_all_references(task, base_dir=str(tmp_path))
        assert "--- File: file1.txt ---" in result
        assert "--- File: file2.txt ---" in result
        assert "Content 1" in result
        assert "Content 2" in result


class TestRealDataset:
    """Integration test against real dataset (skipped if not available)"""

    @pytest.fixture
    def gdpval_path(self):
        path = Path(__file__).parent.parent.parent / "data" / "gdpval-local"
        if not path.exists():
            pytest.skip("GDPVal dataset not available locally")
        return path

    def test_read_real_xlsx(self, gdpval_path):
        """Find and read a real xlsx file from the dataset"""
        ref_dir = gdpval_path / "reference_files"
        for xlsx in ref_dir.rglob("*.xlsx"):
            result = read_reference_file(str(xlsx))
            assert "[Sheet:" in result
            assert len(result) > 0
            break  # Only test one file

    def test_read_real_pdf(self, gdpval_path):
        """Find and read a real pdf file from the dataset"""
        ref_dir = gdpval_path / "reference_files"
        for pdf in ref_dir.rglob("*.pdf"):
            result = read_reference_file(str(pdf))
            assert len(result) > 0
            break

    def test_read_real_docx(self, gdpval_path):
        """Find and read a real docx file from the dataset"""
        ref_dir = gdpval_path / "reference_files"
        for docx in ref_dir.rglob("*.docx"):
            result = read_reference_file(str(docx))
            assert len(result) > 0
            break
