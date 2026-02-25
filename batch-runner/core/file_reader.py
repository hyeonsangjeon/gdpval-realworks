"""Reference File Reader for GDPVal Tasks

Reads reference files (xlsx, pdf, docx, pptx, txt, csv, png, jpg, etc.)
and converts their content to text for inclusion in LLM prompts.

Supported formats:
    Text-extractable:
        .xlsx  → sheet-by-sheet CSV representation
        .pdf   → page-by-page text extraction
        .docx  → paragraph text extraction
        .pptx  → slide-by-slide text extraction
        .txt   → raw text
        .csv   → raw text (treated as .txt)

    Described only (binary/media):
        .png, .jpg, .jpeg  → "[Image file: filename.png]"
        .wav, .mp3, .mp4   → "[Media file: filename.mp4]"
        .zip               → "[Archive: filename.zip]"
        .step              → "[CAD file: filename.step]"
        .psd               → "[Photoshop file: filename.psd]"
        .pages             → "[Apple Pages file: filename.pages]"

Usage:
    from core.file_reader import read_reference_file, read_all_references

    # Single file
    text = read_reference_file("data/gdpval-local/reference_files/abc123/file.xlsx")

    # All files for a task
    combined = read_all_references(task, base_dir="data/gdpval-local")
"""

import os
import csv
import io
from pathlib import Path
from typing import Optional, List

from .config import DEFAULT_LOCAL_PATH


# ─── Text-extractable formats ─────────────────────────────────────────────

def _read_xlsx(filepath: str) -> str:
    """Read Excel file, returning sheet-by-sheet CSV-like representation"""
    import openpyxl

    wb = openpyxl.load_workbook(filepath, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(",".join(cells))
        parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _read_pdf(filepath: str) -> str:
    """Read PDF file, returning page-by-page text"""
    import pdfplumber

    parts = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            # Also try extracting tables
            tables = page.extract_tables()
            table_text = ""
            if tables:
                for t_idx, table in enumerate(tables):
                    table_rows = []
                    for row in table:
                        cells = [str(c) if c is not None else "" for c in row]
                        table_rows.append(" | ".join(cells))
                    table_text += f"\n[Table {t_idx + 1}]\n" + "\n".join(table_rows)

            page_content = text
            if table_text:
                page_content += "\n" + table_text

            if page_content.strip():
                parts.append(f"[Page {i}]\n{page_content.strip()}")
    return "\n\n".join(parts)


def _read_docx(filepath: str) -> str:
    """Read Word document, returning paragraph text"""
    from docx import Document

    doc = Document(filepath)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Also extract tables
    for t_idx, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\n[Table {t_idx + 1}]\n" + "\n".join(rows))

    return "\n".join(parts)


def _read_pptx(filepath: str) -> str:
    """Read PowerPoint file, returning slide-by-slide text"""
    from pptx import Presentation

    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _read_text(filepath: str) -> str:
    """Read plain text file"""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


# ─── Binary/media formats (description only) ──────────────────────────────

def _describe_image(filepath: str) -> str:
    """Describe an image file with basic metadata"""
    filename = os.path.basename(filepath)
    try:
        from PIL import Image
        img = Image.open(filepath)
        w, h = img.size
        mode = img.mode
        return f"[Image file: {filename} ({w}x{h}, {mode})]"
    except Exception:
        return f"[Image file: {filename}]"


def _describe_media(filepath: str) -> str:
    """Describe a media file with ffmpeg metadata extraction"""
    import subprocess
    import json

    filename = os.path.basename(filepath)
    size_kb = os.path.getsize(filepath) / 1024

    # Try to extract metadata using ffprobe
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "error", "-print_format", "json",
             "-show_format", "-show_streams", filepath],
            capture_output=True,
            text=True,
            timeout=10
        )

        if result.returncode == 0:
            data = json.loads(result.stdout)
            streams = data.get("streams", [])
            fmt = data.get("format", {})

            if streams:
                stream = streams[0]
                duration = float(fmt.get("duration", 0))
                sample_rate = stream.get("sample_rate")
                channels = stream.get("channels")
                codec_name = stream.get("codec_name", "unknown")
                bitrate = fmt.get("bit_rate")

                # Format bitrate
                if bitrate:
                    bitrate_kbps = int(bitrate) // 1000
                    bitrate_str = f", Bitrate: {bitrate_kbps}kbps"
                else:
                    bitrate_str = ""

                specs = f"Duration: {duration:.1f}s, Sample Rate: {sample_rate}Hz, Channels: {channels}, Codec: {codec_name}{bitrate_str}"
                return f"[Audio file: {filename} ({size_kb:.0f}KB) - {specs}]"
    except (subprocess.TimeoutExpired, subprocess.CalledProcessError, json.JSONDecodeError, FileNotFoundError):
        # ffprobe not available or error — fall back to basic info
        pass

    # Fallback: basic file info only
    return f"[Media file: {filename} ({size_kb:.0f}KB)]"


def _describe_generic(filepath: str, file_type: str) -> str:
    """Describe a generic binary file"""
    filename = os.path.basename(filepath)
    size_kb = os.path.getsize(filepath) / 1024
    return f"[{file_type}: {filename} ({size_kb:.0f} KB)]"


# ─── Extension → Reader mapping ───────────────────────────────────────────

_READERS = {
    # Text-extractable
    ".xlsx": _read_xlsx,
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".txt": _read_text,
    ".csv": _read_text,
    # Images
    ".png": _describe_image,
    ".jpg": _describe_image,
    ".jpeg": _describe_image,
    # Media
    ".wav": _describe_media,
    ".mp3": _describe_media,
    ".mp4": _describe_media,
    # Other binary
    ".zip": lambda f: _describe_generic(f, "Archive"),
    ".step": lambda f: _describe_generic(f, "CAD file"),
    ".psd": lambda f: _describe_generic(f, "Photoshop file"),
    ".pages": lambda f: _describe_generic(f, "Apple Pages file"),
}


# ─── Public API ────────────────────────────────────────────────────────────

def read_reference_file(filepath: str) -> str:
    """Read a single reference file and return its text content.

    Args:
        filepath: Absolute or relative path to the file

    Returns:
        Extracted text content, or a description placeholder for binary files

    Raises:
        FileNotFoundError: If the file doesn't exist
    """
    path = Path(filepath)
    if not path.exists():
        return f"[File not found: {path.name}]"

    ext = path.suffix.lower()
    reader = _READERS.get(ext)

    if reader is None:
        return _describe_generic(str(path), f"Unsupported file ({ext})")

    try:
        return reader(str(path))
    except Exception as e:
        return f"[Error reading {path.name}: {e}]"


def read_all_references(
    task,
    base_dir: Optional[str] = None,
    max_chars_per_file: int = 50_000,
) -> str:
    """Read all reference files for a GDPValTask and combine into one text block.

    Args:
        task: GDPValTask object (must have reference_files attribute)
        base_dir: Base directory where reference_files are located
                  (default: data/gdpval-local)
        max_chars_per_file: Maximum characters per file to prevent token explosion
                           (default: 50,000 chars ≈ ~12,500 tokens)

    Returns:
        Combined text block with all reference file contents,
        formatted with headers per file. Empty string if no files.
    """
    if not task.reference_files:
        return ""

    base = Path(base_dir) if base_dir else DEFAULT_LOCAL_PATH
    parts = []

    for ref_path in task.reference_files:
        full_path = base / ref_path
        filename = os.path.basename(ref_path)

        content = read_reference_file(str(full_path))

        # Truncate if too long
        if len(content) > max_chars_per_file:
            content = content[:max_chars_per_file] + f"\n... [Truncated at {max_chars_per_file} chars]"

        parts.append(f"--- File: {filename} ---\n{content}")

    return "\n\n".join(parts)


def get_supported_extensions() -> dict:
    """Get mapping of supported extensions and their read mode.

    Returns:
        Dict mapping extension → "text" or "description"
    """
    text_exts = {".xlsx", ".pdf", ".docx", ".pptx", ".txt", ".csv"}
    return {
        ext: "text" if ext in text_exts else "description"
        for ext in _READERS
    }
