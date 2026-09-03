"""``read_deliverable`` — read-only file inspection tool for the v2 grader.

Implements the 6-op surface from SPEC §4.2:

    inspect_structure / read_content / inspect_formatting
    render_to_image   / probe_audio   / probe_video

Design rules (per task spec 201):

1. All ops are read-only. They MUST NOT mutate files.
2. Paths are normalized against a trusted ``base_dir`` to prevent the
   judge from escaping out of the deliverables tree (no ``../`` traversal,
   no absolute paths that leave the base, no symlinks pointing outside).
3. Every op returns a uniform envelope:

       {"ok": True,  "data": <op-specific dict>}
       {"ok": False, "error": "<short reason>", "error_type": "<enum>"}

   Tool-calling judges serialize this directly back to JSON.
4. Large outputs (image bytes, raw content) are size-capped so the judge
   cannot accidentally inflate context by asking for a huge file.
5. PDF rendering uses ``PyMuPDF`` (fitz), not poppler. XLSX/PPTX inputs
    are converted read-only in an isolated temp directory with LibreOffice,
    then rasterized with PyMuPDF. Audio/video probing uses ``PyAV`` (``av``),
    not the ``ffmpeg`` binary. See
   ``tasks/rebuilding_grading_task/PR2_ENV_AUDIT.md`` for the rationale.
"""

from __future__ import annotations

import base64
import io
import json
import math
import os
import shutil
import subprocess
import tempfile
from contextlib import contextmanager
from functools import lru_cache
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple, Union

from core.media_types import (
    GRADER_AUDIO_EXTENSIONS,
    GRADER_SOURCE_CODE_EXTENSIONS,
    GRADER_VISUAL_RENDER_EXTENSIONS,
)

# ── Constants ─────────────────────────────────────────────────────────

#: Full public Python API. The harness uses render_to_image internally;
#: MODEL_READ_DELIVERABLE_OPS below is the smaller model-callable surface.
READ_DELIVERABLE_OPS: Tuple[str, ...] = (
    "inspect_structure",
    "read_content",
    "inspect_formatting",
    "render_to_image",
    "probe_audio",
    "probe_video",
)

#: Output size caps. The tool refuses to return more than this in a
#: single call. These bound prompt growth, not file size on disk.
MAX_CONTENT_CHARS = 200_000   # ``read_content`` text payload cap
MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5MB before/after downsample
MAX_CELLS_FORMATTING = 5_000  # inspect_formatting iterates capped cells

#: Word revision limits. A copy-editing deliverable can carry thousands of
#: tracked edits; the counts stay exact but the sampled bodies do not, so one
#: heavily-marked document cannot flood a judge's context.
MAX_COMMENTS_REPORTED = 50
MAX_COMMENT_CHARS = 500

#: Archive limits. A deliverable submitted as a ``.zip`` is a container of
#: files the other ops already handle, so the archive is listed and single
#: members can be opened through ``scope={"member": ...}``. Both numbers bound
#: what one call can cost: entries bound the listing sent to a judge, bytes
#: bound what is written to temp before an op reads it.
MAX_ZIP_ENTRIES = 2_000
MAX_ZIP_MEMBER_BYTES = 2 * 1024 * 1024 * 1024  # 2 GiB extracted per member

#: Archive members macOS writes alongside the real files. These are resource
#: forks, not deliverables, and listing them buries the actual content -- the
#: one gold answer in the stage-1 corpus that ships as an archive has five
#: real stems and five of these. Hidden from the default listing, still
#: openable by exact name so nothing is unreachable.
_ZIP_NOISE = ("__MACOSX/",)
_ZIP_NOISE_BASENAMES = ("._",)

MAX_PAGES_DEFAULT = 200       # PDF page-iteration safety cap
MAX_SHEETS = 100              # workbook safety cap

#: How much of a file is examined before deciding it is text. The question is
#: whether the bytes are text, and 64 KiB answers it -- a binary format puts
#: its magic number and its first NUL byte far inside this. Bounding the read
#: keeps the decision cheap on a large file.
TEXT_SNIFF_BYTES = 64 * 1024

#: Page geometry. A rubric asks "is it landscape", "is it letter size", "does
#: it fit on one page" -- questions about the page, which a page count alone
#: cannot answer. Measured pages are reported as one row per distinct size, so
#: a uniform 200-page report costs one row and a document that turns landscape
#: halfway costs two. Points are the PDF unit; inches are carried alongside
#: because that is the unit rubrics are written in.
MAX_PDF_PAGE_SIZE_GROUPS = 20
PDF_POINTS_PER_INCH = 72.0
RENDERER_VERSION_TIMEOUT_SEC = 10
RENDERER_VERSION_MAX_CHARS = 200

#: The member scope, said once for every schema that describes ``scope``.
#:
#: A judge meets an archive in three places -- the schema it is handed, the
#: structure listing, and the text read. The last two say it at runtime through
#: ``_ZIP_MEMBER_HINT``; the schemas say it here. It is a constant rather than
#: a sentence typed into each one because this file describes ``scope`` twice,
#: and the two descriptions had already drifted: the model schema carried the
#: member contract and the exported six-op schema did not, while the CHANGELOG
#: offers that schema as drop-in ``tools=[...]``. A caller who took it got a
#: model with no stated way to open anything an archive holds -- which is the
#: 2-of-62 archive answer this text was written to prevent, reintroduced one
#: export over.
_SCOPE_MEMBER_CONTRACT = (
    "To reach a file inside a .zip deliverable, pass "
    "{\"member\": \"<exact name from the listing>\"} to any op -- "
    "probe_audio for a stem, read_content for a document. An archive is a "
    "container of files this tool already reads, not an unreadable binary."
)

#: Full six-op schema retained for public API compatibility and direct
#: harness tests. Do not send this schema to the main model.
READ_DELIVERABLE_TOOL_SCHEMA: Dict[str, Any] = {
    "type": "function",
    "name": "read_deliverable",
    "description": (
        "Read-only inspection of a deliverable file produced by the model "
        "under test. Use this to verify structure, content, formatting, "
        "or to render pages to images for visual judgment. NEVER fabricate "
        "file contents — always call this tool to ground evidence quotes."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "op": {
                "type": "string",
                "enum": list(READ_DELIVERABLE_OPS),
                "description": (
                    "Operation to perform. "
                    "inspect_structure: file type, sheets/slides, page count "
                    "and page size. "
                    "read_content: textual content (no truncation up to cap). "
                    "inspect_formatting: cell fills/fonts/borders/charts/styles, "
                    "plus page count and page size. "
                    "render_to_image: PNG (base64) of an allowed first/page surface. "
                    "probe_audio: sample-rate/channels/duration/peak/silence. "
                    "probe_video: codec/duration/resolution/fps."
                ),
            },
            "path": {
                "type": "string",
                "description": "Path relative to the deliverables base dir.",
            },
            "scope": {
                "type": ["object", "null"],
                "description": (
                    "Optional op-specific scope, e.g. "
                    "{'workbook_page': 1}, {'slide': 1}, {'page': 1}, or "
                    "{'page_start': 1, 'page_end': 3}. "
                    + _SCOPE_MEMBER_CONTRACT
                ),
                "additionalProperties": True,
            },
        },
        "required": ["op", "path"],
        "additionalProperties": False,
    },
}

# The public Python API retains render_to_image for the harness. The model
# only receives this reduced schema and cannot request image bytes directly.
MODEL_READ_DELIVERABLE_OPS: Tuple[str, ...] = tuple(
    op for op in READ_DELIVERABLE_OPS if op != "render_to_image"
)
MODEL_READ_DELIVERABLE_TOOL_SCHEMA: Dict[str, Any] = {
    "type": "function",
    "name": "read_deliverable",
    "description": (
        "Read-only inspection of an allowlisted candidate or reference file. "
        "Use this to verify structure, content, formatting, audio, or video "
        "metadata. Visual rendering is performed only by the grading harness."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "op": {
                "type": "string",
                "enum": list(MODEL_READ_DELIVERABLE_OPS),
                "description": (
                    "inspect_structure: kind, size, sheets/slides, and for a "
                    "document or PDF the page count "
                    "(page_count, or converted_page_count for a .docx, which "
                    "stores no pagination of its own) and the page size "
                    "(width_pt/height_pt, width_in/height_in, orientation). "
                    "How long a document is and which way its pages face are "
                    "answerable ONLY from these fields: paragraph_count and "
                    "char_count are lengths, not page counts, and must never "
                    "be used as one. "
                    "read_content: text. "
                    "inspect_formatting: styles, fills, fonts, borders, and "
                    "the same page count and page size. For a .docx it also "
                    "reports the editing marks: tracked_insertion_count, "
                    "tracked_deletion_count, and comments (author, text, and "
                    "the anchored_text each one is attached to). Whether a "
                    "document carries Track Changes or reviewer comments is "
                    "answerable ONLY from these fields -- read_content shows "
                    "an insertion as ordinary text and omits a deletion "
                    "altogether, so it can never witness either. "
                    "probe_audio / probe_video: media metadata."
                ),
            },
            "path": {
                "type": "string",
                "description": "Exact path from an allowlist in the prompt.",
            },
            "scope": {
                "type": ["object", "null"],
                "description": (
                    "Optional op-specific content/formatting scope. "
                    + _SCOPE_MEMBER_CONTRACT
                ),
                "additionalProperties": True,
            },
        },
        "required": ["op", "path"],
        "additionalProperties": False,
    },
}


# ── Errors / path safety ─────────────────────────────────────────────


class ReadDeliverableError(Exception):
    """Raised for non-recoverable misuse (programmer error)."""


class RendererDependencyError(ReadDeliverableError):
    """Raised when an external renderer required by an op is unavailable."""


class InvalidScope(ReadDeliverableError):
    """Raised when a render scope has invalid keys or values."""


class UnsupportedScope(ReadDeliverableError):
    """Raised when a validly-shaped render scope is not implemented."""


def _envelope_error(msg: str, kind: str = "error") -> Dict[str, Any]:
    return {"ok": False, "error": msg, "error_type": kind}


def _envelope_ok(data: Any) -> Dict[str, Any]:
    return {"ok": True, "data": data}


def _resolve_trusted_path(path: str, base_dir: str) -> Optional[Path]:
    """Return the resolved path iff it lives inside ``base_dir``.

    Rejects: absolute paths outside base, ``..`` traversal, symlinks
    pointing outside base, missing files.
    """
    if not path or not isinstance(path, str):
        return None
    base = Path(base_dir).resolve()
    candidate = Path(path)
    if candidate.is_absolute():
        resolved = candidate.resolve()
    else:
        resolved = (base / candidate).resolve()
    try:
        resolved.relative_to(base)
    except ValueError:
        return None
    if not resolved.exists() or not resolved.is_file():
        return None
    return resolved


# ── File type detection ──────────────────────────────────────────────


_EXT_KIND = {
    ".xlsx": "xlsx", ".xlsm": "xlsx",
    ".docx": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
    ".csv": "csv",
    ".txt": "txt", ".md": "txt",
    ".json": "txt",
    ".ipynb": "notebook",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".bmp": "image", ".webp": "image",
    **{extension: "audio" for extension in GRADER_AUDIO_EXTENSIONS},
    ".mp4": "video", ".mov": "video", ".webm": "video",
    ".mkv": "video", ".avi": "video",
    ".zip": "zip",
}


def _reads_as_text(path: Path) -> bool:
    """Does this file hold text? Answered by opening it, not by its name.

    Only ever asked about an extension ``_EXT_KIND`` does not map, so it
    cannot reinterpret a format the map already names.
    """
    try:
        with path.open("rb") as handle:
            head = handle.read(TEXT_SNIFF_BYTES)
    except OSError:
        # Unreadable is not the same as binary, but nothing downstream can
        # read it either, so the honest answer is the conservative one.
        return False
    if b"\x00" in head:
        return False
    if not head:
        # No bytes, so no byte that is not text. A zero-byte deliverable is a
        # real defect and reads better as an empty file the judge is told is
        # empty than as a file nothing here claims to understand.
        return True
    try:
        head.decode("utf-8")
    except UnicodeDecodeError as exc:
        # A multi-byte character can straddle the cap. That is a fact about
        # where the read stopped, not about the file -- but only when the read
        # actually stopped there, and only within one character of the end.
        if len(head) < TEXT_SNIFF_BYTES or exc.start < len(head) - 4:
            return False
    return True


def _kind_of(path: Path) -> str:
    """Map a file to the kind the ops dispatch on.

    The extension answers first and its answer is final: an ``.xlsx`` is a
    workbook whatever its bytes look like, and letting content override that
    would only invent ways to misread a format already named correctly.

    An unmapped extension is where the name has stopped saying anything, and
    the corpus ships several -- ``.py``, ``.yaml``, ``.overpassql`` are all
    gold deliverables, all plainly readable, and all previously ``unknown``,
    which routes zero rubric items. So the file is opened and looked at, on
    the same reasoning ``has_audio_content`` already uses for containers: the
    extension is a fact about naming, not about whether the bytes are text.
    """
    mapped = _EXT_KIND.get(path.suffix.lower())
    if mapped is not None:
        return mapped
    return "txt" if _reads_as_text(path) else "unknown"


# ── inspect_structure ────────────────────────────────────────────────


def _inspect_xlsx(p: Path) -> Dict[str, Any]:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
    sheets = []
    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        sheets.append({
            "name": name,
            "max_row": ws.max_row,
            "max_col": ws.max_column,
        })
    return {
        "kind": "xlsx",
        "sheet_count": len(wb.sheetnames),
        "sheets": sheets,
        "truncated": len(wb.sheetnames) > MAX_SHEETS,
    }


@lru_cache(maxsize=64)
def _cached_docx_page_count(identity: Tuple[str, int, int]) -> int:
    """Lay a .docx out once per ``(path, size, mtime_ns)``.

    Structure inspection and formatting inspection both want this number, a
    judge asks for each more than once per file, and every Writer conversion
    costs seconds. The key moves whenever the file does, so a cached count can
    never outlive the bytes it was measured from.
    """
    source = Path(identity[0])
    with tempfile.TemporaryDirectory(prefix="gdpval-grade-pagecount-") as temp:
        return _pdf_page_count(_convert_office_to_pdf(source, Path(temp)))


def _docx_converted_page_count(p: Path) -> Dict[str, Any]:
    """How many pages this .docx lays out to, or why that is not knowable.

    A .docx stores no pagination -- the count does not exist until a layout
    engine computes one -- so this converts with the same LibreOffice the
    render path uses and counts what came out. Reported under the render
    path's name for the same reason it is the same number.

    Six stage-1 rubric items asked "does it fit on one page" and were answered
    from ``paragraph_count`` or ``char_count``, neither of which is a page
    count. Five were marked down and the gold answers were the right length.

    Never raises. A missing or failing converter costs the page count and
    nothing else: the paragraph, table and style counts around it stay, and
    the judge is told the number is unknown rather than handed a proxy.
    """
    try:
        stat = p.stat()
        return {
            "converted_page_count": _cached_docx_page_count(
                (str(p), stat.st_size, stat.st_mtime_ns)
            )
        }
    except Exception as exc:  # noqa: BLE001
        return {
            "converted_page_count": None,
            "page_count_error": f"{type(exc).__name__}: {exc}",
            "page_count_note": (
                "page count unavailable for this document; paragraph_count "
                "and char_count are lengths, not page counts, and cannot "
                "stand in for one"
            ),
        }


def _inspect_docx(p: Path) -> Dict[str, Any]:
    from docx import Document  # type: ignore

    doc = Document(str(p))
    return {
        "kind": "docx",
        **_docx_converted_page_count(p),
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "section_count": len(doc.sections),
        "styles_sample": sorted({p.style.name for p in doc.paragraphs[:200]
                                 if p.style is not None})[:20],
    }


def _inspect_pptx(p: Path) -> Dict[str, Any]:
    from pptx import Presentation  # type: ignore

    pres = Presentation(str(p))
    slides = []
    for i, slide in enumerate(pres.slides, 1):
        slides.append({
            "index": i,
            "shape_count": len(slide.shapes),
            "layout_name": getattr(slide.slide_layout, "name", None),
        })
    return {"kind": "pptx", "slide_count": len(pres.slides), "slides": slides}


# ── Archives ─────────────────────────────────────────────────────────
#
# A ``.zip`` deliverable is not an unreadable binary. It is a container of
# files every other op here already handles, and refusing it refuses them all:
# one stage-1 gold answer is five WAV stems in an archive, and it scored 2 of
# 62 because thirty-four rubric items -- sample rate, bit depth, duration, key,
# tempo -- were answered "binary or unsupported for text read" about a file
# nothing ever opened. The single item it did pass was "exactly one top-level
# ZIP archive is submitted".


def _is_zip_noise(name: str) -> bool:
    """AppleDouble resource forks, which are not part of the deliverable."""
    return name.startswith(_ZIP_NOISE) or Path(name).name.startswith(
        _ZIP_NOISE_BASENAMES
    )


def _zip_entries(p: Path) -> Tuple[List[Dict[str, Any]], int, bool]:
    """Real members, count of hidden noise, and whether the list was cut."""
    import zipfile

    entries: List[Dict[str, Any]] = []
    hidden = 0
    truncated = False
    with zipfile.ZipFile(p) as archive:
        for info in archive.infolist():
            if info.is_dir():
                continue
            if _is_zip_noise(info.filename):
                hidden += 1
                continue
            if len(entries) >= MAX_ZIP_ENTRIES:
                truncated = True
                break
            entries.append({
                "name": info.filename,
                "size": info.file_size,
                "compressed_size": info.compress_size,
                "kind": _EXT_KIND.get(Path(info.filename).suffix.lower(), "unknown"),
            })
    return entries, hidden, truncated


#: The runtime half of ``_SCOPE_MEMBER_CONTRACT``: the same fact, said in the
#: two places a judge meets an archive without reading the schema again -- the
#: structure listing and the text read. Worded for a line inside a listing
#: rather than for a parameter description, but it has to name the same key and
#: the same ops, because a judge that does not know a member can be opened
#: reports the archive as unreadable.
_ZIP_MEMBER_HINT = (
    "open a member with scope={\"member\": \"<name>\"} on any op, "
    "e.g. probe_audio for a .wav or read_content for a .docx"
)


def _inspect_zip(p: Path) -> Dict[str, Any]:
    entries, hidden, truncated = _zip_entries(p)
    return {
        "kind": "zip",
        "entry_count": len(entries),
        "entries": entries,
        "hidden_resource_fork_count": hidden,
        "truncated": truncated,
        "note": _ZIP_MEMBER_HINT,
    }


def _read_zip_text(p: Path, _scope: Dict[str, Any]) -> str:
    """The manifest, as text.

    A judge that never learns about the member scope still has to be able to
    answer "does the archive contain a Bass stem in WAV format", so the listing
    itself is the content of an archive. The hint is repeated here rather than
    left to ``inspect_structure`` alone because a judge that came straight to
    ``read_content`` would otherwise see a list of names it has no stated way
    to open, and conclude the files behind them are unreadable.
    """
    entries, hidden, truncated = _zip_entries(p)
    lines = [f"[Archive: {len(entries)} file(s)]"]
    lines += [
        f"{entry['name']} | {entry['kind']} | {entry['size']} bytes"
        for entry in entries
    ]
    if truncated:
        lines.append(f"... listing truncated at {MAX_ZIP_ENTRIES} entries")
    if hidden:
        lines.append(f"({hidden} macOS resource-fork entries hidden)")
    lines.append(_ZIP_MEMBER_HINT)
    return "\n".join(lines)


@contextmanager
def open_archive_member(p: Path, member: str) -> Any:
    """One member on disk, under its own suffix, for the length of one op.

    The name has to already be in the archive, so there is no path a caller can
    name that is not a member -- traversal is refused by lookup rather than by
    sanitising. The extract is streamed to a temp file so a member is bounded
    by disk rather than held in memory, and the declared size is checked first
    so an archive cannot ask for more than the cap by lying cheaply.
    """
    import zipfile

    with zipfile.ZipFile(p) as archive:
        try:
            info = archive.getinfo(member)
        except KeyError:
            names = [
                entry["name"]
                for entry in _zip_entries(p)[0][:20]
            ]
            raise InvalidScope(
                f"no member {member!r} in archive; members: {names}"
            ) from None
        if info.is_dir():
            raise InvalidScope(f"member {member!r} is a directory, not a file")
        if info.file_size > MAX_ZIP_MEMBER_BYTES:
            raise ReadDeliverableError(
                f"member {member!r} is {info.file_size} bytes, over the "
                f"{MAX_ZIP_MEMBER_BYTES}-byte extraction cap"
            )

        temp_dir = tempfile.mkdtemp(prefix="gdpval-zip-")
        try:
            target = Path(temp_dir) / Path(member).name
            written = 0
            with archive.open(info) as source, open(target, "wb") as sink:
                while chunk := source.read(1024 * 1024):
                    written += len(chunk)
                    if written > MAX_ZIP_MEMBER_BYTES:
                        raise ReadDeliverableError(
                            f"member {member!r} exceeded the "
                            f"{MAX_ZIP_MEMBER_BYTES}-byte extraction cap while "
                            "being read; its declared size was understated"
                        )
                    sink.write(chunk)
            yield target
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)


def _orientation(width_pt: float, height_pt: float) -> str:
    if width_pt > height_pt:
        return "landscape"
    if height_pt > width_pt:
        return "portrait"
    return "square"


def _summarize_page_sizes(
    rects: List[Tuple[float, float]],
) -> Tuple[List[Dict[str, Any]], bool]:
    """Group measured page rectangles into one row per distinct size."""
    grouped: Dict[Tuple[float, float], Dict[str, Any]] = {}
    order: List[Tuple[float, float]] = []
    for page_number, (width, height) in enumerate(rects, 1):
        key = (round(float(width), 2), round(float(height), 2))
        entry = grouped.get(key)
        if entry is None:
            grouped[key] = {
                "width_pt": key[0],
                "height_pt": key[1],
                "width_in": round(key[0] / PDF_POINTS_PER_INCH, 2),
                "height_in": round(key[1] / PDF_POINTS_PER_INCH, 2),
                "orientation": _orientation(*key),
                "page_count": 1,
                "first_page": page_number,
            }
            order.append(key)
        else:
            entry["page_count"] += 1
    kept = order[:MAX_PDF_PAGE_SIZE_GROUPS]
    return [grouped[key] for key in kept], len(order) > len(kept)


def _pdf_geometry(
    rects: List[Tuple[float, float]], page_count: int
) -> Dict[str, Any]:
    """The page-geometry block both PDF ops report.

    One stage-1 gold answer lost an orientation item because the only geometry
    a judge could see was ``page_count: 1``. The page was 432x288 -- landscape,
    exactly as the rubric asked -- and the answer was marked wrong for it.
    """
    sizes, truncated = _summarize_page_sizes(rects)
    uniform = len(sizes) == 1 and not truncated
    geometry: Dict[str, Any] = {
        "pages_measured": len(rects),
        "page_sizes": sizes,
        "page_size_uniform": uniform,
    }
    if uniform:
        geometry["orientation"] = sizes[0]["orientation"]
    if truncated:
        geometry["page_sizes_truncated"] = True
    if len(rects) < page_count:
        # Only the measured prefix is described. Say so rather than let
        # "uniform" be read as a claim about pages nothing looked at.
        geometry["pages_measured_capped"] = True
    return geometry


def _fitz_page_rects(doc: Any) -> List[Tuple[float, float]]:
    rects: List[Tuple[float, float]] = []
    for index in range(min(doc.page_count, MAX_PAGES_DEFAULT)):
        rect = doc.load_page(index).rect
        rects.append((float(rect.width), float(rect.height)))
    return rects


def _pdf_page_count(p: Path) -> int:
    try:
        import fitz  # type: ignore
    except ImportError:
        import pdfplumber  # type: ignore
        with pdfplumber.open(p) as pdf:
            return len(pdf.pages)
    doc = fitz.open(str(p))
    try:
        return doc.page_count
    finally:
        doc.close()


def _inspect_pdf(p: Path) -> Dict[str, Any]:
    try:
        import fitz  # type: ignore
    except ImportError:
        import pdfplumber  # type: ignore
        with pdfplumber.open(p) as pdf:
            rects = [
                (float(page.width), float(page.height))
                for page in pdf.pages[:MAX_PAGES_DEFAULT]
            ]
            return {
                "kind": "pdf",
                "page_count": len(pdf.pages),
                **_pdf_geometry(rects, len(pdf.pages)),
            }
    doc = fitz.open(str(p))
    try:
        return {
            "kind": "pdf",
            "page_count": doc.page_count,
            **_pdf_geometry(_fitz_page_rects(doc), doc.page_count),
            "metadata": dict(doc.metadata or {}),
        }
    finally:
        doc.close()


def _inspect_audio(p: Path) -> Dict[str, Any]:
    return _probe_audio_impl(p, basic=True)


def _inspect_video(p: Path) -> Dict[str, Any]:
    return _probe_video_impl(p, basic=True)


def _op_inspect_structure(p: Path, _scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    size = p.stat().st_size
    base = {"kind": kind, "size_bytes": size, "filename": p.name}
    try:
        if kind == "xlsx":
            base.update(_inspect_xlsx(p))
        elif kind == "docx":
            base.update(_inspect_docx(p))
        elif kind == "pptx":
            base.update(_inspect_pptx(p))
        elif kind == "pdf":
            base.update(_inspect_pdf(p))
        elif kind == "audio":
            base["audio"] = _inspect_audio(p)
        elif kind == "video":
            base["video"] = _inspect_video(p)
        elif kind == "zip":
            base.update(_inspect_zip(p))
        elif kind == "notebook":
            base.update(_inspect_notebook(p))
        # txt/csv/image: size + kind only
    except Exception as exc:  # noqa: BLE001
        base["inspection_error"] = f"{type(exc).__name__}: {exc}"
    return base


# ── Notebooks ────────────────────────────────────────────────────────

#: An ``.ipynb`` is JSON, so the sniff above would call it text and a judge
#: would be handed the raw file. That reads badly in a specific way: the
#: corpus notebook is 2,285,938 characters of which 39,478 -- 1.7 per cent --
#: are source, the rest being base64 image payload. A raw read is cut at
#: MAX_CONTENT_CHARS, so 91 per cent of the file never arrives and most of
#: what does is base64. Flattening to source plus text output puts the whole
#: notebook, all 69 cells, inside the window at 48,933 characters.
#:
#: Parsed with ``json`` rather than ``nbformat``: the grading container is
#: pinned by digest, so this cannot add a package to it.

_NOTEBOOK_TEXT_MIMES = ("text/plain", "text/markdown", "application/json")


def _notebook_document(p: Path) -> Optional[Dict[str, Any]]:
    """The parsed notebook, or ``None`` if it is not one.

    A file that does not parse is not an error to report -- it is still text,
    and the caller falls back to reading it as text, which is what a broken
    notebook is.
    """
    try:
        document = json.loads(p.read_text(encoding="utf-8", errors="replace"))
    except ValueError:
        return None
    if not isinstance(document, dict) or not isinstance(document.get("cells"), list):
        return None
    return document


def _joined(value: Any) -> str:
    """Notebook string fields are either a string or a list of lines."""
    if isinstance(value, str):
        return value
    if isinstance(value, list):
        return "".join(part for part in value if isinstance(part, str))
    return ""


def _notebook_output_lines(output: Dict[str, Any]) -> List[str]:
    kind = output.get("output_type")
    if kind == "stream":
        return [f"    [{output.get('name', 'stdout')}] {_joined(output.get('text'))}"]
    if kind == "error":
        traceback = "\n".join(
            line for line in output.get("traceback", []) if isinstance(line, str)
        )
        return [f"    [error] {output.get('ename', '')}: "
                f"{output.get('evalue', '')}\n{traceback}"]
    if kind in ("execute_result", "display_data"):
        data = output.get("data")
        if not isinstance(data, dict):
            return []
        for mime in _NOTEBOOK_TEXT_MIMES:
            if mime in data:
                return [f"    [{mime}] {_joined(data[mime])}"]
        # An image output is named rather than dropped. A rubric asking
        # whether the notebook plots something is answered by the fact that a
        # figure exists here, and the bytes themselves would not help a text
        # read.
        return [f"    [{mime} output, not text]" for mime in sorted(data)]
    return []


def _read_notebook_text(p: Path, _scope: Dict[str, Any]) -> str:
    document = _notebook_document(p)
    if document is None:
        return p.read_text(encoding="utf-8", errors="replace")
    language = (
        (document.get("metadata") or {}).get("language_info", {}).get("name")
        or (document.get("metadata") or {}).get("kernelspec", {}).get("language")
        or "unknown"
    )
    cells = document["cells"]
    lines = [f"[Notebook: {len(cells)} cell(s), language {language}]"]
    length = len(lines[0])
    for index, cell in enumerate(cells, start=1):
        if not isinstance(cell, dict):
            continue
        block = [f"--- cell {index} ({cell.get('cell_type', 'unknown')}) ---",
                 _joined(cell.get("source"))]
        for output in cell.get("outputs") or []:
            if isinstance(output, dict):
                block.extend(_notebook_output_lines(output))
        lines.extend(block)
        length += sum(len(line) + 1 for line in block)
        if length > MAX_CONTENT_CHARS:
            # ``_op_read_content`` truncates and says so; stopping here only
            # avoids building text nobody will see.
            break
    return "\n".join(lines)


def _inspect_notebook(p: Path) -> Dict[str, Any]:
    document = _notebook_document(p)
    if document is None:
        return {"kind": "notebook", "note": "not valid notebook JSON; read_content "
                                            "returns the raw file as text"}
    cells = [cell for cell in document["cells"] if isinstance(cell, dict)]
    metadata = document.get("metadata") or {}
    counts: Dict[str, int] = {}
    outputs = 0
    for cell in cells:
        kind = str(cell.get("cell_type", "unknown"))
        counts[kind] = counts.get(kind, 0) + 1
        outputs += len(cell.get("outputs") or [])
    return {
        "kind": "notebook",
        "cell_count": len(cells),
        "cell_types": dict(sorted(counts.items())),
        "output_count": outputs,
        "language": (metadata.get("language_info", {}).get("name")
                     or metadata.get("kernelspec", {}).get("language")),
        "nbformat": document.get("nbformat"),
    }


# ── read_content ─────────────────────────────────────────────────────


def _read_xlsx_text(p: Path, scope: Dict[str, Any]) -> str:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(p, data_only=True, read_only=True)
    target = scope.get("sheet")
    parts: List[str] = []
    for name in wb.sheetnames[:MAX_SHEETS]:
        if target is not None and name != target:
            continue
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(",".join("" if c is None else str(c) for c in row))
        parts.append(f"[Sheet: {name}]\n" + "\n".join(rows))
        if sum(len(x) for x in parts) > MAX_CONTENT_CHARS:
            break
    return "\n\n".join(parts)


def _pdf_table_blocks(page: Any, page_number: int) -> List[str]:
    """Each ruled table on the page, one row per line, cells kept apart.

    ``extract_text`` reads a page the way a person's eye crosses it, which for
    a table means the header row arrives as one run of prose and the body rows
    as another. A criterion like "the report includes the grade for John Lewis
    Childs School" is then unanswerable from a document that answers it plainly:
    the school and its grade sit side by side on the page and land paragraphs
    apart in the extraction, so the judge reads real evidence and correctly
    concludes the grade is absent.

    Emitting the rows keeps a record together, in the same ``|``-separated
    shape ``_read_docx_text`` and ``_read_xlsx_text`` already use, so a table
    reads the same whoever authored it.

    The linear text stays as well. Table detection is a guess about ruling
    lines, and a wrong guess must not be able to delete content -- the
    structured view is added to the page, never substituted for it.
    """
    try:
        tables = page.extract_tables()
    except Exception:  # noqa: BLE001
        # A malformed table costs its own rows and nothing else: the page's
        # text is already collected and must survive a detector that trips.
        return []

    blocks: List[str] = []
    for table in tables or []:
        rows = [
            " | ".join(
                "" if cell is None else " ".join(str(cell).split())
                for cell in row
            )
            for row in table
            if row and any(cell not in (None, "") for cell in row)
        ]
        if rows:
            blocks.append(f"[Table on page {page_number}]\n" + "\n".join(rows))
    return blocks


def _read_pdf_text(p: Path, scope: Dict[str, Any]) -> str:
    import pdfplumber  # type: ignore

    page_start = int(scope.get("page_start", 1))
    page_end = int(scope.get("page_end", MAX_PAGES_DEFAULT))
    parts: List[str] = []
    with pdfplumber.open(p) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            if i < page_start or i > page_end:
                continue
            txt = page.extract_text() or ""
            if txt.strip():
                parts.append(f"[Page {i}]\n{txt.strip()}")
            parts.extend(_pdf_table_blocks(page, i))
            if sum(len(x) for x in parts) > MAX_CONTENT_CHARS:
                break
    return "\n\n".join(parts)


def _read_docx_text(p: Path, _scope: Dict[str, Any]) -> str:
    from docx import Document  # type: ignore

    doc = Document(str(p))
    parts = [para.text for para in doc.paragraphs if para.text.strip()]
    for tbl in doc.tables:
        rows = []
        for row in tbl.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        parts.append("[Table]\n" + "\n".join(rows))
    return "\n\n".join(parts)


#: Namespace of the DrawingML chart part. Used to read cached category and
#: value points straight out of the XML when python-pptx will not model the
#: plot type -- see ``_pptx_chart_xml_text``.
_CHART_NS = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}


def _pptx_chart_xml_text(chart: Any) -> List[str]:
    """Categories and values read from the chart part itself.

    python-pptx models a subset of plot types and raises ``ValueError:
    unsupported plot type`` on the rest -- ``pie3DChart`` among them, which is
    what four of the five charts in one stage-1 gold answer are. For those, the
    modelled accessors give nothing at all: not the values, not the categories,
    not even the type name.

    The cached points are in the XML regardless, because that is what lets a
    deck render without the workbook it was built from. Reading them there is
    uniform across every plot type, which is the point: the set of chart types
    a judge might meet is not the set python-pptx happens to model.
    """
    lines: List[str] = []
    plot_types = [
        etree_tag.split("}")[-1]
        for plot in chart.findall(".//c:plotArea/*", _CHART_NS)
        for etree_tag in [str(plot.tag)]
        if etree_tag.endswith("Chart")
    ]
    lines.append(f"[Chart: {', '.join(plot_types)}]" if plot_types else "[Chart]")

    for series in chart.findall(".//c:ser", _CHART_NS):
        name = next(
            (v.text for v in series.findall("./c:tx//c:v", _CHART_NS) if v.text),
            "series",
        )
        categories = [
            v.text or "" for v in series.findall("./c:cat//c:pt/c:v", _CHART_NS)
        ]
        values = [
            v.text or "" for v in series.findall("./c:val//c:pt/c:v", _CHART_NS)
        ]
        if categories:
            lines.append("categories: " + ", ".join(categories))
        if values:
            lines.append(f"{name}: " + ", ".join(values))
    return lines


def _pptx_chart_text(shape: Any) -> List[str]:
    """Category and series values behind a chart, when they can be read.

    Rubrics ask things like "no categories outside the specified 12 appear in
    the table or chart", which is unanswerable from a picture of a pie. The
    numbers are in the embedded workbook and in the chart part, and this tries
    the modelled accessors first -- they name the chart type and hand back
    typed values -- then falls back to the raw XML for the plot types
    python-pptx declines to model.

    Every step is guarded. This runs against arbitrary gold answers, so a chart
    that will not describe itself has to cost its own text and nothing else:
    one bad chart must not take the slides behind it down with it.
    """
    try:
        chart = shape.chart
    except Exception:
        return []

    lines: List[str] = []
    try:
        lines.append(f"[Chart: {chart.chart_type}]")
    except Exception:
        lines = []
    if lines:
        try:
            categories = [str(c) for c in chart.plots[0].categories]
            if categories:
                lines.append("categories: " + ", ".join(categories))
        except Exception:
            pass
        try:
            for series in chart.series:
                values = ", ".join("" if v is None else str(v) for v in series.values)
                lines.append(f"{series.name}: {values}")
        except Exception:
            pass

    if len(lines) < 2:
        try:
            lines = _pptx_chart_xml_text(chart._chartSpace)
        except Exception:
            return []
    return ["\n".join(lines)] if len(lines) > 1 else []


def _pptx_shape_text(shape: Any, depth: int = 0) -> List[str]:
    """Everything on one shape a judge should be able to read.

    ``shape.has_text_frame`` on its own misses most of what a real deck carries.
    A table is a graphic frame with no text frame, so every cell of it is
    invisible; a group is a single shape whose children are never visited. Both
    were silently dropped, which is how a five-slide deck read as 186 characters
    of slide titles while the table holding the twelve categories a rubric item
    asked about sat unread.
    """
    out: List[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            out.append(text)
    if getattr(shape, "has_table", False):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in shape.table.rows
        ]
        out.append("[Table]\n" + "\n".join(rows))
    if getattr(shape, "has_chart", False):
        out.extend(_pptx_chart_text(shape))
    # Groups nest, and a table one level down is as unreadable as one at the
    # top. The depth cap is only a guard against a malformed file describing a
    # cycle -- real decks do not nest anywhere near this deep.
    if depth < 8 and hasattr(shape, "shapes"):
        for child in shape.shapes:
            out.extend(_pptx_shape_text(child, depth + 1))
    return out


def _read_pptx_text(p: Path, _scope: Dict[str, Any]) -> str:
    from pptx import Presentation  # type: ignore

    pres = Presentation(str(p))
    parts: List[str] = []
    for i, slide in enumerate(pres.slides, 1):
        chunks = [f"[Slide {i}]"]
        for shape in slide.shapes:
            chunks.extend(_pptx_shape_text(shape))
        parts.append("\n".join(chunks))
        if sum(len(x) for x in parts) > MAX_CONTENT_CHARS:
            break
    return "\n\n".join(parts)


# ── an empty read is a fact about the file, not about its content ────
#
# One stage-1 gold answer is a two-page scan. Its text layer is empty, so
# ``read_content`` returned ``char_count: 0`` and ten rubric items were graded
# "that content is absent" -- about a document that contains all ten things,
# in ink, on pages the harness had already rendered twice. Nothing lied: the
# tool said the text was empty and the judge read empty text as an empty
# document. So the tool now says which of the two it means.


#: Attached to every empty read. The distinction it draws is the whole of
#: this section: "I could not read it here" and "it is not there" are
#: different findings, and only one of them is a reason to fail an item.
_EMPTY_READ_DISCLAIMER = (
    "an empty text read means this file carries no extractable text, NOT "
    "that the content is absent; grade absence only from a file that was "
    "actually read"
)

#: For the kinds that hold no text at all, the op that does read them. A
#: judge told only "unsupported" has been told the file is a dead end.
_NON_TEXT_READ_ROUTES: Dict[str, str] = {
    "audio": (
        "call probe_audio for sample rate, channels, duration, peak level "
        "and silence"
    ),
    "video": (
        "call probe_video for codec, duration, resolution and frame rate"
    ),
    "image": (
        "call inspect_structure for kind and size; an image is judged from "
        "the rendered visual evidence the harness supplies, not from text"
    ),
}


def _pdf_text_layer_facts(p: Path) -> Dict[str, Any]:
    """Pages and embedded images of a PDF that yielded no text.

    A scan and an empty file both read as zero characters. These two numbers
    tell them apart: two pages carrying two images is a scan whose content is
    in the images, and zero pages is a file with nothing in it.

    Never raises -- this runs on the failure path of a read that already
    returned nothing, and losing the read as well would leave the judge with
    strictly less than it has today.
    """
    try:
        import fitz  # type: ignore
    except ImportError:
        try:
            import pdfplumber  # type: ignore

            with pdfplumber.open(p) as pdf:
                return {"page_count": len(pdf.pages)}
        except Exception:  # noqa: BLE001
            return {}
    try:
        doc = fitz.open(str(p))
    except Exception:  # noqa: BLE001
        return {}
    try:
        images = 0
        for index in range(min(doc.page_count, MAX_PAGES_DEFAULT)):
            images += len(doc.load_page(index).get_images(full=True))
        return {"page_count": doc.page_count, "embedded_image_count": images}
    except Exception:  # noqa: BLE001
        return {}
    finally:
        doc.close()


def _empty_read_report(p: Path, kind: str) -> Tuple[Dict[str, Any], str]:
    """Structured facts and a note describing what an unreadable file is."""
    if kind != "pdf":
        return {"has_text_layer": False}, f"no extractable text -- {_EMPTY_READ_DISCLAIMER}"

    facts: Dict[str, Any] = {"has_text_layer": False}
    facts.update(_pdf_text_layer_facts(p))
    images = facts.get("embedded_image_count")
    pages = facts.get("page_count")
    if images:
        what = (
            f"this PDF has no text layer: its {pages} page(s) carry {images} "
            f"embedded image(s), so the content is in the page images -- it "
            f"is a scan or an exported graphic, not an empty document"
        )
    elif pages == 0:
        what = "this PDF has no pages"
    else:
        what = "this PDF has no text layer"
    return facts, f"{what} -- {_EMPTY_READ_DISCLAIMER}"


def _pdf_has_text(p: Path) -> Optional[bool]:
    """Whether any page of a PDF yields a character. ``None`` if unknowable.

    Returns on the first character found, so an ordinary document costs one
    page. A scan costs every page, which is the only way to be sure of a
    negative, and is cheap on pages whose content stream is one image draw.
    """
    try:
        import fitz  # type: ignore
    except ImportError:
        import pdfplumber  # type: ignore

        with pdfplumber.open(p) as pdf:
            for page in pdf.pages[:MAX_PAGES_DEFAULT]:
                if (page.extract_text() or "").strip():
                    return True
            return None if len(pdf.pages) > MAX_PAGES_DEFAULT else False

    doc = fitz.open(str(p))
    try:
        for index in range(min(doc.page_count, MAX_PAGES_DEFAULT)):
            if doc.load_page(index).get_text().strip():
                return True
        return None if doc.page_count > MAX_PAGES_DEFAULT else False
    finally:
        doc.close()


def has_extractable_text(path: Union[str, Path]) -> Optional[bool]:
    """Whether this file yields any text at all. ``None`` means unknowable.

    Routing asks this, not the judge. A file with no text layer answers no
    criterion from its text -- every question about its content is a question
    about its pages -- and the grader uses that to hand the item to the render
    and vision path instead of letting an empty read stand as an answer.

    ``False`` is a positive claim and is only returned when the whole file was
    examined. Anything unexamined, unsupported, or broken is ``None``: routing
    must not escalate on a guess, and the caller treats the two differently.
    """
    p = Path(path)
    if not p.is_file():
        return None
    kind = _kind_of(p)
    if kind == "image":
        # Not a failure to read. An image has no text layer by construction,
        # which is exactly the fact routing wants.
        return False
    try:
        if kind == "pdf":
            return _pdf_has_text(p)
        if kind == "docx":
            return bool(_read_docx_text(p, {}).strip())
        if kind == "pptx":
            return bool(_read_pptx_text(p, {}).strip())
        if kind == "xlsx":
            return bool(_read_xlsx_text(p, {}).strip())
        if kind in ("txt", "csv"):
            return bool(p.read_text(encoding="utf-8", errors="replace").strip())
        if kind == "notebook":
            return bool(_read_notebook_text(p, {}).strip())
    except Exception:  # noqa: BLE001
        return None
    # audio, video, zip, unknown: text is not the medium, and "this file has
    # no text" would be read as a finding about a file nothing can read here.
    return None


def has_audio_content(path: Union[str, Path]) -> Optional[bool]:
    """Whether there is anything here for the listening model to hear.

    Routing asks this, and the interesting answer is the archive. A folder of
    stems delivered as one ``.zip`` is an audio deliverable; the container
    extension is a fact about packaging, not about the medium. Reading it as
    a statement about the medium is how a stage-1 task made entirely of music
    -- tempo, key, vocals, mix -- was graded end to end without a single
    listening call.

    ``True`` for an audio file, an archive holding one, or a video whose
    container carries an audio track. ``False`` is a positive claim and means
    the file was examined and carries no audio. ``None`` is an admission --
    missing file, unreadable archive, absent decoder, or a member list that was
    cut short and may have stopped one entry before the audio -- so that
    routing never promotes on a guess.
    """
    p = Path(path)
    if not p.is_file():
        return None
    kind = _kind_of(p)
    if kind == "audio":
        return True
    if kind == "video":
        # The same mistake as the ``.zip`` above, one container along. A reel
        # is delivered as ``.mp4``; ``.mp4`` is disjoint from the audio
        # extensions; so "the sound effect is audible during the opening
        # shot" was demoted to TEXT and answered by a judge that had read the
        # file's metadata and could not hear it. It scored zero -- not
        # excluded, *failed* -- on evidence reading ``"audio_tracks": 1``.
        #
        # That count is the honest answer to this question and was already
        # being computed one call away.
        count = _video_audio_track_count(p)
        return None if count is None else count > 0
    if kind != "zip":
        return False
    try:
        entries, _hidden, truncated = _zip_entries(p)
    except Exception:  # noqa: BLE001
        return None
    if any(entry["kind"] == "audio" for entry in entries):
        return True
    return None if truncated else False


def has_only_source_code_content(path: Union[str, Path]) -> Optional[bool]:
    """Is this deliverable program text and nothing a renderer could draw?

    The sibling of ``has_audio_content``, asked for the opposite reason. That
    one exists so a criterion about sound is not answered by reading; this one
    exists so a criterion about *code* is not refused for want of a picture.

    ``True`` means every file here is source -- code, stylesheet or markup --
    and none of it can be turned into an image. That is a narrow claim and
    deliberately so. A React component archive is the case it was written for:
    its appearance is not a property of the submission at all, only of
    something that builds and runs it, so the JSX and the CSS are where the
    answer is actually written down. ``False`` says the file was examined and
    is not that -- it renders, or it holds a picture, a recording or data.
    ``None`` is an admission: missing file, unreadable archive, or a member
    list cut short that may have stopped one entry before the image.

    A ``.csv`` answers ``False``, which is the point. Data has a look that a
    reader would see on opening it, and "page layout is visually polished"
    against one is a question reading cannot honestly settle;
    ``test_explicit_visual_item_fails_closed_without_render_target`` pins that
    it keeps failing closed.

    Companion text -- the ``README.md`` and ``package.json`` that ship beside
    any real component -- does not spoil the claim, because nothing about
    those files is looked at either. What does spoil it is one member this
    module cannot name, and an unrecognised extension is refused for exactly
    that reason: it may be the screenshot.
    """
    p = Path(path)
    if not p.is_file():
        return None
    if p.suffix.lower() in GRADER_SOURCE_CODE_EXTENSIONS:
        return True
    if _kind_of(p) != "zip":
        return False
    try:
        entries, _hidden, truncated = _zip_entries(p)
    except Exception:  # noqa: BLE001
        return None

    def _is_source(entry: Dict[str, Any]) -> bool:
        return Path(entry["name"]).suffix.lower() in GRADER_SOURCE_CODE_EXTENSIONS

    def _has_nothing_to_look_at(entry: Dict[str, Any]) -> bool:
        # The renderable test is asked first and separately even though it is
        # redundant today, because it is the one that would stop being
        # redundant: put ``.html`` in the render set tomorrow and the source
        # set would otherwise wave through a file the prepass had just
        # learned to draw.
        if Path(entry["name"]).suffix.lower() in GRADER_VISUAL_RENDER_EXTENSIONS:
            return False
        # ``_zip_entries`` reads member kinds off the extension alone, so
        # "txt" here means ``.txt``, ``.md`` or ``.json`` and nothing that had
        # to be opened to be believed. Everything else -- an image, a
        # recording, a workbook, a ``.csv``, a nested archive, and every
        # extension this module cannot name -- fails.
        return _is_source(entry) or entry["kind"] == "txt"

    if not all(_has_nothing_to_look_at(entry) for entry in entries):
        return False
    if truncated:
        # Unlike the disqualifying member above, "nothing here can be seen" is
        # a claim about the whole archive, and a list cut at the cap cannot
        # support it: the entry that was never listed is precisely the one
        # that would have refuted it.
        return None
    # An archive of prose and configuration is not a source deliverable. It
    # would pass the test above on the strength of having no picture in it,
    # which is the ``.csv`` mistake moved inside a container.
    return any(_is_source(entry) for entry in entries)


def _op_read_content(p: Path, scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    if kind == "xlsx":
        text = _read_xlsx_text(p, scope)
    elif kind == "pdf":
        text = _read_pdf_text(p, scope)
    elif kind == "docx":
        text = _read_docx_text(p, scope)
    elif kind == "pptx":
        text = _read_pptx_text(p, scope)
    elif kind in ("txt", "csv"):
        text = p.read_text(encoding="utf-8", errors="replace")
    elif kind == "notebook":
        text = _read_notebook_text(p, scope)
    elif kind == "zip":
        text = _read_zip_text(p, scope)
    else:
        route = _NON_TEXT_READ_ROUTES.get(kind)
        note = "this file holds no text"
        if route:
            note = f"{note}; {route}"
        return {"kind": kind, "text": "", "char_count": 0, "truncated": False,
                "has_text_layer": False,
                "note": f"{note} -- {_EMPTY_READ_DISCLAIMER}"}
    truncated = len(text) > MAX_CONTENT_CHARS
    if truncated:
        text = text[:MAX_CONTENT_CHARS]
    result: Dict[str, Any] = {"kind": kind, "text": text,
                              "char_count": len(text), "truncated": truncated}
    notes: List[str] = []
    if kind == "docx":
        # The one required item stage 1 failed was failed here: 6,532
        # characters were read as "longer than one page". Length is not
        # layout, and the op that does lay the document out is named.
        notes.append(
            "char_count is a length, not a page count; call inspect_structure "
            "for converted_page_count"
        )
    if kind == "pdf" and text.strip():
        # A page can hold real text *and* carry its numbers as pasted
        # pictures, and extraction drops the pictures without saying so, so
        # the read comes back looking complete. The corpus's worst-scoring
        # gold deliverable is one such page: 572 characters of title, column
        # headers and row labels, with every grade, percentage and ratio
        # living in 30 embedded images. Reporting the count costs nothing and
        # turns a silent omission into the disclosed gap the judge's rule
        # about absence is already written to handle.
        images = _pdf_text_layer_facts(p).get("embedded_image_count")
        if images:
            notes.append(
                f"this PDF also carries {images} embedded image(s) whose "
                "content is not in the text above; a value the criterion asks "
                "for and you cannot find may be drawn inside one of them "
                "rather than absent from the document"
            )
    if not text.strip():
        facts, empty_note = _empty_read_report(p, kind)
        result.update(facts)
        notes.append(empty_note)
    if notes:
        result["note"] = " | ".join(notes)
    return result


# ── inspect_formatting ───────────────────────────────────────────────


def _safe_cell_color(color: Any) -> Optional[str]:
    """Return a stable token for an openpyxl color without reading inactive descriptors."""
    if color is None:
        return None

    color_type = getattr(color, "type", None)
    if color_type == "rgb":
        value = color.rgb
        if isinstance(value, str) and "Values must be of type" not in value:
            return value
    elif color_type == "theme":
        value = color.theme
        if isinstance(value, int) and not isinstance(value, bool):
            token = f"theme:{value}"
            tint = getattr(color, "tint", 0.0)
            if isinstance(tint, (int, float)) and not isinstance(tint, bool) and tint:
                token += f":tint:{tint:g}"
            return token
    elif color_type == "indexed":
        value = color.indexed
        if isinstance(value, int) and not isinstance(value, bool):
            return f"indexed:{value}"
    elif color_type == "auto" and color.auto is True:
        return "auto"
    return None


def _workbook_default_font_color(workbook: Any) -> Any:
    """Find the workbook's default font color across supported openpyxl layouts."""
    named_styles = getattr(workbook, "_named_styles", None)
    if named_styles is not None:
        try:
            normal_style = named_styles["Normal"]
        except (KeyError, TypeError):
            pass
        else:
            normal_font = getattr(normal_style, "font", None)
            if normal_font is not None:
                return getattr(normal_font, "color", None)

    fonts = getattr(workbook, "_fonts", None)
    if fonts:
        return getattr(fonts[0], "color", None)
    return None


def _nondefault_font_color(cell: Any, default_color: Any) -> Any:
    """Return a cell's explicit font color, omitting the workbook default."""
    font = getattr(cell, "font", None)
    color = getattr(font, "color", None)
    if default_color is not None:
        return color if color != default_color else None

    style = getattr(cell, "_style", None)
    font_id = getattr(style, "fontId", None)
    return color if font_id not in (None, 0) else None


def _format_xlsx(p: Path, scope: Dict[str, Any]) -> Dict[str, Any]:
    import openpyxl  # type: ignore

    wb = openpyxl.load_workbook(p, data_only=False)
    default_font_color = _workbook_default_font_color(wb)
    target = scope.get("sheet")
    sheets_meta = []
    for name in wb.sheetnames[:MAX_SHEETS]:
        if target is not None and name != target:
            continue
        ws = wb[name]
        merged = [str(r) for r in ws.merged_cells.ranges]
        col_widths = {k: (d.width if d.width is not None else None)
                      for k, d in ws.column_dimensions.items()}
        styled_cells: List[Dict[str, Any]] = []
        seen = 0
        for row in ws.iter_rows():
            for cell in row:
                if seen >= MAX_CELLS_FORMATTING:
                    break
                seen += 1
                if cell.value is None:
                    continue
                fill_color = _safe_cell_color(
                    cell.fill.fgColor if cell.fill else None
                )
                font_bold = bool(cell.font and cell.font.bold)
                font_color = _safe_cell_color(
                    _nondefault_font_color(cell, default_font_color)
                )
                has_border = bool(cell.border and (
                    (cell.border.top and cell.border.top.style) or
                    (cell.border.bottom and cell.border.bottom.style) or
                    (cell.border.left and cell.border.left.style) or
                    (cell.border.right and cell.border.right.style)
                ))
                if (fill_color and fill_color != "00000000") or font_bold \
                        or font_color or has_border:
                    styled_cells.append({
                        "ref": cell.coordinate,
                        "fill": fill_color,
                        "bold": font_bold,
                        "font_color": font_color,
                        "border": has_border,
                    })
            if seen >= MAX_CELLS_FORMATTING:
                break
        sheets_meta.append({
            "name": name,
            "merged_ranges": merged,
            "column_widths": col_widths,
            "has_charts": bool(getattr(ws, "_charts", [])),
            "styled_cells_sample": styled_cells[:50],
            "styled_cells_count": len(styled_cells),
            "cells_scanned": seen,
            "cells_scan_capped": seen >= MAX_CELLS_FORMATTING,
        })
    return {"kind": "xlsx", "sheets": sheets_meta}


#: WordprocessingML namespace. Tracked revisions and comments live in the
#: document XML and have no python-docx accessor, so they are read from there.
_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _docx_revisions_and_comments(p: Path) -> Dict[str, Any]:
    """Tracked insertions, tracked deletions, and reviewer comments.

    A copy-editing brief is graded on its marks: "substantive edits appear as
    Track Changes insertions", "the .docx includes reviewer comments anchored
    to specific text ranges". python-docx models none of this -- ``para.text``
    silently folds an insertion into the running text and drops a deletion
    entirely -- so a formatting report built from it can only describe a
    marked-up document as though it were clean. The criteria are then
    unanswerable, and an item nothing can see fails whatever the file contains.

    The marks are in ``word/document.xml`` (``w:ins``, ``w:del``) and
    ``word/comments.xml``, which is where this reads them.

    Counts are exact. Comment bodies and their anchors are sampled and
    truncated, because how many marks there are is the gradeable fact and a
    thousand of them must not cost a judge its context.

    A file that cannot be opened reports zeros rather than raising: this
    enriches a formatting report that is useful without it, so it must never be
    the reason the whole report fails.
    """
    blank: Dict[str, Any] = {
        "tracked_insertion_count": 0,
        "tracked_deletion_count": 0,
        "comment_count": 0,
        "comments": [],
    }
    import zipfile
    from xml.etree import ElementTree

    try:
        with zipfile.ZipFile(p) as archive:
            names = set(archive.namelist())
            if "word/document.xml" not in names:
                return blank
            body = ElementTree.fromstring(archive.read("word/document.xml"))
            comment_root = (
                ElementTree.fromstring(archive.read("word/comments.xml"))
                if "word/comments.xml" in names
                else None
            )
    except Exception:  # noqa: BLE001
        return blank

    def q(tag: str) -> str:
        return f"{{{_W_NS}}}{tag}"

    # Walk the body once in document order, carrying the set of comment ranges
    # that are currently open. Text seen while a range is open is that
    # comment's anchor -- the "specific text range" the rubric asks about.
    anchors: Dict[str, List[str]] = {}
    open_ids: set[str] = set()
    for element in body.iter():
        tag = str(element.tag).rsplit("}", 1)[-1]
        if tag == "commentRangeStart":
            open_ids.add(element.get(q("id"), ""))
        elif tag == "commentRangeEnd":
            open_ids.discard(element.get(q("id"), ""))
        elif tag in ("t", "delText") and open_ids and element.text:
            for comment_id in open_ids:
                anchors.setdefault(comment_id, []).append(element.text)

    comments: List[Dict[str, Any]] = []
    if comment_root is not None:
        for comment in comment_root.findall(f".//{q('comment')}"):
            comment_id = comment.get(q("id"), "")
            text = "".join(
                node.text or "" for node in comment.findall(f".//{q('t')}")
            )
            comments.append({
                "id": comment_id,
                "author": comment.get(q("author")),
                "text": text.strip()[:MAX_COMMENT_CHARS],
                "anchored_text": (
                    " ".join(anchors.get(comment_id, [])).strip()[:MAX_COMMENT_CHARS]
                ),
            })

    return {
        "tracked_insertion_count": len(body.findall(f".//{q('ins')}")),
        "tracked_deletion_count": len(body.findall(f".//{q('del')}")),
        "comment_count": len(comments),
        "comments": comments[:MAX_COMMENTS_REPORTED],
    }


def _format_docx(p: Path, _scope: Dict[str, Any]) -> Dict[str, Any]:
    from docx import Document  # type: ignore

    doc = Document(str(p))
    styles = {}
    for para in doc.paragraphs:
        if para.style is None:
            continue
        styles[para.style.name] = styles.get(para.style.name, 0) + 1
    return {
        "kind": "docx",
        **_docx_converted_page_count(p),
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "section_count": len(doc.sections),
        "style_histogram": styles,
        **_docx_revisions_and_comments(p),
    }


def _format_pptx(p: Path, _scope: Dict[str, Any]) -> Dict[str, Any]:
    from pptx import Presentation  # type: ignore

    pres = Presentation(str(p))
    slides = []
    for i, slide in enumerate(pres.slides, 1):
        shape_types: Dict[str, int] = {}
        for shape in slide.shapes:
            t = str(shape.shape_type)
            shape_types[t] = shape_types.get(t, 0) + 1
        slides.append({
            "index": i,
            "layout_name": getattr(slide.slide_layout, "name", None),
            "shape_types": shape_types,
        })
    return {"kind": "pptx", "slides": slides}


def _op_inspect_formatting(p: Path, scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    if kind == "xlsx":
        return _format_xlsx(p, scope)
    if kind == "docx":
        return _format_docx(p, scope)
    if kind == "pptx":
        return _format_pptx(p, scope)
    if kind == "pdf":
        # PDF formatting probe stays light — PyMuPDF font enumeration.
        try:
            import fitz  # type: ignore
            doc = fitz.open(str(p))
            try:
                fonts = set()
                for page in doc:
                    for f in page.get_fonts(full=False):
                        fonts.add(f[3])
                return {"kind": "pdf", "page_count": doc.page_count,
                        **_pdf_geometry(_fitz_page_rects(doc), doc.page_count),
                        "fonts": sorted(fonts)[:50]}
            finally:
                doc.close()
        except ImportError:
            return {"kind": "pdf", "note": "PyMuPDF not available"}
    return {"kind": kind, "note": "formatting inspection not supported for this kind"}


# ── render_to_image ──────────────────────────────────────────────────


def _png_bytes_from_pil(img) -> bytes:
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def _downsample_to_cap(png: bytes) -> bytes:
    if len(png) <= MAX_IMAGE_BYTES:
        return png
    from PIL import Image  # type: ignore
    img = Image.open(io.BytesIO(png))
    while len(png) > MAX_IMAGE_BYTES and max(img.size) > 64:
        scale = min(0.9, (MAX_IMAGE_BYTES / len(png)) ** 0.5 * 0.95)
        new_size = (
            max(64, int(img.width * scale)),
            max(64, int(img.height * scale)),
        )
        if new_size == img.size:
            break
        img = img.resize(new_size, Image.LANCZOS)
        png = _png_bytes_from_pil(img)
    if len(png) > MAX_IMAGE_BYTES:
        raise ReadDeliverableError(
            f"rendered image exceeds {MAX_IMAGE_BYTES} byte cap after downsampling"
        )
    return png


def _positive_int(value: Any, label: str) -> int:
    if isinstance(value, bool):
        raise InvalidScope(f"{label} must be a positive 1-based integer")
    if isinstance(value, int):
        parsed = value
    elif isinstance(value, str) and value.strip().isdigit():
        parsed = int(value.strip())
    else:
        raise InvalidScope(f"{label} must be a positive 1-based integer")
    if parsed < 1:
        raise InvalidScope(f"{label} must be a positive 1-based integer")
    return parsed


def _validate_scope_keys(
    scope: Dict[str, Any],
    *,
    allowed: set[str],
    source_kind: str,
) -> None:
    unknown = sorted(set(scope) - allowed)
    if unknown:
        raise InvalidScope(
            f"unknown scope keys for {source_kind}: {unknown}; "
            f"allowed keys: {sorted(allowed)}"
        )


def _render_pdf_page(p: Path, page: int) -> Tuple[bytes, int]:
    import fitz  # type: ignore
    doc = fitz.open(str(p))
    try:
        if page < 1 or page > doc.page_count:
            raise InvalidScope(
                f"page {page} out of range 1..{doc.page_count}")
        pix = doc.load_page(page - 1).get_pixmap(dpi=150)
        return pix.tobytes("png"), doc.page_count
    finally:
        doc.close()


def _render_image(p: Path) -> bytes:
    from PIL import Image  # type: ignore
    return _png_bytes_from_pil(Image.open(p))


#: How many frames one contact sheet samples, and how wide each tile is
#: drawn before the shared byte cap gets a say. Twelve at 480px lands a
#: 16:9 source in a 1920x900-ish sheet: coarse enough to stay under the cap,
#: dense enough that a ten-second clip is sampled slightly under once a
#: second. Both are part of the render contract rather than tuning knobs --
#: they set the temporal resolution every judgement about *when* something
#: happens is limited to, so they are reported in the renderer metadata.
#:
#: The frame count is public because the judge plans the render before
#: calling it: ``_VISUAL_RENDER_SCOPES`` names a ``frames`` scope per video
#: suffix, and if that number and this one were written down separately they
#: could disagree, leaving the sheet's own docstring describing a resolution
#: no judged item was actually rendered at.
VIDEO_CONTACT_SHEET_FRAMES = 12
_VIDEO_CONTACT_SHEET_TILE_WIDTH = 480
_VIDEO_CONTACT_SHEET_COLUMNS = 4


def _video_sample_times(duration_s: float, frames: int) -> List[float]:
    """Evenly spaced sample times, first frame included, last one excluded.

    The final instant is skipped deliberately: seeking to exactly the
    duration lands past the last decodable frame in most containers, which
    costs a tile and yields nothing. Sampling is uniform rather than
    criterion-aware on purpose -- the scope is per-suffix and shared by every
    item, so a sampler that chased one criterion's timings would quietly
    change what a different criterion was shown.
    """
    if frames < 1:
        raise InvalidScope("frames must be a positive integer")
    if duration_s <= 0:
        return [0.0]
    return [round(duration_s * i / frames, 3) for i in range(frames)]


def _video_label_font(tile_width: int):
    """A legible label font, or the bitmap default if no truetype is found.

    The labels are an aid, not the contract: ``sampled_timestamps_s`` in the
    renderer metadata carries the same numbers in frame order, so a host with
    no fonts installed degrades to a smaller caption rather than to a sheet
    whose frames cannot be placed in time.
    """
    from PIL import ImageFont  # type: ignore

    size = max(12, tile_width // 24)
    for candidate in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ):
        if Path(candidate).is_file():
            try:
                return ImageFont.truetype(candidate, size)
            except OSError:
                continue
    try:
        return ImageFont.load_default(size=size)
    except TypeError:  # Pillow < 10.1 has no scalable default
        return ImageFont.load_default()


def _render_video_contact_sheet(
    p: Path, frames: int
) -> Tuple[bytes, Dict[str, Any]]:
    """Tile evenly spaced, timestamped frames of a video into one image.

    One sheet per video, not one call per frame: the per-file visual budget
    (``file_cap_per_item``, ``call_cap_per_task``) counts renders and vision
    calls, and a video that spent a call per frame would price a single
    deliverable like twelve of them.

    Tiles are scaled uniformly and never padded. That is load-bearing rather
    than tidy -- two of the criteria this exists to answer are about aspect
    ratio and letterboxing, so a montage that boxed frames into a fixed cell
    would manufacture the very artefact being judged. Source dimensions go to
    the metadata so those two are answered from numbers, not from pixels.
    """
    import av  # type: ignore
    from PIL import Image, ImageDraw  # type: ignore

    container = av.open(str(p))
    try:
        if not container.streams.video:
            raise ReadDeliverableError("video has no video stream")
        stream = container.streams.video[0]
        time_base = stream.time_base
        duration_s = 0.0
        if stream.duration and time_base:
            duration_s = float(stream.duration * time_base)
        elif container.duration:
            duration_s = float(container.duration) / 1_000_000.0
        fps = float(stream.average_rate) if stream.average_rate else 0.0
        width = int(stream.codec_context.width or 0)
        height = int(stream.codec_context.height or 0)
        # Read before the container closes. Every attribute here is backed by
        # the open container, and reading one afterwards returns whatever is
        # left at that address -- which is how this first reported a
        # twelve-of-14-billion-frames coverage ratio rather than 12 of 360.
        #
        # Not every container stores a frame count (Matroska and WebM
        # routinely report 0), so fall back to duration x rate and leave it
        # None when neither is knowable. None reads as "length unknown"
        # downstream; a 0 would divide the coverage ratio by nothing.
        stored_frames = int(stream.frames or 0)
        if stored_frames > 0:
            frame_count: Optional[int] = stored_frames
        elif duration_s > 0 and fps > 0:
            frame_count = round(duration_s * fps)
        else:
            frame_count = None

        sampled: List[Tuple[float, Any]] = []
        for target in _video_sample_times(duration_s, frames):
            if time_base and duration_s > 0:
                try:
                    container.seek(int(target / time_base), stream=stream)
                except (av.AVError, OSError, ValueError):
                    continue
            decoded = None
            for frame in container.decode(video=0):
                decoded = frame
                break
            if decoded is None:
                continue
            sampled.append((target, decoded.to_image()))
        if not sampled:
            raise ReadDeliverableError(
                "no decodable video frames; the file has a video stream but "
                "nothing could be sampled from it"
            )
    finally:
        container.close()

    tile_w = _VIDEO_CONTACT_SHEET_TILE_WIDTH
    first = sampled[0][1]
    tile_h = max(1, round(first.height * tile_w / first.width))
    font = _video_label_font(tile_w)
    strip_h = max(14, tile_w // 20)
    columns = min(_VIDEO_CONTACT_SHEET_COLUMNS, len(sampled))
    rows = math.ceil(len(sampled) / columns)
    cell_h = tile_h + strip_h

    sheet = Image.new("RGB", (columns * tile_w, rows * cell_h), (32, 32, 32))
    draw = ImageDraw.Draw(sheet)
    for index, (timestamp, image) in enumerate(sampled):
        x = (index % columns) * tile_w
        y = (index // columns) * cell_h
        sheet.paste(image.convert("RGB").resize(
            (tile_w, tile_h), Image.LANCZOS), (x, y))
        draw.text(
            (x + 4, y + tile_h + 1),
            f"t={timestamp:.2f}s",
            fill=(255, 255, 255),
            font=font,
        )

    metadata: Dict[str, Any] = {
        "source_width": width or None,
        "source_height": height or None,
        "source_duration_s": round(duration_s, 3) if duration_s else None,
        "source_fps": round(fps, 3) if fps else None,
        "source_frame_count": frame_count,
        "sampled_frame_count": len(sampled),
        "sampled_timestamps_s": [t for t, _ in sampled],
        "contact_sheet_grid": f"{columns}x{rows}",
    }
    return _png_bytes_from_pil(sheet), metadata


def _find_soffice() -> Optional[str]:
    for executable in ("soffice", "libreoffice"):
        found = shutil.which(executable)
        if found:
            return found
    return None


def _minimal_libreoffice_env(work_dir: Path) -> Dict[str, str]:
    environment = {
        "PATH": os.environ.get("PATH", "/usr/bin:/bin"),
        "HOME": str(work_dir),
        "TMPDIR": str(work_dir),
    }
    for key, value in os.environ.items():
        if key == "LANG" or key.startswith("LC_"):
            environment[key] = value
    environment.setdefault("LC_ALL", "C.UTF-8")
    return environment


def _pymupdf_runtime_version() -> str:
    try:
        import fitz  # type: ignore
    except ImportError as exc:
        raise RendererDependencyError(
            "PyMuPDF is required for grading renders"
        ) from exc

    for candidate in (
        getattr(fitz, "__version__", None),
        getattr(fitz, "VersionBind", None),
    ):
        if candidate is not None and str(candidate).strip():
            return str(candidate).strip()
    version_tuple = getattr(fitz, "version", None)
    if isinstance(version_tuple, (tuple, list)) and version_tuple:
        candidate = str(version_tuple[0]).strip()
        if candidate:
            return candidate
    raise RendererDependencyError("PyMuPDF runtime version is unavailable")


@lru_cache(maxsize=8)
def _renderer_fingerprint_for_executable(executable: str) -> Dict[str, str]:
    with tempfile.TemporaryDirectory(
        prefix="gdpval-grade-renderer-version-"
    ) as temp:
        work_dir = Path(temp)
        try:
            completed = subprocess.run(
                [executable, "--headless", "--version"],
                capture_output=True,
                text=True,
                timeout=RENDERER_VERSION_TIMEOUT_SEC,
                check=False,
                env=_minimal_libreoffice_env(work_dir),
                shell=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise RendererDependencyError(
                "LibreOffice version probe timed out"
            ) from exc
        except OSError as exc:
            raise RendererDependencyError(
                f"LibreOffice version probe failed: {exc}"
            ) from exc

    if completed.returncode != 0:
        detail = (completed.stderr or completed.stdout or "no output")[-300:]
        raise RendererDependencyError(
            "LibreOffice version probe failed "
            f"(exit {completed.returncode}): {detail}"
        )
    output = completed.stdout or completed.stderr or ""
    first_line = output.splitlines()[0].strip() if output.splitlines() else ""
    if not first_line:
        raise RendererDependencyError(
            "LibreOffice version probe returned no version"
        )
    return {
        "libreoffice_binary": Path(executable).name,
        "libreoffice_version": first_line[:RENDERER_VERSION_MAX_CHARS],
        "pymupdf_version": _pymupdf_runtime_version(),
    }


def get_renderer_fingerprint() -> Dict[str, str]:
    """Return the exact fail-closed Office renderer dependency fingerprint."""
    executable = _find_soffice()
    if executable is None:
        raise RendererDependencyError(
            "LibreOffice executable not found; install soffice/libreoffice "
            "to render XLSX or PPTX deliverables"
        )
    try:
        return dict(_renderer_fingerprint_for_executable(executable))
    except ReadDeliverableError:
        raise
    except Exception as exc:  # noqa: BLE001
        raise RendererDependencyError(
            f"renderer fingerprint probe failed: {type(exc).__name__}: {exc}"
        ) from exc


def _convert_office_to_pdf(source: Path, out_dir: Path) -> Path:
    soffice = _find_soffice()
    if soffice is None:
        raise RendererDependencyError(
            "LibreOffice executable not found; install soffice/libreoffice "
            "to render XLSX or PPTX deliverables"
        )
    profile_dir = out_dir / "libreoffice-profile"
    profile_dir.mkdir()
    command = [
        soffice,
        "--headless",
        "--safe-mode",
        "--norestore",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--nolockcheck",
        f"-env:UserInstallation={profile_dir.as_uri()}",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(source),
    ]
    try:
        completed = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
            env=_minimal_libreoffice_env(out_dir),
            shell=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise ReadDeliverableError(
            "LibreOffice conversion timed out after 120 seconds"
        ) from exc
    pdf_path = out_dir / f"{source.stem}.pdf"
    if completed.returncode != 0 or not pdf_path.is_file():
        detail = (completed.stderr or completed.stdout or "no converter output")[-300:]
        raise ReadDeliverableError(
            f"LibreOffice conversion failed (exit {completed.returncode}): {detail}"
        )
    return pdf_path


def _xlsx_sheet_count(source: Path) -> int:
    import openpyxl  # type: ignore

    workbook = openpyxl.load_workbook(
        source,
        data_only=True,
        read_only=True,
        keep_vba=source.suffix.lower() == ".xlsm",
    )
    try:
        sheet_count = len(workbook.sheetnames)
        if sheet_count == 0:
            raise ReadDeliverableError("workbook has no sheets to render")
        return sheet_count
    finally:
        workbook.close()


def _op_render_to_image(p: Path, scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    if kind == "pdf":
        _validate_scope_keys(scope, allowed={"page"}, source_kind="pdf")
        page = _positive_int(scope.get("page", 1), "page")
        png, page_count = _render_pdf_page(p, page)
        png = _downsample_to_cap(png)
        return {
            "kind": "image_png_base64",
            "page": page,
            "source_kind": "pdf",
            "scope": {"page": page},
            "source_page_count": page_count,
            "renderer": {
                "rasterizer": "pymupdf",
                "pymupdf_version": _pymupdf_runtime_version(),
                "dpi": 150,
            },
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
        }
    if kind == "xlsx":
        unknown_keys = sorted(
            set(scope) - {"workbook_page", "sheet", "sheet_page"}
        )
        if unknown_keys:
            raise InvalidScope(
                f"unknown scope keys for {p.suffix.lower().lstrip('.')}: "
                f"{unknown_keys}; allowed keys: ['workbook_page']"
            )
        legacy_keys = sorted(set(scope).intersection({"sheet", "sheet_page"}))
        if legacy_keys:
            raise UnsupportedScope(
                "named sheet/sheet_page rendering is not supported; "
                "use {'workbook_page': 1} to sample the first converted "
                "workbook page"
            )
        _validate_scope_keys(
            scope, allowed={"workbook_page"}, source_kind=p.suffix.lower().lstrip(".")
        )
        workbook_page = _positive_int(
            scope.get("workbook_page", 1), "workbook_page"
        )
        if workbook_page != 1:
            raise UnsupportedScope(
                "only workbook_page=1 is supported; generic workbook page "
                "discovery is deferred"
            )
        sheet_count = _xlsx_sheet_count(p)
        renderer_fingerprint = get_renderer_fingerprint()
        with tempfile.TemporaryDirectory(prefix="gdpval-grade-render-") as temp:
            out_dir = Path(temp)
            converted = _convert_office_to_pdf(p, out_dir)
            png, converted_page_count = _render_pdf_page(
                converted, workbook_page
            )
        png = _downsample_to_cap(png)
        return {
            "kind": "image_png_base64",
            "source_kind": "xlsx",
            "scope": {"workbook_page": workbook_page},
            "source_sheet_count": sheet_count,
            "converted_page_count": converted_page_count,
            "renderer": {
                "converter": "libreoffice",
                "rasterizer": "pymupdf",
                "dpi": 150,
                **renderer_fingerprint,
            },
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
        }
    if kind == "pptx":
        from pptx import Presentation  # type: ignore

        _validate_scope_keys(scope, allowed={"slide"}, source_kind="pptx")
        slide = _positive_int(scope.get("slide", 1), "slide")
        slide_count = len(Presentation(str(p)).slides)
        if slide > slide_count:
            raise InvalidScope(
                f"slide {slide} out of range 1..{slide_count}"
            )
        renderer_fingerprint = get_renderer_fingerprint()
        with tempfile.TemporaryDirectory(prefix="gdpval-grade-render-") as temp:
            out_dir = Path(temp)
            converted = _convert_office_to_pdf(p, out_dir)
            png, converted_page_count = _render_pdf_page(converted, slide)
        png = _downsample_to_cap(png)
        return {
            "kind": "image_png_base64",
            "source_kind": "pptx",
            "scope": {"slide": slide},
            "source_slide_count": slide_count,
            "converted_page_count": converted_page_count,
            "renderer": {
                "converter": "libreoffice",
                "rasterizer": "pymupdf",
                "dpi": 150,
                **renderer_fingerprint,
            },
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
        }
    if kind == "docx":
        # Unlike the two branches above there is no pre-conversion count to
        # validate `page` against: a .docx stores no pagination, so its page
        # count does not exist until a layout engine computes one. Nothing
        # here opens the file with python-docx for that reason -- it could
        # only add a way to fail on a document LibreOffice would still
        # render. `_render_pdf_page` range-checks against the converted PDF,
        # which is the only authoritative page count there is.
        #
        # The corollary, which the returned key name is careful about: page N
        # is page N *of this conversion*. LibreOffice's line breaking is not
        # guaranteed to match Word's, so a page index is only as stable as
        # the converter. Page 1 -- the only page the judge asks for -- is
        # stable under any engine.
        _validate_scope_keys(scope, allowed={"page"}, source_kind="docx")
        page = _positive_int(scope.get("page", 1), "page")
        renderer_fingerprint = get_renderer_fingerprint()
        with tempfile.TemporaryDirectory(prefix="gdpval-grade-render-") as temp:
            out_dir = Path(temp)
            converted = _convert_office_to_pdf(p, out_dir)
            png, converted_page_count = _render_pdf_page(converted, page)
        png = _downsample_to_cap(png)
        return {
            "kind": "image_png_base64",
            "source_kind": "docx",
            "scope": {"page": page},
            "converted_page_count": converted_page_count,
            "renderer": {
                "converter": "libreoffice",
                "rasterizer": "pymupdf",
                "dpi": 150,
                **renderer_fingerprint,
            },
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
        }
    if kind == "image":
        _validate_scope_keys(scope, allowed=set(), source_kind="image")
        png = _downsample_to_cap(_render_image(p))
        return {
            "kind": "image_png_base64",
            "source_kind": "image",
            "scope": {},
            "renderer": {"rasterizer": "pillow"},
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
        }
    if kind == "video":
        _validate_scope_keys(scope, allowed={"frames"}, source_kind="video")
        frames = _positive_int(
            scope.get("frames", VIDEO_CONTACT_SHEET_FRAMES), "frames"
        )
        png, video_metadata = _render_video_contact_sheet(p, frames)
        png = _downsample_to_cap(png)
        return {
            "kind": "image_png_base64",
            "source_kind": "video",
            "scope": {"frames": frames},
            "renderer": {"rasterizer": "pyav+pillow", "tile_width": (
                _VIDEO_CONTACT_SHEET_TILE_WIDTH)},
            "base64": base64.b64encode(png).decode("ascii"),
            "byte_size": len(png),
            **video_metadata,
        }
    raise UnsupportedScope(
        f"render_to_image not supported for kind={kind}; "
        "supported: pdf, xlsx, pptx, docx, image, video. "
        "Use inspect_structure/read_content instead."
    )


# ── probe_audio / probe_video ────────────────────────────────────────


def _probe_audio_impl(p: Path, basic: bool = False) -> Dict[str, Any]:
    """Use PyAV (wheel bundles ffmpeg) per env audit decision."""
    try:
        import av  # type: ignore
    except ImportError:
        return {"backend": "none", "note": "PyAV not installed"}

    container = av.open(str(p))
    try:
        if not container.streams.audio:
            return {"backend": "av", "error": "no audio stream"}
        stream = container.streams.audio[0]
        duration_s: Optional[float] = None
        if container.duration:
            duration_s = float(container.duration) / 1_000_000.0
        elif stream.duration and stream.time_base:
            duration_s = float(stream.duration * stream.time_base)
        info: Dict[str, Any] = {
            "backend": "av",
            "sample_rate": stream.sample_rate,
            "channels": stream.channels,
            "duration_s": duration_s,
            "codec": stream.codec.name if stream.codec else None,
        }
        if basic:
            return info

        # Peak + silence ratio across decoded frames (cap at ~120s scan)
        max_amp = 0.0
        silent_samples = 0
        total_samples = 0
        scan_limit_s = 120
        for frame in container.decode(audio=0):
            arr = frame.to_ndarray()
            if arr.size == 0:
                continue
            import math
            try:
                amp = float(abs(arr).max())
                if frame.format.is_packed:
                    # int16 → normalize
                    if arr.dtype.kind == "i":
                        amp /= float(2 ** (8 * arr.dtype.itemsize - 1))
                max_amp = max(max_amp, amp)
                # silence threshold = -40 dBFS
                threshold = 10 ** (-40 / 20)
                if arr.dtype.kind == "i":
                    norm = abs(arr) / float(2 ** (8 * arr.dtype.itemsize - 1))
                else:
                    norm = abs(arr)
                silent_samples += int((norm < threshold).sum())
                total_samples += int(arr.size)
            except Exception:  # noqa: BLE001
                pass
            if duration_s and frame.pts and frame.time_base \
                    and float(frame.pts * frame.time_base) > scan_limit_s:
                break
        info["peak_amplitude_normalized"] = max_amp
        info["clipping_suspected"] = max_amp >= 0.999
        info["silence_ratio"] = (silent_samples / total_samples) if total_samples else None
        info["scan_capped_seconds"] = scan_limit_s
        return info
    finally:
        container.close()


def _op_probe_audio(p: Path, _scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    if kind != "audio":
        return {"kind": kind, "note": "not an audio file"}
    info = _probe_audio_impl(p, basic=False)
    info["kind"] = "audio"
    return info


def _video_audio_track_count(p: Path) -> Optional[int]:
    """How many audio streams a video container carries, or ``None``.

    Deliberately not routed through ``_probe_video_impl``: that helper returns
    early when a container has no *video* stream, which would report ``None``
    for a file that is nothing but sound. Counting the audio streams directly
    answers the question that was asked.
    """
    try:
        import av  # type: ignore
    except ImportError:
        return None
    try:
        container = av.open(str(p))
    except Exception:  # noqa: BLE001
        return None
    try:
        return len(container.streams.audio)
    except Exception:  # noqa: BLE001
        return None
    finally:
        container.close()


def _probe_video_impl(p: Path, basic: bool = False) -> Dict[str, Any]:
    try:
        import av  # type: ignore
    except ImportError:
        return {"backend": "none", "note": "PyAV not installed"}
    container = av.open(str(p))
    try:
        if not container.streams.video:
            return {"backend": "av", "error": "no video stream"}
        vs = container.streams.video[0]
        duration_s = None
        if container.duration:
            duration_s = float(container.duration) / 1_000_000.0
        fps = None
        try:
            if vs.average_rate:
                fps = float(vs.average_rate)
        except Exception:  # noqa: BLE001
            pass
        return {
            "backend": "av",
            "codec": vs.codec.name if vs.codec else None,
            "width": vs.width,
            "height": vs.height,
            "fps": fps,
            "duration_s": duration_s,
            "audio_tracks": len(container.streams.audio),
            "video_tracks": len(container.streams.video),
            "audio": _audio_stream_summary(container),
        }
    finally:
        container.close()


def _audio_stream_summary(container: Any) -> Optional[Dict[str, Any]]:
    """How a video container's first audio stream is encoded, or ``None``.

    A video probe used to answer with two track counts and nothing else, so a
    criterion like "the audio sample rate is 44.1 kHz or 48 kHz" reached a
    judge that had no way to know and marked it failed rather than excluded.

    That is the same defect PR #276 fixed one field over. There, a container
    reporting ``"audio_tracks": 1`` was sent to a judge that could not listen;
    the count was right there in the evidence and unusable. Here the count is
    right there and the *rate* is missing, so the criterion fails on a fact
    the file states plainly. Fixing the routing and leaving this would have
    meant the same bug survived in the field next door.

    The field names deliberately match ``_probe_audio_impl`` so the judge sees
    one vocabulary whether the sound arrived in a ``.wav`` or an ``.mp4``.
    """
    try:
        streams = container.streams.audio
    except Exception:  # noqa: BLE001
        return None
    if not streams:
        return None
    stream = streams[0]

    def _field(name: str) -> Any:
        # Each read is guarded separately: a container can describe its rate
        # and not its channel layout, and one missing attribute must not turn
        # the whole summary back into the silence this replaced.
        try:
            return getattr(stream, name, None)
        except Exception:  # noqa: BLE001
            return None

    codec = _field("codec")
    return {
        "sample_rate": _field("sample_rate"),
        "channels": _field("channels"),
        "codec": getattr(codec, "name", None) if codec else None,
    }


def _op_probe_video(p: Path, _scope: Dict[str, Any]) -> Dict[str, Any]:
    kind = _kind_of(p)
    if kind != "video":
        return {"kind": kind, "note": "not a video file"}
    info = _probe_video_impl(p, basic=False)
    info["kind"] = "video"
    return info


# ── Public entrypoint ────────────────────────────────────────────────


_OP_TABLE = {
    "inspect_structure": _op_inspect_structure,
    "read_content": _op_read_content,
    "inspect_formatting": _op_inspect_formatting,
    "render_to_image": _op_render_to_image,
    "probe_audio": _op_probe_audio,
    "probe_video": _op_probe_video,
}


def read_deliverable(
    op: str,
    path: str,
    *,
    base_dir: str,
    scope: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Dispatch a read-only inspection op on a deliverable file.

    Args:
        op: One of ``READ_DELIVERABLE_OPS``.
        path: Path to the file. Resolved against ``base_dir``; traversal
              outside the base is rejected.
        base_dir: Trusted base directory (required keyword).
        scope: Op-specific scope dict (sheet name, page range, …).
               On a ``.zip``, ``{"member": "<name>"}`` runs the op against
               that member of the archive instead of the archive itself; the
               rest of the scope is passed through to it unchanged.

    Returns:
        Envelope dict — see module docstring.
    """
    if op not in _OP_TABLE:
        return _envelope_error(
            f"unknown op {op!r}; allowed: {list(READ_DELIVERABLE_OPS)}",
            kind="bad_op",
        )
    if scope is not None and not isinstance(scope, dict):
        return _envelope_error("scope must be an object or null", kind="bad_scope")
    resolved = _resolve_trusted_path(path, base_dir)
    if resolved is None:
        return _envelope_error(
            f"path {path!r} not found or escapes base_dir",
            kind="bad_path",
        )
    fn = _OP_TABLE[op]
    try:
        scope = dict(scope or {})
        # A member is addressed by the same op as a loose file, so the
        # extraction happens here rather than inside each op: `probe_audio` on
        # a WAV inside an archive is `probe_audio` on a WAV.
        if "member" in scope:
            if _kind_of(resolved) != "zip":
                raise InvalidScope(
                    f"scope key 'member' is only valid on a zip archive, "
                    f"not on a {_kind_of(resolved)} file"
                )
            member = scope.pop("member")
            if not isinstance(member, str) or not member:
                raise InvalidScope("scope key 'member' must be a non-empty string")
            with open_archive_member(resolved, member) as member_path:
                return _envelope_ok(fn(member_path, scope))
        data = fn(resolved, scope)
        return _envelope_ok(data)
    except RendererDependencyError as exc:
        return _envelope_error(str(exc), kind="dependency_missing")
    except InvalidScope as exc:
        return _envelope_error(str(exc), kind="bad_scope")
    except UnsupportedScope as exc:
        return _envelope_error(str(exc), kind="unsupported_scope")
    except ReadDeliverableError as exc:
        return _envelope_error(str(exc), kind="op_error")
    except ImportError as exc:
        package = "PyMuPDF" if exc.name == "fitz" else exc.name
        return _envelope_error(
            f"required Python library missing: {package}",
            kind="dependency_missing",
        )
    except Exception as exc:  # noqa: BLE001
        return _envelope_error(
            f"{type(exc).__name__}: {exc}",
            kind="exception",
        )
