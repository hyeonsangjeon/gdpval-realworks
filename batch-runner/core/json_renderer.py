"""
JSON Renderer for fair cross-model comparison.

LLM outputs JSON specification only (no code), and a fixed renderer generates
the actual files. This ensures all models use identical file generation logic,
making comparisons fair.

Supported file types:
- docx: Word documents
- xlsx: Excel spreadsheets
- pdf: PDF documents
- pptx: PowerPoint presentations
- png: Charts or diagrams
- md: Markdown files
- html: HTML files
- json: JSON files
"""

import json
import io
import re
from typing import Optional

from core.config import DEFAULT_TOKENS
from core.llm_client import complete


class JsonRenderer:
    """LLM JSON spec → fixed renderer for portable file generation"""

    JSON_PROMPT = """Create deliverables by outputting JSON specification only (no code).

FORMAT:
{{
  "deliverable_text": "Your main text response here",
  "deliverables": [
    {{
      "type": "docx",
      "filename": "report.docx",
      "content": {{
        "title": "Report Title",
        "sections": [
          {{"heading": "Section 1", "body": "Content here..."}}
        ]
      }}
    }},
    {{
      "type": "xlsx",
      "filename": "data.xlsx",
      "content": {{
        "sheets": [
          {{
            "name": "Sheet1",
            "headers": ["Col A", "Col B"],
            "rows": [["val1", "val2"]]
          }}
        ]
      }}
    }},
    {{
      "type": "pdf",
      "filename": "document.pdf",
      "content": {{
        "title": "Document Title",
        "sections": [
          {{"heading": "Section 1", "body": "Content..."}}
        ]
      }}
    }},
    {{
      "type": "pptx",
      "filename": "presentation.pptx",
      "content": {{
        "slides": [
          {{"title": "Slide 1", "body": "Content..."}}
        ]
      }}
    }},
    {{
      "type": "png",
      "filename": "chart.png",
      "content": {{
        "chart_type": "bar",
        "title": "Sales by Quarter",
        "labels": ["Q1", "Q2", "Q3", "Q4"],
        "values": [100, 150, 120, 180]
      }}
    }},
    {{
      "type": "md",
      "filename": "notes.md",
      "content": {{
        "text": "# Heading\\n\\nContent here..."
      }}
    }}
  ]
}}

SUPPORTED TYPES: docx, xlsx, pdf, pptx, png, md, html, json

CHART TYPES (for png):
- "bar": Bar chart
- "line": Line chart
- "pie": Pie chart

TASK:
{task_prompt}

OUTPUT: Only JSON in ```json``` blocks, no explanations.
"""

    def __init__(self, llm_client, max_completion_tokens: Optional[int] = None):
        """
        Initialize JSON renderer with LLM client.

        Args:
            llm_client: AzureOpenAI client instance for JSON spec generation
            max_completion_tokens: Completion token cap override
        """
        self.llm_client = llm_client
        self.max_completion_tokens = (
            max_completion_tokens
            if max_completion_tokens is not None
            else DEFAULT_TOKENS["json_render"]
        )

    def run(self, task_prompt: str, model: str) -> dict:
        """
        Generate JSON spec via LLM and render files with fixed renderers.

        Args:
            task_prompt: The task instruction
            model: Model deployment name

        Returns:
            dict with keys:
                - success (bool): Whether execution succeeded
                - text (str): deliverable_text from JSON
                - files (list): List of generated files [{filename, content}]
                - error (str, optional): Error message if failed
        """
        try:
            # Step 1: Generate JSON spec using LLM
            prompt = self.JSON_PROMPT.format(task_prompt=task_prompt)

            messages = [
                {"role": "system", "content": "You are a JSON specification generator. Output only JSON in ```json``` blocks."},
                {"role": "user", "content": prompt}
            ]

            response, latency = complete(
                client=self.llm_client,
                model=model,
                messages=messages,
                max_completion_tokens=self.max_completion_tokens
            )

            response_text = response.choices[0].message.content

            # Step 2: Extract and parse JSON
            try:
                spec = self._extract_json(response_text)
            except json.JSONDecodeError as e:
                return {
                    "success": False,
                    "text": "",
                    "files": [],
                    "error": f"JSON parse error: {e}. Response: {response_text[:200]}..."
                }

            # Step 3: Render files using fixed renderers
            output_files = []
            for item in spec.get("deliverables", []):
                try:
                    file_data = self._render(item)
                    if file_data:
                        output_files.append(file_data)
                except Exception as e:
                    print(f"Warning: Failed to render {item.get('filename', 'unknown')}: {e}")

            return {
                "success": True,
                "text": spec.get("deliverable_text", ""),
                "files": output_files
            }

        except Exception as e:
            return {
                "success": False,
                "text": "",
                "files": [],
                "error": f"JSON rendering failed: {str(e)}"
            }

    def _extract_json(self, text: str) -> dict:
        """
        Extract JSON from ```json``` code blocks.

        Args:
            text: LLM response text

        Returns:
            Parsed JSON dict

        Raises:
            json.JSONDecodeError: If JSON is invalid
        """
        # Try to extract from ```json blocks
        pattern = r"```json\n(.*?)```"
        matches = re.findall(pattern, text, re.DOTALL)

        if matches:
            return json.loads(matches[0].strip())

        # Fallback: Try to parse entire text as JSON
        return json.loads(text.strip())

    def _render(self, item: dict) -> Optional[dict]:
        """
        Render a single file using appropriate renderer.

        Args:
            item: Deliverable spec with type, filename, content

        Returns:
            dict with filename and content (bytes), or None if unsupported type
        """
        file_type = item.get("type")
        filename = item.get("filename")
        content = item.get("content")

        if not all([file_type, filename, content]):
            return None

        renderers = {
            "docx": self._render_docx,
            "xlsx": self._render_xlsx,
            "pdf": self._render_pdf,
            "pptx": self._render_pptx,
            "png": self._render_png,
            "md": self._render_text,
            "html": self._render_text,
            "json": self._render_json,
        }

        renderer = renderers.get(file_type)
        if renderer:
            try:
                file_content = renderer(content, filename)
                return {
                    "filename": filename,
                    "content": file_content
                }
            except Exception as e:
                print(f"Renderer error for {filename}: {e}")
                return None

        return None

    def _render_docx(self, content: dict, filename: str) -> bytes:
        """Render DOCX using python-docx"""
        from docx import Document

        doc = Document()

        # Add title
        if content.get("title"):
            doc.add_heading(content["title"], 0)

        # Add sections
        for section in content.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], 1)
            if section.get("body"):
                doc.add_paragraph(section["body"])

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    def _render_xlsx(self, content: dict, filename: str) -> bytes:
        """Render XLSX using openpyxl"""
        from openpyxl import Workbook

        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet

        for i, sheet_data in enumerate(content.get("sheets", [])):
            ws = wb.create_sheet(title=sheet_data.get("name", f"Sheet{i+1}"))

            # Add headers
            if sheet_data.get("headers"):
                ws.append(sheet_data["headers"])

            # Add rows
            for row in sheet_data.get("rows", []):
                ws.append(row)

        buffer = io.BytesIO()
        wb.save(buffer)
        return buffer.getvalue()

    def _render_pdf(self, content: dict, filename: str) -> bytes:
        """Render PDF using reportlab"""
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Add title
        if content.get("title"):
            story.append(Paragraph(content["title"], styles["Title"]))
            story.append(Spacer(1, 12))

        # Add sections
        for section in content.get("sections", []):
            if section.get("heading"):
                story.append(Paragraph(section["heading"], styles["Heading1"]))
                story.append(Spacer(1, 6))
            if section.get("body"):
                story.append(Paragraph(section["body"], styles["Normal"]))
                story.append(Spacer(1, 12))

        doc.build(story)
        return buffer.getvalue()

    def _render_pptx(self, content: dict, filename: str) -> bytes:
        """Render PPTX using python-pptx"""
        from pptx import Presentation

        prs = Presentation()

        for slide_data in content.get("slides", []):
            # Use title and content layout (index 1)
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            # Set title
            if slide_data.get("title"):
                slide.shapes.title.text = slide_data["title"]

            # Set body
            if slide_data.get("body") and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data["body"]

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _render_png(self, content: dict, filename: str) -> bytes:
        """Render PNG using matplotlib or Pillow"""
        # Chart-based PNG (matplotlib)
        if content.get("chart_type"):
            import matplotlib
            matplotlib.use('Agg')  # Non-interactive backend
            import matplotlib.pyplot as plt

            fig, ax = plt.subplots(figsize=(10, 6))

            chart_type = content["chart_type"]
            labels = content.get("labels", [])
            values = content.get("values", [])

            if chart_type == "bar":
                ax.bar(labels, values)
            elif chart_type == "line":
                ax.plot(labels, values, marker='o')
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct='%1.1f%%')

            if content.get("title"):
                ax.set_title(content["title"])

            buffer = io.BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
            plt.close()
            return buffer.getvalue()

        # Diagram-based PNG (Pillow)
        else:
            from PIL import Image, ImageDraw

            width = content.get("width", 800)
            height = content.get("height", 600)
            bg_color = content.get("background", "white")

            img = Image.new("RGB", (width, height), bg_color)
            draw = ImageDraw.Draw(img)

            # Add text
            if content.get("text"):
                draw.text((50, 50), content["text"], fill="black")

            # Add shapes
            for shape in content.get("shapes", []):
                if shape["type"] == "rectangle":
                    draw.rectangle(shape["coords"], outline=shape.get("color", "black"))
                elif shape["type"] == "ellipse":
                    draw.ellipse(shape["coords"], outline=shape.get("color", "black"))

            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            return buffer.getvalue()

    def _render_text(self, content: dict, filename: str) -> bytes:
        """Render text files (MD, HTML)"""
        text = content.get("text", "")
        return text.encode("utf-8")

    def _render_json(self, content: dict, filename: str) -> bytes:
        """Render JSON files"""
        return json.dumps(content, indent=2, ensure_ascii=False).encode("utf-8")
